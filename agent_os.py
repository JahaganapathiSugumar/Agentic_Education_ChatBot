# ==============================================================================
# --- 1. IMPORTS (Cleaned Version)
# ==============================================================================

import os
from dotenv import load_dotenv
import torch
import requests
import traceback
from flask import Flask, request
from collections import defaultdict
import re
from datetime import datetime, timedelta
from io import BytesIO
import tempfile
import numpy as np
import whisper 
import base64
import moviepy
import moviepy.editor
from fpdf import FPDF
from pydub import AudioSegment
from moviepy.editor import ImageClip, concatenate_videoclips, AudioFileClip
from gtts import gTTS
from PIL import Image, ImageDraw
import json
import threading

# --- AI & Machine Learning ---
import google.generativeai as genai
from sklearn.metrics.pairwise import cosine_similarity
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_community.document_loaders import UnstructuredFileLoader
from langchain_community.vectorstores import FAISS
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.schema import Document

# --- Google Cloud APIs ---
import google.auth.transport.requests
from google.oauth2.credentials import Credentials
from google_auth_oauthlib.flow import InstalledAppFlow
from googleapiclient.discovery import build
from googleapiclient.errors import HttpError
from googleapiclient.http import MediaFileUpload

# --- Utility & File Generation ---
from fpdf import FPDF
from pydub import AudioSegment
from gtts import gTTS
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches

# --- PDF and Document Processing ---
import fitz  # PyMuPDF
from docx import Document as DocxDocument

# --- Firebase Imports ---
import firebase_admin
from firebase_admin import credentials, firestore
import uuid
from datetime import datetime

#
# ---------------------------------
# --- MODEL LOADING (at startup)
# ---------------------------------
#

print("Loading Whisper model...")
whisper_model = whisper.load_model("base")
print("Whisper model loaded successfully.")
os.environ["KMP_DUPLICATE_LIB_OK"] = "TRUE"

#
# ---------------------------------
# --- ENVIRONMENT & CONFIG
# ---------------------------------
#
app = Flask(__name__)

load_dotenv()
# --- WhatsApp & Meta Config ---
VERIFY_TOKEN = os.getenv("VERIFY_TOKEN")
ACCESS_TOKEN = os.getenv("ACCESS_TOKEN")
PHONE_NUMBER_ID = os.getenv("PHONE_NUMBER_ID")
TEMPLATE_NAME = os.getenv("TEMPLATE_NAME")

# Google & GenAI
GENAI_API_KEY = os.getenv("GENAI_API_KEY")
GOOGLE_FORMS_SERVICE_ACCOUNT_KEY_PATH = os.getenv("GOOGLE_FORMS_SERVICE_ACCOUNT_KEY_PATH")
GOOGLE_FORMS_TARGET_FOLDER_ID = os.getenv("GOOGLE_FORMS_TARGET_FOLDER_ID")
CLASSROOM_COURSE_ID = os.getenv("CLASSROOM_COURSE_ID")

print("Verify Token:", VERIFY_TOKEN)
print("Configs loaded successfully")

# --- List of users to notify on startup ---
student_phone_numbers = [
    "916379613654",
    "918870420449",
]

# --- Google AI & API Config ---
genai.configure(api_key=GENAI_API_KEY)
GOOGLE_FORMS_SERVICE_ACCOUNT_KEY_PATH = 'sahayak-465916-8bf5ddce5515.json'
GOOGLE_FORMS_TARGET_FOLDER_ID = '11vL5AgJiLYbX6fgAJFoMEjea9jiR1WMk'
CLASSROOM_COURSE_ID = '791255014049'

# --- State & Memory Management ---
user_states = defaultdict(str)
user_temp_data = defaultdict(dict)
user_memory = defaultdict(list)
MAX_HISTORY = 50

# --- Dynamic RAG Configuration ---
DYNAMIC_VECTOR_INDEX_PATH = "vector_index/dynamic_uploads"
os.makedirs(DYNAMIC_VECTOR_INDEX_PATH, exist_ok=True)

# Add this to your IMPORTS section
QUESTION_PAPER_VECTOR_INDEX_PATH = "vector_index/question_paper_materials"
os.makedirs(QUESTION_PAPER_VECTOR_INDEX_PATH, exist_ok=True)

# Separate embeddings for question paper materials
qp_embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")

def add_to_qp_vector_index(file_bytes, filename, material_type, metadata=None):
    """Add documents to question paper specific vector index"""
    try:
        # Extract text from uploaded file
        text_content = extract_text_from_file(file_bytes, filename)
        if not text_content:
            return False, "Could not extract text from the uploaded file."
        
        # Prepare metadata
        if metadata is None:
            metadata = {}
        
        base_metadata = {
            'source': filename,
            'material_type': material_type,  # syllabus, chapter_notes, sample_papers
            'uploaded_at': datetime.now().isoformat(),
            'type': 'question_paper_material'
        }
        base_metadata.update(metadata)
        
        # Create document object
        document = Document(
            page_content=text_content,
            metadata=base_metadata
        )
        
        # Split text into chunks
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=100
        )
        split_documents = text_splitter.split_documents([document])
        
        # Load existing QP vector store or create new one
        if os.path.exists(QUESTION_PAPER_VECTOR_INDEX_PATH) and os.listdir(QUESTION_PAPER_VECTOR_INDEX_PATH):
            try:
                vector_store = FAISS.load_local(
                    QUESTION_PAPER_VECTOR_INDEX_PATH, 
                    qp_embeddings, 
                    allow_dangerous_deserialization=True
                )
                vector_store.add_documents(split_documents)
                print(f"Added {len(split_documents)} chunks to QP vector store")
            except Exception as e:
                print(f"Error loading QP index, creating new: {e}")
                vector_store = FAISS.from_documents(split_documents, qp_embeddings)
        else:
            vector_store = FAISS.from_documents(split_documents, qp_embeddings)
            print(f"Created new QP vector store with {len(split_documents)} chunks")
        
        vector_store.save_local(QUESTION_PAPER_VECTOR_INDEX_PATH)
        return True, f"Successfully added {filename} to question paper materials with {len(split_documents)} chunks."
        
    except Exception as e:
        print(f"Error adding document to QP vector index: {e}")
        return False, f"Error processing file: {str(e)}"

def query_qp_rag(question, material_type=None, k=3):
    """Query question paper specific RAG"""
    try:
        if not os.path.exists(QUESTION_PAPER_VECTOR_INDEX_PATH) or not os.listdir(QUESTION_PAPER_VECTOR_INDEX_PATH):
            return {"result": "No question paper materials uploaded yet.", "source": "no_content"}
        
        qp_db = FAISS.load_local(QUESTION_PAPER_VECTOR_INDEX_PATH, qp_embeddings, allow_dangerous_deserialization=True)
        
        # Filter by material type if specified
        if material_type:
            def filter_by_type(doc):
                return doc.metadata.get('material_type') == material_type
            retriever = qp_db.as_retriever(
                search_kwargs={"k": k, "filter": filter_by_type}
            )
        else:
            retriever = qp_db.as_retriever(search_kwargs={"k": k})
        
        docs = retriever.invoke(question)
        
        if not docs:
            return {"result": f"No relevant {material_type} found.", "source": "no_matches"}
        
        context = "\n\n".join([doc.page_content for doc in docs])
        
        # Generate answer
        prompt = f"""Based on the following context from question paper materials, answer the question:

Context:
{context}

Question: {question}"""

        model = genai.GenerativeModel('gemini-2.0-flash')
        response = model.generate_content(prompt)
        
        return {"result": response.text, "source": "qp_rag"}
        
    except Exception as e:
        print(f"Error in QP RAG query: {e}")
        return {"result": "Error during question paper materials lookup.", "source": "error"}
    



SUBJECT_EXAMPLES = [
    "What is this topic about?", "Explain the main concepts.", "Tell me about the key principles.",
    "How does this work?", "What are the types of approaches?", "Describe the fundamental ideas.",
    "What are the main components?", "How does the system function?", "What is the basic theory?",
    "Explain the core concepts.", "What are the important terms?", "Discuss the main topics.",
    "What are the fundamental principles?", "Explain the basic structure.", "What is the overview?",
    "What topics are covered?", "Show me the syllabus.", "List the main units?",
    "What are the key areas?", "How is the content organized?", "Tell me about the course outline."
]

FORM_GENERATION_EXAMPLES = [
    "Generate a quiz with 10 MCQ questions.", "Create a quiz on this topic.",
    "Make 5 multiple choice questions.", "I need a short quiz.",
    "Can you generate a Google Form for me?", "Form with 7 questions.",
    "Build a form for 20 questions.", "Generate a quiz form."
]

CLASSROOM_ASSIGNMENT_EXAMPLES = [
    "Create an assignment for my class.", "Post a new assignment.",
    "Make an assignment for the course.", "Create a quiz in Classroom.",
    "New assignment for my students.", "Assign a new quiz.",
    "Create homework for the class."
]

WORKSHEET_GENERATION_EXAMPLES = [
    "Generate a worksheet.", "Create a PDF worksheet.",
    "Give me a short answer worksheet.", "Make a worksheet with 5 fill-in-the-blanks.",
    "Can you generate a practice sheet?", "Generate a PDF review.",
    "Generate 10 multiple choice questions worksheet.", "Create an MCQ worksheet.",
    "Give me 5 short answer questions.", "Generate a short answer worksheet.",
    "Make 3 long answer questions.", "Give me an essay question worksheet.",
    "Generate 2 numerical questions.", "Create a worksheet with calculation problems."
]

embeddings_model_for_classification = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
dynamic_embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")

subject_example_embeddings = None
form_generation_example_embeddings = None
classroom_assignment_example_embeddings = None
worksheet_generation_example_embeddings = None

SIMILARITY_THRESHOLD_SUBJECT = 0.65
SIMILARITY_THRESHOLD_FORM = 0.70
SIMILARITY_THRESHOLD_CLASSROOM = 0.70
SIMILARITY_THRESHOLD_WORKSHEET = 0.75


# --- API Service Clients (Initialized later) ---
forms_service = None
drive_service = None
classroom_service = None
speech_client = None
creds = None
calendar_service = None

#
# ==============================================================================
# --- DYNAMIC RAG FUNCTIONS
# ==============================================================================
#

def extract_text_from_file(file_bytes, filename):
    """Extract text from various file formats."""
    try:
        file_extension = filename.lower().split('.')[-1]
        
        # Create temporary file
        with tempfile.NamedTemporaryFile(delete=False, suffix=f'.{file_extension}') as temp_file:
            temp_file.write(file_bytes)
            temp_file_path = temp_file.name
        
        text = ""
        
        if file_extension == 'pdf':
            # Using PyMuPDF for PDF extraction
            doc = fitz.open(temp_file_path)
            for page in doc:
                text += page.get_text()
            doc.close()
            
        elif file_extension in ['docx', 'doc']:
            # Using python-docx for Word documents
            doc = DocxDocument(temp_file_path)
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
                
        elif file_extension in ['pptx', 'ppt']:
            # Using python-pptx for PowerPoint
            prs = Presentation(temp_file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                        
        elif file_extension == 'txt':
            with open(temp_file_path, 'r', encoding='utf-8') as f:
                text = f.read()
                
        else:
            # Fallback to unstructured for other formats
            try:
                loader = UnstructuredFileLoader(temp_file_path)
                documents = loader.load()
                text = "\n".join([doc.page_content for doc in documents])
            except Exception as e:
                print(f"Error with unstructured loader: {e}")
                return None
        
        # Clean up temporary file
        os.unlink(temp_file_path)
        
        return text.strip() if text else None
        
    except Exception as e:
        print(f"Error extracting text from {filename}: {e}")
        # Clean up temporary file if it exists
        if 'temp_file_path' in locals() and os.path.exists(temp_file_path):
            try:
                os.unlink(temp_file_path)
            except:
                pass
        return None

def add_documents_to_vector_index(file_bytes, filename, metadata=None):
    """Add new documents to the dynamic vector index."""
    try:
        # Extract text from uploaded file
        text_content = extract_text_from_file(file_bytes, filename)
        if not text_content:
            return False, "Could not extract text from the uploaded file."
        
        # Prepare metadata
        if metadata is None:
            metadata = {}
        
        base_metadata = {
            'source': filename,
            'uploaded_at': datetime.now().isoformat(),
            'type': 'user_upload'
        }
        base_metadata.update(metadata)
        
        # Create document object
        document = Document(
            page_content=text_content,
            metadata=base_metadata
        )
        
        # Split text into chunks
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=100
        )
        split_documents = text_splitter.split_documents([document])
        
        # Load existing vector store or create new one
        if os.path.exists(DYNAMIC_VECTOR_INDEX_PATH) and os.listdir(DYNAMIC_VECTOR_INDEX_PATH):
            try:
                vector_store = FAISS.load_local(
                    DYNAMIC_VECTOR_INDEX_PATH, 
                    dynamic_embeddings, 
                    allow_dangerous_deserialization=True
                )
                # Add new documents
                vector_store.add_documents(split_documents)
                print(f"Added {len(split_documents)} chunks to existing vector store")
            except Exception as e:
                print(f"Error loading existing index, creating new: {e}")
                vector_store = FAISS.from_documents(split_documents, dynamic_embeddings)
        else:
            # Create new vector store
            vector_store = FAISS.from_documents(split_documents, dynamic_embeddings)
            print(f"Created new vector store with {len(split_documents)} chunks")
        
        # Save the updated vector store
        vector_store.save_local(DYNAMIC_VECTOR_INDEX_PATH)
        
        return True, f"Successfully added {filename} to knowledge base with {len(split_documents)} chunks."
        
    except Exception as e:
        print(f"Error adding document to vector index: {e}")
        traceback.print_exc()
        return False, f"Error processing file: {str(e)}"

# Update the function that was calling the OS-specific one
def query_dynamic_rag(question, k=3):
        """Query ONLY dynamic user uploads (no static database)"""
        try:
            # Check if dynamic database exists and has content
            if not os.path.exists(DYNAMIC_VECTOR_INDEX_PATH) or not os.listdir(DYNAMIC_VECTOR_INDEX_PATH):
                return {
                    "result": "❌ No materials uploaded yet. Please upload documents first using 'Upload Materials' option.", 
                    "source": "no_content"
                }
            
            # Load only dynamic database
            dynamic_db = FAISS.load_local(
                DYNAMIC_VECTOR_INDEX_PATH, 
                dynamic_embeddings, 
                allow_dangerous_deserialization=True
            )
            dynamic_retriever = dynamic_db.as_retriever(search_kwargs={"k": k})
            dynamic_docs = dynamic_retriever.invoke(question)
            
            if not dynamic_docs:
                return {
                    "result": "❌ No relevant information found in your uploaded documents. Try uploading more relevant materials.", 
                    "source": "no_matches"
                }
            
            # Combine context from dynamic docs only
            context = "\n\n".join([doc.page_content for doc in dynamic_docs])
            
            # Identify sources
            sources = []
            for doc in dynamic_docs:
                if 'source' in doc.metadata:
                    source_name = os.path.basename(doc.metadata['source'])
                    if source_name not in sources:
                        sources.append(source_name)
            
            source_info = "\n\n📚 Sources: " + ", ".join(sources) if sources else ""
            
            # Generate answer using Gemini
            prompt = f"""Based ONLY on the following context from user-uploaded documents, answer the question accurately:

    Context:
    {context}

    Question: {question}

    If the context doesn't contain enough information to answer the question, clearly state that and suggest uploading relevant documents."""

            model = genai.GenerativeModel('gemini-2.0-flash')
            response = model.generate_content(prompt)
            
            final_answer = response.text + source_info
            return {"result": final_answer, "source": "dynamic_rag"}
            
        except Exception as e:
            print(f"Error in dynamic RAG query: {e}")
            return {"result": "❌ Error during knowledge base lookup. Please try again.", "source": "error"}

def get_uploaded_files_list():
    """Get list of files that ACTUALLY exist in the dynamic vector database"""
    try:
        if not os.path.exists(DYNAMIC_VECTOR_INDEX_PATH) or not os.listdir(DYNAMIC_VECTOR_INDEX_PATH):
            return []
        
        vector_store = FAISS.load_local(
            DYNAMIC_VECTOR_INDEX_PATH, 
            dynamic_embeddings, 
            allow_dangerous_deserialization=True
        )
        
        # Extract unique sources from metadata that actually exist in the vector store
        sources = set()
        
        # Method 1: Check docstore
        if hasattr(vector_store, 'docstore') and hasattr(vector_store.docstore, '_dict'):
            for doc_id, doc in vector_store.docstore._dict.items():
                if hasattr(doc, 'metadata') and 'source' in doc.metadata:
                    sources.add(doc.metadata['source'])
        
        # Method 2: Check index if available
        if hasattr(vector_store, 'index_to_docstore_id'):
            for doc_id in vector_store.index_to_docstore_id.values():
                doc = vector_store.docstore.search(doc_id)
                if doc and hasattr(doc, 'metadata') and 'source' in doc.metadata:
                    sources.add(doc.metadata['source'])
        
        return list(sources)
        
    except Exception as e:
        print(f"Error getting uploaded files list from vector DB: {e}")
        return []

def handle_document_upload(sender, media_id, filename):
    """Worker function to handle single document upload and processing."""
    print(f"Starting document upload task for {sender}, file: {filename}")
    send_whatsapp_message(sender, f"Received {filename}. Processing and adding to knowledge base...")
    
    if media_url := get_media_url(media_id):
        if file_bytes := download_media_file(media_url):
            # Add document to vector index
            success, message = add_documents_to_vector_index(file_bytes, filename)
            
            if success:
                # Save file metadata to Firebase ONLY if successfully added to vector DB
                save_uploaded_file_metadata(sender, filename, len(file_bytes), message.split(' ')[-2] if 'chunks' in message else 0)
                send_whatsapp_message(sender, f"✅ Success! {message}\n\nI can now answer questions based on this content.")
            else:
                send_whatsapp_message(sender, f"❌ Failed to process file: {message}")
        else:
            send_whatsapp_message(sender, "❌ Sorry, I couldn't download the file.")
    else:
        send_whatsapp_message(sender, "❌ Sorry, I couldn't access the file.")
    
    # ✅ Clear state and show menu ONLY AFTER processing is complete
    if sender in user_states:
        del user_states[sender]
        clear_user_state_from_firebase(sender)
    end_conversation_and_show_menu(sender, None)
    print(f"Document upload task for {sender} finished.")

def get_upload_queue_status(sender):
    """Get current upload queue status"""
    if sender in user_temp_data and 'bulk_upload_items' in user_temp_data[sender]:
        items = user_temp_data[sender]['bulk_upload_items']
        return len(items), [item['filename'] for item in items]
    return 0, []

# Optional: Add a command to check queue status
def check_upload_queue(sender):
    """Check current upload queue"""
    count, files = get_upload_queue_status(sender)
    if count > 0:
        file_list = "\n".join([f"• {f}" for f in files[:3]])
        if count > 3:
            file_list += f"\n• ... and {count - 3} more"
        send_whatsapp_message(sender, f"📋 Upload Queue: {count} files\n{file_list}")
    else:
        send_whatsapp_message(sender, "📋 No files in upload queue")
#
# ---------------------------------
# --- API INITIALIZATION
# ---------------------------------
#
SCOPES = [
    'https://www.googleapis.com/auth/classroom.courses',
    'https://www.googleapis.com/auth/classroom.coursework.students',
    'https://www.googleapis.com/auth/classroom.announcements',
    'https://www.googleapis.com/auth/drive',
    'https://www.googleapis.com/auth/forms.body',
    'https://www.googleapis.com/auth/calendar.events'
]

def init_google_apis():
    """Authenticates user and initializes Google API clients with proper scopes"""
    global classroom_service, drive_service, forms_service, calendar_service, creds
    
    # Delete the old token to force re-authentication with new scopes
    #if os.path.exists('token.json'):
    #    os.remove('token.json')
    #    print("Removed old token to force re-authentication with new scopes")
    
    if os.path.exists('token.json'):
        creds = Credentials.from_authorized_user_file('token.json', SCOPES)
    if not creds or not creds.valid:
        if creds and creds.expired and creds.refresh_token:
            creds.refresh(google.auth.transport.requests.Request())
        else:
            flow = InstalledAppFlow.from_client_secrets_file('credentials.json', SCOPES)
            creds = flow.run_local_server(port=0)
        with open('token.json', 'w') as token:
            token.write(creds.to_json())
    
    try:
        drive_service = build('drive', 'v3', credentials=creds)
        classroom_service = build('classroom', 'v1', credentials=creds)
        forms_service = build('forms', 'v1', credentials=creds)
        calendar_service = build('calendar', 'v3', credentials=creds)
        print("Google APIs initialized successfully with proper scopes.")
    except Exception as e:
        print(f"Error initializing Google APIs: {e}")

# Add this function to your WHATSAPP & MEMORY FUNCTIONS section
def send_start_template(to_number):
    """Sends the approved 'start' template message to a phone number."""
    url = f"https://graph.facebook.com/v19.0/{PHONE_NUMBER_ID}/messages"
    headers = {"Authorization": f"Bearer {ACCESS_TOKEN}", "Content-Type": "application/json"}
    payload = {
        "messaging_product": "whatsapp",
        "to": to_number,
        "type": "template",
        "template": {
            "name": "start_conversation_prompt", # Or your approved template name
            "language": {
                "code": "en"
            }
        }
    }
    try:
        response = requests.post(url, headers=headers, json=payload)
        response.raise_for_status()
        print(f"Successfully sent 'START' template to {to_number}.")
    except requests.exceptions.RequestException as e:
        print(f"Failed to send template to {to_number}: {e}")
        print(f"Response: {response.text}")

# ---------------------------------
# --- WHATSAPP & MEMORY FUNCTIONS
# ---------------------------------
#
def append_to_memory(user_id, role, content):
    """
    Appends a message to a user's conversation history in the correct
    format for the Gemini API.
    """
    # The Gemini API expects "user" and "model" as roles.
    # This maps your internal "assistant" role to "model".
    api_role = "model" if role == "assistant" else "user"
    
    # This uses the 'parts' key required by the API instead of 'content'.
    user_memory[user_id].append({"role": api_role, "parts": [str(content)]})
    
    if len(user_memory[user_id]) > MAX_HISTORY:
        user_memory[user_id].pop(0)

def send_whatsapp_message(to, message):
    url = f"https://graph.facebook.com/v19.0/{PHONE_NUMBER_ID}/messages"
    headers = {"Authorization": f"Bearer {ACCESS_TOKEN}", "Content-Type": "application/json"}
    payload = {"messaging_product": "whatsapp", "to": to, "text": {"body": message}}
    try:
        res = requests.post(url, headers=headers, json=payload)
        res.raise_for_status()
        print(f"WhatsApp text message sent to {to}: {message[:75]}...")
    except Exception as e:
        print(f"Failed to send WhatsApp text message: {e}")

def send_whatsapp_document(to, file_bytes, filename):
    media_url = f"https://graph.facebook.com/v19.0/{PHONE_NUMBER_ID}/media"
    headers = {"Authorization": f"Bearer {ACCESS_TOKEN}"}
    files_and_data = {
        'file': (filename, file_bytes, 'application/pdf'),
        'messaging_product': (None, 'whatsapp'), 'type': (None, 'document')
    }
    try:
        media_res = requests.post(media_url, headers=headers, files=files_and_data)
        media_res.raise_for_status()
        media_id = media_res.json().get("id")
        if not media_id:
            raise Exception("Failed to get media ID from upload response.")
    except Exception as e:
        print(f"Failed to upload PDF to WhatsApp: {e}")
        return
    
    message_url = f"https://graph.facebook.com/v19.0/{PHONE_NUMBER_ID}/messages"
    message_headers = {"Authorization": f"Bearer {ACCESS_TOKEN}", "Content-Type": "application/json"}
    message_payload = {"messaging_product": "whatsapp", "to": to, "type": "document", "document": {"id": media_id, "filename": filename}}
    try:
        msg_res = requests.post(message_url, headers=message_headers, json=message_payload)
        msg_res.raise_for_status()
        print(f"WhatsApp document sent to {to}: {filename}")
    except Exception as e:
        print(f"Failed to send WhatsApp document message: {e}")

def send_menu_message(to, text, options):
    """Sends an interactive List Message."""
    url = f"https://graph.facebook.com/v19.0/{PHONE_NUMBER_ID}/messages"
    headers = {"Authorization": f"Bearer {ACCESS_TOKEN}", "Content-Type": "application/json"}
    
    rows = [{"id": f"option_{i}", "title": option} for i, option in enumerate(options)]
    
    payload = {
        "messaging_product": "whatsapp",
        "to": to,
        "type": "interactive",
        "interactive": {
            "type": "list",
            "header": {"type": "text", "text": "Main Menu"},
            "body": {"text": text},
            "footer": {"text": "Please choose an option"},
            "action": {
                "button": "Choose an Option",
                "sections": [{"title": "Available Actions", "rows": rows}]
            }
        }
    }
    try:
        requests.post(url, headers=headers, data=json.dumps(payload)).raise_for_status()
        print(f"List Menu sent to {to}.")
    except Exception as e:
        print(f"Failed to send List Menu: {e}")

def send_whatsapp_video(to, video_bytes, filename):
    media_url = f"https://graph.facebook.com/v19.0/{PHONE_NUMBER_ID}/media"
    headers = {"Authorization": f"Bearer {ACCESS_TOKEN}"}
    files = {'file': (filename, video_bytes, 'video/mp4'), 'messaging_product': (None, 'whatsapp')}
    try:
        media_res = requests.post(media_url, headers=headers, files=files)
        media_res.raise_for_status()
        media_id = media_res.json().get("id")
        if not media_id: return "failure"
        
        msg_url = f"https://graph.facebook.com/v19.0/{PHONE_NUMBER_ID}/messages"
        msg_payload = {"messaging_product": "whatsapp", "to": to, "type": "video", "video": {"id": media_id}}
        msg_res = requests.post(msg_url, headers=headers, json=msg_payload)
        msg_res.raise_for_status()
        return "success"
    except Exception as e:
        print(f"Failed to send video: {e}")
        return "failure"
    
def handle_audio_task(sender, media_id):
    """Worker function to handle audio processing in a background thread."""
    print(f"Starting audio task for {sender}")
    send_whatsapp_message(sender, "Got your voice message, listening... 🎤")
    if media_url := get_media_url(media_id):
        if audio_bytes := download_media_file(media_url):
            if transcribed_text := transcribe_audio(audio_bytes):
                process_message(sender, transcribed_text)
            else:
                send_whatsapp_message(sender, "Sorry, I couldn't understand your audio.")
        else:
            send_whatsapp_message(sender, "Sorry, I couldn't download your voice message.")
    print(f"Audio task for {sender} finished.")
    
def handle_image_task(sender, media_id, prompt):
    """Worker function to handle image processing in a background thread."""
    print(f"Starting image task for {sender}")
    send_whatsapp_message(sender, "Analyzing your image... 🖼️")
    if media_url := get_media_url(media_id):
        if image_bytes := download_media_file(media_url):
            if 'podcast' in prompt.lower() or 'audio' in prompt.lower():
                ocr_prompt = "Extract all the text from this image. Do not summarize or explain, just provide the raw text."
                ocr_response = query_gemini_vision(ocr_prompt, image_bytes)
                
                if ocr_response.get("source") == "gemini_vision":
                    extracted_text = ocr_response["result"]
                    if audio_bytes := generate_voiceover(extracted_text):
                        send_whatsapp_audio(sender, audio_bytes, "podcast.mp3")
                    else:
                        send_whatsapp_message(sender, "I extracted the text, but couldn't create the audio.")
                else:
                    send_whatsapp_message(sender, ocr_response["result"])
            else:
                response = query_gemini_vision(prompt, image_bytes)
                send_whatsapp_message(sender, response["result"])
        else:
            send_whatsapp_message(sender, "Sorry, I couldn't download the image.")
    print(f"Image task for {sender} finished.")

def build_subject_vector_index():
    """Build a generic vector index from documents in the data folder"""
    print("Loading documents for subject vector index...")
    
    # Load all documents from the data folder dynamically
    data_folder = "data"
    documents = []
    
    if os.path.exists(data_folder):
        for filename in os.listdir(data_folder):
            file_path = os.path.join(data_folder, filename)
            if os.path.isfile(file_path):
                try:
                    loader = UnstructuredFileLoader(file_path)
                    raw_docs = loader.load()
                    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=100)
                    split_docs = text_splitter.split_documents(raw_docs)
                    documents.extend(split_docs)
                    print(f"Loaded and split {len(raw_docs)} documents from {filename}")
                except Exception as e:
                    print(f"Error loading {filename}: {e}")
    
    if not documents:
        print("No documents loaded, FAISS index will be empty.")
        # Create a minimal generic document
        generic_doc = Document(
            page_content="This is a knowledge base for various subjects. Upload documents to add more content.",
            metadata={'source': 'generic_base', 'type': 'base'}
        )
        documents = [generic_doc]
    
    embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
    db = FAISS.from_documents(documents, embeddings)
    db.save_local("vector_index/subject_docs")
    print("Generic Subject Vector Index built successfully.")


# ==============================================================================
# --- FIREBASE CONFIGURATION
# ==============================================================================

# Initialize Firebase
try:
    if not firebase_admin._apps:
        cred = credentials.Certificate('firebase_config.json')
        firebase_admin.initialize_app(cred)
    db = firestore.client()
    print("Firebase initialized successfully")
except Exception as e:
    print(f"Firebase initialization error: {e}")
    db = None

# ==============================================================================
# --- FIREBASE MEMORY MANAGEMENT
# ==============================================================================

def get_user_memory_from_firebase(user_id):
    """Retrieve user memory from Firebase"""
    if not db:
        return []
    
    try:
        doc_ref = db.collection('user_memories').document(user_id)
        doc = doc_ref.get()
        
        if doc.exists:
            data = doc.to_dict()
            return data.get('conversation_history', [])
        else:
            # Create new user document
            doc_ref.set({
                'user_id': user_id,
                'conversation_history': [],
                'created_at': datetime.now(),
                'updated_at': datetime.now()
            })
            return []
    except Exception as e:
        print(f"Error getting user memory from Firebase: {e}")
        return []

def save_user_memory_to_firebase(user_id, conversation_history):
    """Save user memory to Firebase"""
    if not db:
        return False
    
    try:
        doc_ref = db.collection('user_memories').document(user_id)
        doc_ref.set({
            'user_id': user_id,
            'conversation_history': conversation_history,
            'updated_at': datetime.now()
        }, merge=True)
        return True
    except Exception as e:
        print(f"Error saving user memory to Firebase: {e}")
        return False

def append_to_memory(user_id, role, content):
    """
    Appends a message to a user's conversation history and saves to Firebase
    """
    # The Gemini API expects "user" and "model" as roles
    api_role = "model" if role == "assistant" else "user"
    
    # Get existing memory from Firebase
    if user_id not in user_memory:
        user_memory[user_id] = get_user_memory_from_firebase(user_id)
    
    # Append new message
    user_memory[user_id].append({"role": api_role, "parts": [str(content)]})
    
    # Limit history length
    if len(user_memory[user_id]) > MAX_HISTORY:
        user_memory[user_id] = user_memory[user_id][-MAX_HISTORY:]
    
    # Save to Firebase
    save_user_memory_to_firebase(user_id, user_memory[user_id])

def get_user_state_from_firebase(user_id):
    """Retrieve user state from Firebase"""
    if not db:
        return "", {}
    
    try:
        doc_ref = db.collection('user_states').document(user_id)
        doc = doc_ref.get()
        
        if doc.exists:
            data = doc.to_dict()
            return data.get('current_state', ''), data.get('temp_data', {})
        else:
            # Create new user state document
            doc_ref.set({
                'user_id': user_id,
                'current_state': '',
                'temp_data': {},
                'created_at': datetime.now(),
                'updated_at': datetime.now()
            })
            return "", {}
    except Exception as e:
        print(f"Error getting user state from Firebase: {e}")
        return "", {}

def save_user_state_to_firebase(user_id, current_state, temp_data):
    """Save user state to Firebase"""
    if not db:
        return False
    
    try:
        doc_ref = db.collection('user_states').document(user_id)
        doc_ref.set({
            'user_id': user_id,
            'current_state': current_state,
            'temp_data': temp_data,
            'updated_at': datetime.now()
        }, merge=True)
        return True
    except Exception as e:
        print(f"Error saving user state to Firebase: {e}")
        return False

def clear_user_state_from_firebase(user_id):
    """Clear user state from Firebase"""
    if not db:
        return False
    
    try:
        doc_ref = db.collection('user_states').document(user_id)
        doc_ref.set({
            'user_id': user_id,
            'current_state': '',
            'temp_data': {},
            'updated_at': datetime.now()
        }, merge=True)
        return True
    except Exception as e:
        print(f"Error clearing user state from Firebase: {e}")
        return False

def save_uploaded_file_metadata(user_id, filename, file_size, chunks_count):
    """Save uploaded file metadata to Firebase"""
    if not db:
        return False
    
    try:
        file_id = str(uuid.uuid4())
        doc_ref = db.collection('uploaded_files').document(file_id)
        doc_ref.set({
            'file_id': file_id,
            'user_id': user_id,
            'filename': filename,
            'file_size': file_size,
            'chunks_count': chunks_count,
            'uploaded_at': datetime.now()
        })
        return True
    except Exception as e:
        print(f"Error saving file metadata to Firebase: {e}")
        return False

def get_user_uploaded_files(user_id):
    """Get list of files uploaded by a user that exist in BOTH Firebase AND vector DB"""
    if not db:
        return []
    
    try:
        # First get files from Firebase
        docs = db.collection('uploaded_files').where('user_id', '==', user_id).stream()
        firebase_files = []
        for doc in docs:
            file_data = doc.to_dict()
            firebase_files.append({
                'filename': file_data.get('filename'),
                'uploaded_at': file_data.get('uploaded_at'),
                'file_size': file_data.get('file_size'),
                'file_id': file_data.get('file_id')
            })
        
        # Now get files that actually exist in vector DB
        vector_db_files = get_uploaded_files_list()
        vector_db_filenames = [os.path.basename(f) for f in vector_db_files]
        
        # Only return files that exist in both Firebase AND vector DB
        valid_files = []
        for file in firebase_files:
            if file['filename'] in vector_db_filenames:
                valid_files.append(file)
        
        return valid_files
        
    except Exception as e:
        print(f"Error getting user uploaded files: {e}")
        return []
    

# Add these helper functions first
def check_syllabus_document(topic):
    """Check if syllabus document exists in DYNAMIC RAG only"""
    try:
        if not os.path.exists(DYNAMIC_VECTOR_INDEX_PATH) or not os.listdir(DYNAMIC_VECTOR_INDEX_PATH):
            return False, None
            
        syllabus_queries = [
            f"syllabus {topic}",
            f"course outline {topic}",
            f"curriculum {topic}",
            f"course structure {topic}"
        ]
        
        for query in syllabus_queries:
            # Use the dynamic-only query function
            result = query_dynamic_rag(query, k=2)
            if result["source"] == "dynamic_rag" and "no relevant information" not in result["result"].lower():
                return True, result["result"]
        return False, None
    except Exception as e:
        print(f"Error checking syllabus: {e}")
        return False, None

def check_chapter_notes_enhanced(topic, chapters=5):
    """PROPERLY check if chapter notes actually contain substantial content"""
    chapter_status = {}
    
    for chapter_num in range(1, chapters + 1):
        has_substantial_content = False
        content_length = 0
        
        chapter_queries = [
            f"chapter {chapter_num} {topic}",
            f"unit {chapter_num} {topic}",
            f"module {chapter_num} {topic}",
            f"part {chapter_num} {topic}",
            f"lecture {chapter_num} {topic}",
            f"notes {chapter_num} {topic}"
        ]
        
        best_content = ""
        for query in chapter_queries:
            result = query_dynamic_rag(query, k=3)
            if (result["source"] == "dynamic_rag" and 
                "no relevant information" not in result["result"].lower() and
                "no materials uploaded" not in result["result"].lower()):
                
                content = result["result"]
                # Check if content is substantial (not just error messages)
                if len(content) > 200 and not any(phrase in content.lower() for phrase in [
                    "sorry", "error", "could not", "no relevant", "upload materials"
                ]):
                    has_substantial_content = True
                    best_content = content
                    content_length = len(content)
                    break
        
        chapter_status[chapter_num] = {
            'has_content': has_substantial_content,
            'content': best_content,
            'content_length': content_length,
            'coverage_quality': 'good' if content_length > 500 else 'basic' if content_length > 200 else 'poor'
        }
    
    return chapter_status

def add_sample_question_papers():
    """Add sample question papers to guide difficulty levels"""
    sample_qps = {
        "easy_level": """
        SAMPLE EASY QUESTION PATTERN:
        - Define basic terms (2 marks)
        - List components/functions (2 marks)
        - Simple calculations (2 marks)
        - Explain basic concepts (16 marks)
        - Compare simple ideas (16 marks)
        """,
        
        "medium_level": """
        SAMPLE MEDIUM QUESTION PATTERN:
        - Explain with examples (2 marks)
        - Differentiate between concepts (2 marks)
        - Solve typical problems (2 marks)
        - Analyze scenarios (16 marks)
        - Design simple systems (16 marks)
        """,
        
        "hard_level": """
        SAMPLE DIFFICULT QUESTION PATTERN:
        - Critical analysis (2 marks)
        - Complex problem solving (2 marks)
        - Application in real-world (2 marks)
        - Evaluate and justify (16 marks)
        - Design complex systems (16 marks)
        """
    }
    
    # Add sample QPs to vector database
    for level, content in sample_qps.items():
        sample_doc = Document(
            page_content=content,
            metadata={
                'source': f'sample_qp_{level}.txt',
                'type': 'sample_qp',
                'difficulty': level,
                'uploaded_at': datetime.now().isoformat()
            }
        )
        
        # Add to vector index
        try:
            if os.path.exists(DYNAMIC_VECTOR_INDEX_PATH) and os.listdir(DYNAMIC_VECTOR_INDEX_PATH):
                vector_store = FAISS.load_local(DYNAMIC_VECTOR_INDEX_PATH, dynamic_embeddings, allow_dangerous_deserialization=True)
                vector_store.add_documents([sample_doc])
            else:
                vector_store = FAISS.from_documents([sample_doc], dynamic_embeddings)
            
            vector_store.save_local(DYNAMIC_VECTOR_INDEX_PATH)
            print(f"Added sample QP: {level}")
        except Exception as e:
            print(f"Error adding sample QP: {e}")

def generate_questions_with_difficulty(topic, question_type, count, marks, context, difficulty="medium"):
    """Generate questions with proper difficulty levels"""
    
    difficulty_prompts = {
        "easy": {
            "short": "Generate simple recall and basic understanding questions",
            "long": "Generate questions that test basic concept explanation"
        },
        "medium": {
            "short": "Generate application and analysis questions", 
            "long": "Generate questions that require problem-solving and examples"
        },
        "hard": {
            "short": "Generate evaluation and critical thinking questions",
            "long": "Generate questions that require design, evaluation, and justification"
        }
    }
    
    prompt = f"""
    Based EXACTLY on the following study material, generate {count} {question_type} answer questions worth {marks} marks each.
    
    DIFFICULTY LEVEL: {difficulty.upper()}
    {difficulty_prompts[difficulty][question_type]}
    
    STUDY MATERIAL:
    {context}
    
    REQUIREMENTS:
    - Questions MUST be based ONLY on the provided study material
    - For {difficulty} difficulty: {difficulty_prompts[difficulty][question_type]}
    - Each question should be clear and unambiguous
    - Format: "1. Question text"
    - Cover different aspects of the material
    
    Generate {count} {difficulty}-level questions:
    """
    
    try:
        model = genai.GenerativeModel('gemini-2.0-flash')
        response = model.generate_content(prompt)
        return response.text
    except Exception as e:
        print(f"Error generating {difficulty} questions: {e}")
        return None

def get_material_coverage_summary(chapter_status):
    """Get summary of material coverage"""
    chapters_with_content = [chap for chap, status in chapter_status.items() if status['has_content']]
    chapters_without_content = [chap for chap, status in chapter_status.items() if not status['has_content']]
    coverage_percentage = (len(chapters_with_content) / len(chapter_status)) * 100
    
    return {
        'total_chapters': len(chapter_status),
        'chapters_with_material': chapters_with_content,
        'chapters_without_material': chapters_without_content,
        'coverage_percentage': coverage_percentage,
        'sufficient_for_paper': len(chapters_with_content) >= 3
    }

# ==============================================================================
# --- UPDATED STATE MANAGEMENT WITH FIREBASE
# ==============================================================================

def initialize_user_session(sender):
    """Initialize user session with data from Firebase"""
    # Load memory from Firebase
    if sender not in user_memory:
        user_memory[sender] = get_user_memory_from_firebase(sender)
    
    # Load state from Firebase
    if sender not in user_states or sender not in user_temp_data:
        current_state, temp_data = get_user_state_from_firebase(sender)
        user_states[sender] = current_state
        user_temp_data[sender] = temp_data

def update_user_state(sender, state=None, temp_data=None):
    """Update user state and save to Firebase"""
    if state is not None:
        user_states[sender] = state
    
    if temp_data is not None:
        user_temp_data[sender].update(temp_data)
    
    # Save to Firebase
    save_user_state_to_firebase(sender, user_states.get(sender, ''), user_temp_data.get(sender, {}))

def clear_user_session(sender):
    """Clear user session and update Firebase"""
    user_states.pop(sender, None)
    user_temp_data.pop(sender, None)
    clear_user_state_from_firebase(sender)

# ---------------------------------
# --- VOICE INPUT HELPER FUNCTIONS
# ---------------------------------
#
def get_media_url(media_id):
    """Retrieves the downloadable URL for a piece of media from WhatsApp."""
    url = f"https://graph.facebook.com/v19.0/{media_id}/"
    headers = {"Authorization": f"Bearer {ACCESS_TOKEN}"}
    try:
        res = requests.get(url, headers=headers)
        res.raise_for_status()
        return res.json().get("url")
    except Exception as e:
        print(f"Failed to get media URL for ID {media_id}: {e}")
        return None
    
def handle_summary_task(sender, media_id):
    """Worker for the Image-to-Summary feature."""
    print(f"Starting summary task for {sender}")
    send_whatsapp_message(sender, "Processing your image for the summary...")
    if media_url := get_media_url(media_id):
        if image_bytes := download_media_file(media_url):
            ocr_prompt = "Extract all text from this image. Do not summarize or explain."
            ocr_response = query_gemini_vision(ocr_prompt, image_bytes)
            if ocr_response.get("source") == "gemini_vision":
                extracted_text = ocr_response["result"]
                summary_prompt = f"Please provide a concise summary of the following text: {extracted_text}"
                summary_response = query_gemini(summary_prompt, [])
                end_conversation_and_show_menu(sender, summary_response["result"])
            else:
                end_conversation_and_show_menu(sender, ocr_response["result"])
        else:
            end_conversation_and_show_menu(sender, "Sorry, I couldn't download the image.")
    print(f"Summary task for {sender} finished.")

def download_media_file(media_url):
    """Downloads the media file from the given URL."""
    headers = {"Authorization": f"Bearer {ACCESS_TOKEN}"}
    try:
        res = requests.get(media_url, headers=headers)
        res.raise_for_status()
        return res.content
    except Exception as e:
        print(f"Failed to download media from URL: {e}")
        return None

def transcribe_audio(audio_bytes):
    """Transcribes audio bytes using the self-hosted Whisper model."""
    print("Transcribing audio with local Whisper model...")
    
    temp_file_path = None
    try:
        # Create a temporary file, get its path, and then close it
        # so that FFmpeg can access it without a permission error.
        with tempfile.NamedTemporaryFile(suffix=".ogg", delete=False) as temp_file:
            temp_file.write(audio_bytes)
            temp_file_path = temp_file.name
        
        # Load the audio from the file path to get a NumPy array
        audio_np_array = whisper.load_audio(temp_file_path)

        # Transcribe the NumPy array
        result = whisper_model.transcribe(audio_np_array, fp16=False)
        
        transcript = result.get("text", "")
        print(f"Whisper transcription successful: '{transcript}'")
        return transcript

    except Exception as e:
        print(f"Error during local Whisper transcription: {e}")
        traceback.print_exc()
        return None
    finally:
        # Manually delete the temporary file after we're done with it
        if temp_file_path and os.path.exists(temp_file_path):
            os.remove(temp_file_path)

def end_conversation_and_show_menu(sender, final_message):
    """Sends a final message, clears user state safely, and shows the main menu."""
    if final_message:
        send_whatsapp_message(sender, final_message)
    
    clear_user_session(sender)

    menu_text = "What would you like to do next?"
    options = [
        "Ask Question", 
        "Create Worksheet", 
        "Create PPT", 
        "Generate Question Paper",
        "Upload Materials",  # Updated option
        "View Uploaded Files", 
        "Podcast from Image", 
        "Summary from Image"
    ]
    send_menu_message(sender, menu_text, options)
    append_to_memory(sender, "assistant", "Task complete. Displayed main menu.")

def check_material_coverage(topic):
    """Check if sufficient material is available for question paper generation"""
    try:
        chapter_materials = {}
        chapters_with_content = 0
        
        for chapter_num in range(1, 6):  # Check 5 chapters
            # Simple query for each chapter
            queries = [
                f"chapter {chapter_num}",
                f"unit {chapter_num}",
                f"module {chapter_num}",
                f"part {chapter_num}"
            ]
            
            has_content = False
            for query in queries:
                context = get_relevant_material_for_questions(f"{topic} {query}", num_chunks=1)
                if context and len(context) > 50 and "No relevant material" not in context:
                    has_content = True
                    break
            
            chapter_materials[chapter_num] = has_content
            if has_content:
                chapters_with_content += 1
        
        coverage_percentage = (chapters_with_content / 5) * 100
        
        return {
            'total_chapters': 5,
            'chapters_with_material': [chap for chap, has_mat in chapter_materials.items() if has_mat],
            'chapters_without_material': [chap for chap, has_mat in chapter_materials.items() if not has_mat],
            'coverage_percentage': coverage_percentage,
            'sufficient_for_paper': chapters_with_content >= 3  # Need at least 3 chapters
        }
        
    except Exception as e:
        print(f"Error checking material coverage: {e}")
        return {
            'total_chapters': 5,
            'chapters_with_material': [],
            'chapters_without_material': [1, 2, 3, 4, 5],
            'coverage_percentage': 0,
            'sufficient_for_paper': False
        }

def handle_bulk_document_upload(sender, media_items):
    """Handle multiple document uploads at once and send menu only after ALL processing"""
    print(f"Starting bulk upload for {sender}, {len(media_items)} files")
    
    successful_uploads = 0
    failed_uploads = 0
    processed_files = []
    
    send_whatsapp_message(sender, f"📦 Processing {len(media_items)} files... This may take a few minutes.")
    
    for i, media_item in enumerate(media_items, 1):
        media_id = media_item.get('id')
        filename = media_item.get('filename', f'document_{i}')
        
        send_whatsapp_message(sender, f"📄 Processing ({i}/{len(media_items)}): {filename}")
        
        if media_url := get_media_url(media_id):
            if file_bytes := download_media_file(media_url):
                success, message = add_documents_to_vector_index(file_bytes, filename)
                
                if success:
                    successful_uploads += 1
                    processed_files.append(filename)
                    # Save file metadata to Firebase if available
                    if db:
                        save_uploaded_file_metadata(sender, filename, len(file_bytes), message.split(' ')[-2] if 'chunks' in message else 0)
                else:
                    failed_uploads += 1
                    send_whatsapp_message(sender, f"❌ Failed to process {filename}: {message}")
            else:
                failed_uploads += 1
                send_whatsapp_message(sender, f"❌ Couldn't download {filename}")
        else:
            failed_uploads += 1
            send_whatsapp_message(sender, f"❌ Couldn't access {filename}")
    
    # ✅ Send summary ONLY after ALL files are processed
    summary_msg = f"✅ Bulk Upload Complete!\n\n"
    summary_msg += f"📊 Success: {successful_uploads} files\n"
    if failed_uploads > 0:
        summary_msg += f"❌ Failed: {failed_uploads} files\n"
    
    if successful_uploads > 0:
        summary_msg += f"\n📚 Added to knowledge base:\n"
        for file in processed_files[:5]:  # Show first 5 files
            summary_msg += f"• {file}\n"
        if len(processed_files) > 5:
            summary_msg += f"• ... and {len(processed_files) - 5} more\n"
        
        summary_msg += f"\nYou can now generate question papers using this content!"
    
    send_whatsapp_message(sender, summary_msg)
    
    # ✅ Clear state and show menu ONLY AFTER all processing is done
    if sender in user_states:
        del user_states[sender]
    end_conversation_and_show_menu(sender, None)
    
    print(f"Bulk upload for {sender} finished. Success: {successful_uploads}, Failed: {failed_uploads}")

# Update your QUESTION_PAPER_TEMPLATES
QUESTION_PAPER_TEMPLATES = {
    "semester": {
        "sections": [
            {
                "name": "PART A",
                "instructions": "Answer ALL questions (10 x 2 = 20 marks)",
                "questions": [
                    {"type": "short", "count": 10, "marks": 2, "coverage": "all_chapters"}
                ]
            },
            {
                "name": "PART B", 
                "instructions": "Answer ANY FIVE questions (5 x 16 = 80 marks)\n(Either / Or Choice from each unit)",
                "questions": [
                    {"type": "long", "count": 5, "marks": 16, "coverage": "per_chapter", "choice": "either_or"}
                ]
            }
        ],
        "total_marks": 100,
        "duration": "3 hours",
        "header": {
            "college_name": "YOUR COLLEGE NAME",
            "department": "Computer Science & Engineering",
            "exam_type": "B.E. / B.Tech. Degree Examination",
            "course_code": "CSXXX"
        },
        "chapters": 5  # Number of chapters/units
    },
    "cat": {
        "coming_soon": True,
        "message": "CAT question paper generation coming soon!"
    }
}

def get_relevant_material_for_questions(topic, num_chunks=5):
    """Get the most relevant material from RAG for question generation"""
    try:
        # Query both static and dynamic RAG
        static_embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
        all_results = []
        
        # Query static database
        try:
            static_db = FAISS.load_local("vector_index/subject_docs", static_embeddings, allow_dangerous_deserialization=True)
            static_retriever = static_db.as_retriever(search_kwargs={"k": num_chunks})
            static_docs = static_retriever.invoke(topic)
            all_results.extend(static_docs)
        except: pass
        
        # Query dynamic uploads
        try:
            if os.path.exists(DYNAMIC_VECTOR_INDEX_PATH) and os.listdir(DYNAMIC_VECTOR_INDEX_PATH):
                dynamic_db = FAISS.load_local(DYNAMIC_VECTOR_INDEX_PATH, dynamic_embeddings, allow_dangerous_deserialization=True)
                dynamic_retriever = dynamic_db.as_retriever(search_kwargs={"k": num_chunks})
                dynamic_docs = dynamic_retriever.invoke(topic)
                all_results.extend(dynamic_docs)
        except: pass
        
        if not all_results:
            return "No relevant material found for question generation."
        
        # Combine context
        context = "\n\n".join([doc.page_content for doc in all_results])
        return context
        
    except Exception as e:
        print(f"Error getting material for questions: {e}")
        return "Error retrieving material."

def generate_questions_from_material(topic, question_type, count, marks_per_question, context_chunks):
    """Generate questions based on specific material from RAG"""
    
    if question_type == "short":
        prompt = f"""Based EXACTLY on the following context material, generate {count} short answer questions worth {marks_per_question} marks each.

CONTEXT MATERIAL:
{context_chunks}

Requirements:
- Each question should be directly based on the provided context
- Questions should test basic understanding and recall
- Each question should be answerable in 2-3 sentences
- Format: "1. Question text"
- Cover different aspects of the material

Generate {count} questions:"""
    
    else:  # long questions
        prompt = f"""Based EXACTLY on the following context material, generate {count} long answer questions worth {marks_per_question} marks each.

CONTEXT MATERIAL:
{context_chunks}

Requirements:
- Each question should require analytical thinking based on the context
- Questions should test application, analysis, or evaluation
- Each question should be answerable in about 1 page
- Format: "1. Question text"
- Cover major concepts from the material

Generate {count} questions:"""
    
    try:
        model = genai.GenerativeModel('gemini-2.0-flash')
        response = model.generate_content(prompt)
        return response.text
    except Exception as e:
        print(f"Error generating questions: {e}")
        return None

def get_chapter_wise_material(topic, total_chapters=5):
    """Get material organized by chapters/units"""
    chapter_materials = {}
    
    for chapter_num in range(1, total_chapters + 1):
        # Search for chapter-specific content
        chapter_queries = [
            f"chapter {chapter_num} {topic}",
            f"unit {chapter_num} {topic}", 
            f"module {chapter_num} {topic}",
            f"part {chapter_num} {topic}"
        ]
        
        best_material = ""
        for query in chapter_queries:
            context = get_relevant_material_for_questions(query, num_chunks=3)
            if context and len(context) > len(best_material):
                best_material = context
        
        chapter_materials[chapter_num] = best_material
    
    return chapter_materials

def generate_semester_questions(topic, chapter_materials):
    """Generate semester pattern questions with either/or choices"""
    
    all_questions = {
        "part_a": [],  # 2-mark questions
        "part_b": []   # 16-mark questions with choices
    }
    
    # Generate PART A questions (2 marks each from all chapters)
    print("Generating PART A questions...")
    part_a_context = "\n\n".join([f"Chapter {chap}: {mat}" for chap, mat in chapter_materials.items() if mat])
    
    part_a_questions = generate_questions_from_material(
        topic=topic,
        question_type="short",
        count=10,
        marks_per_question=2,
        context_chunks=part_a_context
    )
    
    if part_a_questions:
        # Parse and format questions
        questions = []
        for line in part_a_questions.split('\n'):
            if re.match(r'^\d+\.', line.strip()):
                questions.append(line.strip())
        all_questions["part_a"] = questions[:10]  # Take first 10
    
    # Generate PART B questions (16 marks with either/or choices per chapter)
    print("Generating PART B questions...")
    for chapter_num, material in chapter_materials.items():
        if material:  # Only if we have material for this chapter
            # Generate two questions for either/or choice
            long_questions = generate_questions_from_material(
                topic=f"{topic} Chapter {chapter_num}",
                question_type="long", 
                count=2,
                marks_per_question=16,
                context_chunks=material
            )
            
            if long_questions:
                questions = []
                for line in long_questions.split('\n'):
                    if re.match(r'^\d+\.', line.strip()):
                        questions.append(line.strip())
                
                if len(questions) >= 2:
                    # Format as either/or
                    either_or_question = {
                        'chapter': chapter_num,
                        'option_a': questions[0],
                        'option_b': questions[1]
                    }
                    all_questions["part_b"].append(either_or_question)
    
    return all_questions

def generate_enhanced_semester_question_paper(topic):
    """Generate high-quality semester question paper with proper validation"""
    
    # Step 1: Validate materials properly
    send_whatsapp_message("SYSTEM", "🔍 Validating uploaded materials...")
    
    # Check syllabus
    has_syllabus, syllabus_content = check_syllabus_document(topic)
    if not has_syllabus:
        return None, "❌ No syllabus document found. Please upload the syllabus first."
    
    # Check chapter notes PROPERLY
    chapter_status = check_chapter_notes_enhanced(topic)
    coverage = get_material_coverage_summary(chapter_status)
    
    # Require at least 3 chapters with substantial content
    substantial_chapters = [chap for chap, status in chapter_status.items() 
                          if status['coverage_quality'] in ['good', 'basic']]
    
    if len(substantial_chapters) < 3:
        missing_chapters = [chap for chap in range(1, 6) if chap not in substantial_chapters]
        return None, f"❌ Insufficient chapter notes. Need substantial content for at least 3 chapters. Missing: {', '.join(map(str, missing_chapters))}"
    
    send_whatsapp_message("SYSTEM", f"✅ Found materials for {len(substantial_chapters)}/5 chapters. Generating question paper...")
    
    # Step 2: Generate questions with mixed difficulty
    all_questions = {
        "part_a": [],  # 2-mark questions
        "part_b": []   # 16-mark questions
    }
    
    # PART A: Mixed difficulty (2 marks each)
    part_a_context = "\n\n".join([
        f"Chapter {chap}: {chapter_status[chap]['content'][:1000]}" 
        for chap in substantial_chapters
    ])
    
    # Generate easy, medium, hard questions for Part A
    difficulties = ["easy", "medium", "hard"]
    questions_per_difficulty = max(3, 10 // len(difficulties))
    
    for difficulty in difficulties:
        questions_text = generate_questions_with_difficulty(
            topic=topic,
            question_type="short",
            count=questions_per_difficulty,
            marks=2,
            context=part_a_context,
            difficulty=difficulty
        )
        
        if questions_text:
            questions = [q.strip() for q in questions_text.split('\n') if re.match(r'^\d+\.', q.strip())]
            all_questions["part_a"].extend(questions[:questions_per_difficulty])
    
    # PART B: Chapter-wise questions (16 marks each)
    for chapter_num in substantial_chapters[:5]:  # Use first 5 substantial chapters
        chapter_content = chapter_status[chapter_num]['content']
        
        # Generate two questions for either/or choice
        long_questions = generate_questions_with_difficulty(
            topic=f"{topic} Chapter {chapter_num}",
            question_type="long",
            count=2,
            marks=16,
            context=chapter_content,
            difficulty="medium"  # Part B should be medium difficulty
        )
        
        if long_questions:
            questions = [q.strip() for q in long_questions.split('\n') if re.match(r'^\d+\.', q.strip())]
            if len(questions) >= 2:
                all_questions["part_b"].append({
                    'chapter': chapter_num,
                    'option_a': questions[0],
                    'option_b': questions[1]
                })
    
    # Step 3: Build question paper
    return build_question_paper_content(topic, all_questions, substantial_chapters), "success"

def build_question_paper_content(topic, questions, available_chapters):
    """Build the final question paper content"""
    
    content_parts = [
        "YOUR COLLEGE NAME",
        "Department of Computer Science & Engineering",
        "B.E. / B.Tech. Degree Examination",
        f"Course: CSXXX - {topic.title()}",
        "Time: 3 hours | Maximum Marks: 100",
        "=" * 60,
        f"NOTE: Generated from uploaded materials (Chapters: {', '.join(map(str, available_chapters))})",
        "=" * 60
    ]
    
    # PART A
    content_parts.extend([
        "\nPART A",
        "Answer ALL questions (10 x 2 = 20 marks)",
        "-" * 40
    ])
    
    if questions["part_a"]:
        for i, question in enumerate(questions["part_a"][:10], 1):
            # Clean question text
            clean_question = re.sub(r'^\d+\.\s*', '', question)
            content_parts.append(f"{i}. {clean_question} [2 marks]")
    else:
        for i in range(1, 11):
            content_parts.append(f"{i}. [Question based on uploaded materials] [2 marks]")
    
    # PART B
    content_parts.extend([
        "\nPART B",
        "Answer ANY FIVE questions (5 x 16 = 80 marks)",
        "(Either / Or Choice from each unit)",
        "-" * 40
    ])
    
    for i, q_pair in enumerate(questions["part_b"][:5], 1):
        content_parts.extend([
            f"{i}. a) {q_pair['option_a']}",
            f"   b) {q_pair['option_b']}",
            ""
        ])
    
    # Add remaining chapters if needed
    remaining_slots = 5 - len(questions["part_b"])
    for i in range(remaining_slots):
        chapter_num = len(questions["part_b"]) + i + 1
        content_parts.extend([
            f"{chapter_num}. a) [Upload notes for Chapter {chapter_num}]",
            f"   b) [Upload notes for Chapter {chapter_num}]", 
            ""
        ])
    
    content_parts.extend([
        "=" * 60,
        "Note: All questions carry equal marks as indicated.",
        f"Generated from uploaded materials covering {len(available_chapters)} chapters.",
        "Total number of printed pages: 2"
    ])
    
    return "\n".join(content_parts)



#
# ---------------------------------
# --- RAG & VECTOR DB FUNCTIONS
# ---------------------------------
#
def build_os_vector_index():
    print("Loading documents for OS vector index...")
    loaders = [
        UnstructuredFileLoader("data/OS_Syllabus.pdf"),
        UnstructuredFileLoader("data/OS_Unit_I_Overview.pptx"),
        UnstructuredFileLoader("data/OS_Unit_II_Process_Management.pptx"),
        UnstructuredFileLoader("data/OS_Unit_II_Process_Scheduling.pptx"),
        UnstructuredFileLoader("data/OS_Unit_III_Deadlocks.pptx"),
        UnstructuredFileLoader("data/OS_Unit_IV_Memory_management.pptx"),
        UnstructuredFileLoader("data/OS_Unit_IV_Virtual_Memory.pptx"),
        UnstructuredFileLoader("data/OS_Unit_V_File_System_Implementation.pptx"),
        UnstructuredFileLoader("data/OS_Unit_V_Storage_management.pptx"),
    ]
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=100)
    documents = []
    for loader in loaders:
        try:
            raw_docs = loader.load()
            split_docs = text_splitter.split_documents(raw_docs)
            documents.extend(split_docs)
            print(f"Loaded and split {len(raw_docs)} documents from {loader.file_path}")
        except Exception as e:
            print(f"Error loading {loader.file_path}: {e}")
    if not documents:
        print("No documents loaded, FAISS index will be empty.")
        return
    embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
    db = FAISS.from_documents(documents, embeddings)
    db.save_local("vector_index/os_docs")
    print("OS Vector Index built successfully.")

def query_subject_knowledge(question):
    """Generic query function that searches both base knowledge and uploaded materials"""
    return query_dynamic_rag(question)


def upload_file_to_drive(file_path, folder_id=None):
    """Uploads a file to Google Drive, optionally to a specific folder"""
    if not drive_service:
        print("Google Drive API not initialized.")
        return None
    try:
        file_metadata = {'name': os.path.basename(file_path)}
        
        # Only add parent folder if it exists and is valid
        if folder_id and folder_id.strip():
            try:
                # Verify the folder exists before trying to use it
                drive_service.files().get(fileId=folder_id).execute()
                file_metadata['parents'] = [folder_id]
                print(f"Uploading to folder: {folder_id}")
            except HttpError as e:
                print(f"Folder {folder_id} not found, uploading to root directory")
                # Don't add parents if folder doesn't exist
        
        media = MediaFileUpload(file_path, mimetype='application/pdf')
        file = drive_service.files().create(body=file_metadata, media_body=media, fields='id').execute()
        print(f"File uploaded to Drive with ID: {file.get('id')}")
        return file.get('id')
    except Exception as e:
        print(f"Error uploading file to Drive: {e}")
        return None

    
def post_announcement(course_id, text, drive_file_id):
    if not classroom_service:
        return {"result": "Google Classroom API not initialized.", "source": "api_error"}
    
    announcement = {
        'text': text,
        'materials': [{'driveFile': {'driveFile': {'id': drive_file_id}, 'shareMode': 'VIEW'}}],
        'state': 'PUBLISHED'
    }
    try:
        announcement = classroom_service.courses().announcements().create(courseId=course_id, body=announcement).execute()
        return {"result": f"Successfully posted announcement in Google Classroom!", "source": "classroom_success"}
    except HttpError as api_error:
        print(f"Google Classroom API Error: {api_error.content}")
        return {"result": "Couldn't post announcement. Please check permissions and Course ID.", "source": "api_error"}

def post_assignment(course_id, title, description, drive_file_id, due_date=None):
    """Create a classroom assignment without due dates"""
    if not classroom_service:
        return {"result": "Google Classroom API not initialized.", "source": "api_error"}
    
    coursework = {
        'title': title,
        'description': description,
        'materials': [{'driveFile': {'driveFile': {'id': drive_file_id}, 'shareMode': 'VIEW'}}],
        'workType': 'ASSIGNMENT',
        'state': 'PUBLISHED'
    }
    
    # Skip due dates entirely as requested
    try:
        assignment = classroom_service.courses().courseWork().create(courseId=course_id, body=coursework).execute()
        return {"result": f"Successfully created assignment '{assignment.get('title')}' in Google Classroom!", "source": "classroom_success"}
    except HttpError as api_error:
        print(f"Google Classroom API Error: {api_error.content}")
        return {"result": "Couldn't create Classroom assignment. Please check permissions and Course ID.", "source": "api_error"}
    except Exception as e:
        print(f"General error creating Classroom assignment: {e}")
        return {"result": "An unexpected error occurred while creating the Classroom assignment.", "source": "error"}
def query_gemini(question, history):
    try:
        model = genai.GenerativeModel('gemini-2.0-flash')
        chat = model.start_chat(history=history)
        response = chat.send_message(question)
        return {"result": response.text}
    except Exception as e:
        print("Gemini general query error:", e)
        return {"result": "Sorry, I encountered an issue while processing your general question."}
    
def generate_video_script(topic, history):
    print(f"Generating video script for topic: {topic}")
    prompt = f"Create a short, simple video script explaining the basics of {topic} for a beginner. The script should have 3 key points. For each key point, provide a narration sentence and a simple visual description on a new line starting with 'VISUAL:'."
    try:
        model = genai.GenerativeModel('gemini-2.0-flash')
        chat = model.start_chat(history=history)
        response = chat.send_message(prompt)
        return response.text
    except Exception as e:
        print(f"Error generating video script: {e}")
        return None

def generate_images_for_script(script):
    print("Generating images for script...")
    image_paths = []
    visual_prompts = re.findall(r'VISUAL:\s*(.*)', script)
    for i, prompt in enumerate(visual_prompts):
        try:
            img = Image.new('RGB', (1280, 720), color = 'darkblue')
            d = ImageDraw.Draw(img)
            try:
                font = ImageFont.truetype("arial.ttf", 40)
            except IOError:
                font = ImageFont.load_default()
            
            # Simple text wrapping
            lines = []
            words = prompt.split()
            current_line = ""
            for word in words:
                if len(current_line + " " + word) < 50:
                    current_line += " " + word
                else:
                    lines.append(current_line)
                    current_line = word
            lines.append(current_line)

            y_text = 300
            for line in lines:
                d.text((100, y_text), line.strip(), font=font, fill=(255,255,255))
                y_text += 50

            img_path = f"temp_image_{i}.png"
            img.save(img_path)
            image_paths.append(img_path)
        except Exception as e:
            print(f"Error creating image {i}: {e}")
    return image_paths

def generate_video_with_modelscope(prompt, output_path):
    """Generates a video from a text prompt using the ModelScope model."""
    print(f"Generating video for prompt: '{prompt}'...")
    try:
        video_frames = video_pipe(prompt, num_inference_steps=25, num_frames=16).frames
        export_to_video(video_frames, output_path)
        return True
    except Exception as e:
        print(f"Error generating video with ModelScope: {e}")
        return False
    
def generate_voiceover(text):
    """Generates a voiceover from text and speeds it up."""
    print("Generating voiceover...")
    try:
        # Generate the initial audio
        tts = gTTS(text=text, lang='en', slow=False)
        fp = BytesIO()
        tts.write_to_fp(fp)
        fp.seek(0)
        
        # Load the audio with pydub and speed it up
        audio = AudioSegment.from_file(fp, format="mp3")
        sped_up_audio = audio.speedup(playback_speed=1.25)
        
        # Export the sped-up audio to a new BytesIO object
        output_fp = BytesIO()
        sped_up_audio.export(output_fp, format="mp3")
        output_fp.seek(0)
        
        return output_fp.read()
    except Exception as e:
        print(f"Error generating voiceover: {e}")
        return None
    
def create_calendar_event(title, due_date):
    if not calendar_service: return
    try:
        event = {
            'summary': f'Due: {title}',
            'description': 'An assignment is due.',
            'start': {'date': due_date.isoformat()},
            'end': {'date': due_date.isoformat()},
        }
        calendar_service.events().insert(calendarId='primary', body=event).execute()
        print("Calendar event created.")
    except Exception as e:
        print(f"Error creating calendar event: {e}")

#
# ---------------------------------
# --- CONTENT GENERATION FUNCTIONS
# ---------------------------------
#
def generate_mcq_questions_text(topic, num_questions, history):
    print(f"Attempting to generate {num_questions} TEXT MCQs on topic: '{topic}'")
    context, source_info = "", ""
    # Simplified RAG check for brevity
    topic_embedding = embeddings_model_for_classification.embed_query(topic)
    topic_embedding_reshaped = np.array(topic_embedding).reshape(1, -1)
    similarities = cosine_similarity(topic_embedding_reshaped, os_example_embeddings)
    if np.max(similarities) > SIMILARITY_THRESHOLD_OS:
        try:
            embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
            db = FAISS.load_local("vector_index/os_docs", embeddings, allow_dangerous_deserialization=True)
            retriever = db.as_retriever(search_kwargs={"k": 3})
            docs = retriever.invoke(topic)
            if docs:
                context = "\n\n".join([doc.page_content for doc in docs])
                unique_sources = list(set([os.path.basename(d.metadata['source']) for d in docs if 'source' in d.metadata]))
                if unique_sources:
                    source_info = f"\n\n(Context from: {', '.join(unique_sources)})"
        except Exception as e:
            print(f"Error retrieving RAG context for MCQ: {e}")
    mcq_prompt = f"Generate {num_questions} ONLY multiple-choice questions on the topic of '{topic}'. For each question, provide 4 options (A, B, C, D) and clearly indicate the correct answer. Format each as: 'Question Number. Question Text\\n A) Option A\\n B) Option B\\n C) Option C\\n D) Option D\\n Correct Answer: X'. Provide ONLY the questions and answers."
    if context:
        mcq_prompt = f"Based on the following context:\n\n{context}\n\n" + mcq_prompt
    try:
        model = genai.GenerativeModel('gemini-2.0-flash')
        chat = model.start_chat(history=history)
        response = chat.send_message(mcq_prompt)
        return {"result": response.text + source_info, "source": "generated_form_questions_text"}
    except Exception as e:
        print(f"Error generating TEXT MCQs for '{topic}': {e}")
        return {"result": "Sorry, I couldn't generate the quiz questions.", "source": "error"}

def create_google_form_mcq(title, mcq_questions_text):
    """Creates a Google Form quiz from a string of MCQ questions."""
    if not forms_service:
        return {"result": "Google Forms API not initialized.", "source": "api_error"}

    try:
        form_metadata = {'info': {'title': title}}
        form = forms_service.forms().create(body=form_metadata).execute()
        form_id = form['formId']
        
        update = {"requests": [{"updateSettings": {"settings": {"quizSettings": {'isQuiz': True}}, "updateMask": "quizSettings"}}]}
        forms_service.forms().batchUpdate(formId=form_id, body=update).execute()

        question_requests = []
        # A more robust way to split into question blocks
        question_blocks = re.split(r'\n\s*\d+\.\s*', mcq_questions_text)
        
        # Filter out any empty strings or introductory text
        question_blocks = [b.strip() for b in question_blocks if b.strip() and "Instructions:" not in b and "Worksheet" not in b]

        for i, q_block in enumerate(question_blocks):
            lines = [line.strip() for line in q_block.split('\n') if line.strip()]
            if not lines:
                continue

            question_text = lines[0]
            options = []
            correct_answer_text = ""

            for line in lines[1:]:
                if re.match(r'^[a-dA-D][\.\)]', line):
                    options.append(line[2:].strip())
                elif "Correct Answer:" in line:
                    correct_answer_match = re.search(r'Correct Answer:\s*([a-dA-D])', line, re.IGNORECASE)
                    if correct_answer_match:
                        correct_answer_letter = correct_answer_match.group(1).upper()
                        correct_answer_index = ord(correct_answer_letter) - ord('A')
                        if 0 <= correct_answer_index < len(options):
                            correct_answer_text = options[correct_answer_index]

            if not question_text or len(options) < 2 or not correct_answer_text:
                print(f"Skipping malformed block: {q_block[:50]}...")
                continue

            choices = [{"value": opt} for opt in options]
            
            question_requests.append({
                "createItem": {
                    "item": {
                        "title": question_text,
                        "questionItem": {
                            "question": {
                                "required": True,
                                "choiceQuestion": {"type": "RADIO", "options": choices},
                                "grading": {"pointValue": 1, "correctAnswers": {"answers": [{"value": correct_answer_text}]}}
                            }
                        }
                    }, "location": {"index": i}
                }
            })

        if question_requests:
            forms_service.forms().batchUpdate(formId=form_id, body={"requests": question_requests}).execute()

        return {"result": f"I've created a Google Form quiz for you! Access it here: {form['responderUri']}", "source": "google_form_created"}

    except Exception as e:
        print(f"Error creating Google Form: {e}")
        traceback.print_exc()
        return {"result": "Sorry, I encountered an error while creating the Google Form.", "source": "error"}

def generate_worksheet_content_text(topic, num_items, worksheet_type, history):
    """
    Generates worksheet questions AND answers, separated by a unique string.
    """
    print(f"Generating {num_items} {worksheet_type} items and answers on: '{topic}'")
    
    # New prompt that asks for both questions and a clearly separated answer key
    prompt = (
        f"Generate a student worksheet with {num_items} {worksheet_type} questions on the topic of '{topic}'. "
        "The questions should be clear and suitable for a student. Do not provide the answers immediately after each question. "
        "After the questions, add a separator line exactly like this: '--- ANSWERS ---'. "
        "Then, after the separator, provide a numbered list of the corresponding answers. "
        "The final output must contain both the questions and the answers, separated by the '--- ANSWERS ---' line."
    )

    try:
        model = genai.GenerativeModel('gemini-2.0-flash')
        chat = model.start_chat(history=history)
        response = chat.send_message(prompt)
        return {"result": response.text, "source": "generated_worksheet_text"}
    except Exception as e:
        print(f"Error generating worksheet content: {e}")
        return {"result": "Sorry, I couldn't generate the worksheet content.", "source": "error"}
    
def generate_ppt_content(topic):
    prompt = f"Create the content for a 10-slide presentation on the topic of '{topic}'. The first slide should be a title slide. The next 9 slides should be content slides, each with a title and 3-4 bullet points. Format the output clearly, using 'SLIDE:' to mark the beginning of each slide."
    try:
        model = genai.GenerativeModel('gemini-2.0-flash')
        response = model.generate_content(prompt)
        return response.text
    except Exception as e:
        print(f"Error generating PPT content: {e}")
        return None
    
def create_ppt_file(title, content):
    try:
        prs = Presentation()
        # More robustly split the content into slides
        slides = re.split(r'\n*SLIDE:?\s*\d*\s*:?\n*', content)
        slides = [s.strip() for s in slides if s.strip()] # Clean up empty entries

        if not slides:
            print("Could not parse any slides from the generated content.")
            return None

        # Handle the first slide as the title slide
        title_slide_content = slides[0].split('\n')
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title_slide_content[0]
        if len(title_slide_content) > 1:
            slide.placeholders[1].text = "\n".join(title_slide_content[1:])
        
        # Handle the rest as content slides
        for slide_content in slides[1:]:
            lines = [line.strip() for line in slide_content.strip().split('\n') if line.strip()]
            if not lines: continue

            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = lines[0] # Assume first line is the title
            
            body_shape = slide.shapes.placeholders[1]
            tf = body_shape.text_frame
            tf.clear()
            
            for point in lines[1:]:
                # Remove leading bullet points like '*' or '-'
                cleaned_point = re.sub(r'^\s*[\*\-]\s*', '', point)
                p = tf.add_paragraph()
                p.text = cleaned_point
                p.level = 0

        fp = BytesIO()
        prs.save(fp)
        fp.seek(0)
        return fp.read()
    except Exception as e:
        print(f"Error creating PPT file: {e}")
        traceback.print_exc()
        return None

def send_whatsapp_ppt(to, file_bytes, filename):
    media_url = f"https://graph.facebook.com/v19.0/{PHONE_NUMBER_ID}/media"
    headers = {"Authorization": f"Bearer {ACCESS_TOKEN}"}
    files = {'file': (filename, file_bytes, 'application/vnd.openxmlformats-officedocument.presentationml.presentation'), 'messaging_product': (None, 'whatsapp')}
    try:
        media_res = requests.post(media_url, headers=headers, files=files)
        media_res.raise_for_status()
        media_id = media_res.json().get("id")
        if not media_id: return "failure"
        
        msg_url = f"https://graph.facebook.com/v19.0/{PHONE_NUMBER_ID}/messages"
        msg_headers = {"Authorization": f"Bearer {ACCESS_TOKEN}", "Content-Type": "application/json"}
        msg_payload = {"messaging_product": "whatsapp", "to": to, "type": "document", "document": {"id": media_id, "filename": filename}}
        requests.post(msg_url, headers=msg_headers, json=msg_payload).raise_for_status()
        return "success"
    except Exception as e:
        print(f"Failed to send PPT: {e}")
        return "failure"

def create_pdf_locally(title, content):
    """
    Creates a PDF file from a title and content, saving it to a temporary file
    and returning the raw bytes.
    """
    pdf_bytes = None
    temp_pdf_path = None
    try:
        # Create a temporary file to save the PDF to disk first
        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as temp_file:
            temp_pdf_path = temp_file.name

        # Initialize the PDF document
        pdf = FPDF()
        pdf.add_page()
        pdf.set_auto_page_break(auto=True, margin=15)
        
        # Add the title
        pdf.set_font("Helvetica", 'B', 16)
        pdf.multi_cell(0, 10, title, align='C')
        pdf.ln(10)
        
        # Add the main content
        pdf.set_font("Helvetica", size=11)
        # Clean the content to handle special characters that FPDF might not support
        cleaned_content = content.encode('latin-1', 'replace').decode('latin-1')
        pdf.multi_cell(0, 7, cleaned_content)
        
        # Save the PDF to the temporary file path
        pdf.output(temp_pdf_path)

        # Read the bytes from the saved file to ensure it's not empty
        with open(temp_pdf_path, 'rb') as f:
            pdf_bytes = f.read()

        return pdf_bytes
        
    except Exception as e:
        print(f"Error creating PDF: {e}")
        return None
    finally:
        # Clean up the temporary file from disk after we're done with it
        if temp_pdf_path and os.path.exists(temp_pdf_path):
            os.remove(temp_pdf_path)

#
# ---------------------------------
# --- CORE MESSAGE PROCESSOR
# ---------------------------------
#

def handle_final_classroom_post(sender, title):
    """Helper function to handle the final steps of posting to Classroom."""
    try:
        post_choice = user_temp_data[sender]['post_choice']
        questions = user_temp_data[sender]['questions_text']
        answers = user_temp_data[sender]['answers_text']
        
        pdf_content = questions
        pdf_filename = f"{title.replace(' ', '_')}_worksheet.pdf"
        
        if post_choice.lower() == "post with answers":
            pdf_content += f"\n\n--- ANSWERS ---\n{answers}"
            pdf_filename = f"{title.replace(' ', '_')}_with_answers.pdf"

        pdf_bytes = create_pdf_locally(title, pdf_content)
        
        if pdf_bytes:
            with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as temp_file:
                temp_file.write(pdf_bytes)
                temp_file_path = temp_file.name
            
            # Try to upload to the specified folder, but fall back to root if needed
            drive_file_id = upload_file_to_drive(temp_file_path, GOOGLE_FORMS_TARGET_FOLDER_ID)
            os.remove(temp_file_path)

            if drive_file_id:
                if post_choice.lower() != "don't post":
                    # Create assignment with the PDF
                    description = f"Worksheet on {title}"
                    if post_choice.lower() == "post with answers":
                        description += " (Includes answer key)"
                    
                    # Try to post to classroom
                    result = post_assignment(CLASSROOM_COURSE_ID, title, description, drive_file_id, due_date=None)
                    
                    if result["source"] == "classroom_success":
                        send_whatsapp_message(sender, result['result'])
                    else:
                        # If classroom posting fails, provide the Drive link instead
                        drive_link = f"https://drive.google.com/file/d/{drive_file_id}/view"
                        send_whatsapp_message(sender, f"I couldn't post to Classroom due to permission issues. Here's the direct link to the file: {drive_link}")
                else:
                    send_whatsapp_message(sender, "Okay, I won't post to Classroom. The PDFs have been sent to you directly.")
            else:
                send_whatsapp_message(sender, "I created the PDF, but failed to upload it to Google Drive.")
        else:
            send_whatsapp_message(sender, "Sorry, I failed to create the final PDF for upload.")

    except KeyError:
        send_whatsapp_message(sender, "An error occurred. I've lost the worksheet details. Please start over.")
    except Exception as e:
        print(f"Unexpected error in handle_final_classroom_post: {e}")
        send_whatsapp_message(sender, "An unexpected error occurred. Please try again.")
    finally:
        if sender in user_states: del user_states[sender]
        if sender in user_temp_data: del user_temp_data[sender]
        
def query_gemini_vision(prompt, image_bytes):
    """Sends a prompt and an image to the Gemini Vision model."""
    print("Querying Gemini Vision...")
    try:
        image_parts = [{"mime_type": "image/jpeg", "data": base64.b64encode(image_bytes).decode('utf-8')}]
        prompt_parts = [prompt] + image_parts
        
        model = genai.GenerativeModel('gemini-2.0-flash')
        response = model.generate_content(prompt_parts)
        
        if response.parts:
            return {"result": response.text, "source": "gemini_vision"}
        else:
            return {"result": "I'm sorry, I couldn't understand the image.", "source": "error"}
            
    except Exception as e:
        print(f"Error querying Gemini Vision: {e}")
        return {"result": "Sorry, there was an error processing the image.", "source": "error"}
    
def send_whatsapp_audio(to, audio_bytes, filename):
    media_url = f"https://graph.facebook.com/v19.0/{PHONE_NUMBER_ID}/media"
    headers = {"Authorization": f"Bearer {ACCESS_TOKEN}"}
    files = {'file': (filename, audio_bytes, 'audio/mpeg'), 'messaging_product': (None, 'whatsapp')}
    try:
        media_res = requests.post(media_url, headers=headers, files=files)
        media_res.raise_for_status()
        media_id = media_res.json().get("id")
        if not media_id: return "failure"

        msg_url = f"https://graph.facebook.com/v19.0/{PHONE_NUMBER_ID}/messages"
        msg_payload = {"messaging_product": "whatsapp", "to": to, "type": "audio", "audio": {"id": media_id}}
        requests.post(msg_url, headers=headers, json=msg_payload).raise_for_status()
        return "success"
    except Exception as e:
        print(f"Failed to send audio: {e}")
        return "failure"
    
def handle_podcast_task(sender, media_id):
    """Worker for the Image-to-Podcast feature."""
    print(f"Starting podcast task for {sender}")
    send_whatsapp_message(sender, "Processing your image for the podcast...")
    if media_url := get_media_url(media_id):
        if image_bytes := download_media_file(media_url):
            ocr_prompt = "Extract all the text from this image. Do not summarize or explain, just provide the raw text."
            ocr_response = query_gemini_vision(ocr_prompt, image_bytes)
            if ocr_response.get("source") == "gemini_vision":
                if audio_bytes := generate_voiceover(ocr_response["result"]):
                    send_whatsapp_audio(sender, audio_bytes, "podcast.mp3")
                    end_conversation_and_show_menu(sender, "Here is your podcast.")
                else:
                    end_conversation_and_show_menu(sender, "I extracted the text, but couldn't create the audio.")
            else:
                end_conversation_and_show_menu(sender, ocr_response["result"])
        else:
            end_conversation_and_show_menu(sender, "Sorry, I couldn't download the image.")
    print(f"Podcast task for {sender} finished.")
    
def generate_voiceover(text):
    print("Generating voiceover...")
    try:
        tts = gTTS(text=text, lang='en')
        fp = BytesIO()
        tts.write_to_fp(fp)
        fp.seek(0)
        return fp.read()
    except Exception as e:
        print(f"Error generating voiceover: {e}")
        return None
    
def create_question_paper_pdf(topic, content):
    """Create PDF for the generated question paper"""
    try:
        pdf = FPDF()
        pdf.add_page()
        pdf.set_auto_page_break(auto=True, margin=15)
        
        # Title
        pdf.set_font("Helvetica", 'B', 16)
        pdf.multi_cell(0, 10, f"Question Paper - {topic}", align='C')
        pdf.ln(5)
        
        # Content
        pdf.set_font("Helvetica", size=11)
        
        # Split content into lines and add to PDF
        lines = content.split('\n')
        for line in lines:
            if line.strip() == "="*50:
                pdf.set_draw_color(0, 0, 0)
                pdf.line(10, pdf.get_y(), 200, pdf.get_y())
                pdf.ln(2)
            elif line.strip() == "-"*30:
                pdf.set_draw_color(150, 150, 150)
                pdf.line(10, pdf.get_y(), 200, pdf.get_y())
                pdf.ln(2)
            elif any(keyword in line.upper() for keyword in ['SECTION', 'COURSE:', 'TIME:']):
                pdf.set_font("Helvetica", 'B', 12)
                pdf.multi_cell(0, 8, line)
                pdf.set_font("Helvetica", size=11)
            elif line.strip().startswith(tuple(str(i) for i in range(10))):
                pdf.set_font("Helvetica", 'B', 11)
                pdf.multi_cell(0, 7, line)
                pdf.set_font("Helvetica", size=11)
            else:
                pdf.multi_cell(0, 6, line)
        
        # Save to bytes
        pdf_bytes = pdf.output(dest='S').encode('latin1')
        return pdf_bytes
        
    except Exception as e:
        print(f"Error creating question paper PDF: {e}")
        return None
    
def recheck_syllabus_after_upload(sender):
    """Re-check syllabus after user uploads a document"""
    if sender in user_temp_data and 'qp_topic' in user_temp_data[sender]:
        topic = user_temp_data[sender]['qp_topic']
        send_whatsapp_message(sender, "✅ Document uploaded! Re-checking syllabus...")
        
        has_syllabus, syllabus_content = check_syllabus_document(topic)
        if has_syllabus:
            # Move to chapter check
            user_states[sender] = "awaiting_question_paper_topic"
            process_message(sender, topic)  # This will trigger chapter check
        else:
            send_whatsapp_message(sender, 
                "Still no syllabus found. Please upload a proper syllabus document.\n"
                "I'll continue checking for more uploads..."
            )

def recheck_chapters_after_upload(sender):
    """Re-check chapters after user uploads documents"""
    if sender in user_temp_data and 'qp_topic' in user_temp_data[sender]:
        topic = user_temp_data[sender]['qp_topic']
        send_whatsapp_message(sender, "✅ Documents uploaded! Re-checking chapter coverage...")
        
        chapter_status = check_chapter_notes_enhanced(topic)
        coverage = get_material_coverage_summary(chapter_status)
        
        user_temp_data[sender]['chapter_status'] = chapter_status
        user_temp_data[sender]['coverage'] = coverage
        
        if coverage['sufficient_for_paper']:
            user_states[sender] = "ready_to_generate_qp"
            process_message(sender, "ready")  # Trigger ready state
        else:
            missing_chapters = ", ".join(map(str, coverage['chapters_without_material']))
            send_whatsapp_message(sender, 
                f"Still missing chapters: {missing_chapters}. "
                f"Please upload more notes. I'll continue checking for uploads..."
            )

def check_qp_syllabus(topic):
    """Check if syllabus exists in QP RAG"""
    try:
        result = query_qp_rag(f"syllabus {topic}", material_type="syllabus", k=2)
        return result["source"] == "qp_rag", result["result"]
    except Exception as e:
        print(f"Error checking QP syllabus: {e}")
        return False, None

def check_qp_chapter_notes(topic, chapters=5):
    """Check chapter notes in QP RAG"""
    chapter_status = {}
    
    for chapter_num in range(1, chapters + 1):
        has_content = False
        
        result = query_qp_rag(f"chapter {chapter_num} {topic}", material_type="chapter_notes", k=1)
        if result["source"] == "qp_rag":
            has_content = True
        
        chapter_status[chapter_num] = {'has_content': has_content}
    
    return chapter_status

def check_qp_sample_papers(topic):
    """Check if sample papers exist in QP RAG"""
    try:
        result = query_qp_rag(f"question paper sample {topic}", material_type="sample_papers", k=1)
        return result["source"] == "qp_rag", result["result"]
    except Exception as e:
        print(f"Error checking QP sample papers: {e}")
        return False, None
    
# Add these new states to your state management

def start_question_paper_flow(sender):
    """Start the dedicated SEMESTER question paper generation flow"""
    # ✅ Clear any existing state first
    if sender in user_states:
        del user_states[sender]
    if sender in user_temp_data:
        del user_temp_data[sender]
    
    # ✅ Initialize fresh state
    user_states[sender] = "awaiting_qp_title"
    user_temp_data[sender] = {
        'qp_materials': {},
        'current_step': 0,
        'total_steps': 4,
        'exam_type': 'semester'
    }
    
    send_whatsapp_message(sender,
        "🎓 SEMESTER QUESTION PAPER GENERATION\n\n"
        "I'll help you create a **semester exam** question paper (100 marks).\n\n"
        "📋 Required Materials:\n"
        "1. Syllabus Document\n"
        "2. Chapter Notes (All 5 chapters)\n" 
        "3. Sample Question Papers\n\n"
        "Let's start! What should be the title of your semester question paper?"
    )

def handle_qp_document_upload(sender, media_id, filename, material_type):
    """Handle document upload for question paper materials"""
    print(f"Processing {material_type} for QP: {filename}")
    
    if media_url := get_media_url(media_id):
        if file_bytes := download_media_file(media_url):
            success, message = add_to_qp_vector_index(file_bytes, filename, material_type)
            
            if success:
                # Store in user temp data
                if material_type not in user_temp_data[sender]['qp_materials']:
                    user_temp_data[sender]['qp_materials'][material_type] = []
                
                user_temp_data[sender]['qp_materials'][material_type].append({
                    'filename': filename,
                    'uploaded_at': datetime.now().isoformat()
                })
                
                return True, f"✅ {material_type.replace('_', ' ').title()} added: {filename}"
            else:
                return False, f"❌ Failed to process {filename}: {message}"
    
    return False, f"❌ Could not download {filename}"

def generate_enhanced_qp_from_materials(title):
    """Generate question paper using only QP RAG materials"""
    try:
        # Check all materials exist
        has_syllabus, syllabus_content = check_qp_syllabus(title)
        if not has_syllabus:
            return None, "Syllabus not found in uploaded materials."
        
        chapter_status = check_qp_chapter_notes(title)
        substantial_chapters = [chap for chap, status in chapter_status.items() if status['has_content']]
        
        if len(substantial_chapters) < 3:
            return None, f"Insufficient chapter notes. Need at least 3 chapters. Currently have: {len(substantial_chapters)}"
        
        has_samples, sample_content = check_qp_sample_papers(title)
        
        # Generate using the enhanced generator (from previous implementation)
        qp_content = build_enhanced_question_paper(title, substantial_chapters, has_samples)
        return qp_content, "success"
        
    except Exception as e:
        print(f"Error generating QP from materials: {e}")
        return None, f"Generation error: {str(e)}"
    
def debug_user_state(sender, action):
    """Enhanced debug user state transitions"""
    current_state = user_states.get(sender, "NO_STATE")
    temp_data = user_temp_data.get(sender, {})
    print(f"🔍 [{action}] User: {sender}, State: {current_state}, Temp Keys: {list(temp_data.keys())}")

def clear_user_session(sender):
    """Completely clear user session"""
    if sender in user_states:
        del user_states[sender]
    if sender in user_temp_data:
        del user_temp_data[sender]
    clear_user_state_from_firebase(sender)

def get_qp_progress(sender):
    """Get current progress in QP generation"""
    temp_data = user_temp_data.get(sender, {})
    qp_materials = temp_data.get('qp_materials', {})
    
    progress = {
        'syllabus': len(qp_materials.get('syllabus', [])) > 0,
        'chapters_uploaded': len(qp_materials.get('chapter_notes', [])),
        'sample_papers': len(qp_materials.get('sample_papers', [])) > 0,
        'current_step': temp_data.get('current_step', 0),
        'total_steps': temp_data.get('total_steps', 4)
    }
    
    return progress
    
def process_message(sender, text):
    # Initialize user session from Firebase
    initialize_user_session(sender)
    
    append_to_memory(sender, "user", text)
    current_state = user_states.get(sender)

    # --- Universal Cancel/Menu ---
    if text.lower().strip() in ['cancel', 'stop', 'menu', 'start', 'exit']:
        if current_state:
            clear_user_session(sender)
            send_whatsapp_message(sender, "Okay, I've canceled the current operation.")
        menu_text = "Hello! I'm your AI teaching assistant. Please choose an option:"
        options = ["Ask Question", "Create Worksheet", "Create PPT","Generate Question Paper", "Upload Material", "View Uploaded Files", "Podcast from Image", "Summary from Image", "Create Video"]
        send_menu_message(sender, menu_text, options)
        append_to_memory(sender, "assistant", "Displayed main menu.")
        return

    # --- Initial Triggers for New Conversations (if no active state) ---
    if not current_state:
        lower_text = text.lower().strip()

        

        if lower_text == "ask question":
            update_user_state(sender, state="awaiting_question")
            send_whatsapp_message(sender, "Of course! What is your question? I'll search through all uploaded documents and knowledge base.")
            return
        
        if lower_text == "create worksheet":
            update_user_state(sender, state="awaiting_worksheet_topic")
            send_whatsapp_message(sender, "Let's create a worksheet! What topic would you like the worksheet to be about?")
            return
            
        if lower_text == "create ppt":
            update_user_state(sender, state="awaiting_ppt_topic")
            send_whatsapp_message(sender, "Excellent! What topic would you like the presentation to be about?")
            return
            
        if lower_text == "upload materials" or lower_text == "upload material":
            update_user_state(sender, state="awaiting_material_file")
            send_whatsapp_message(sender, "Please send the file you'd like to upload (PDF, Word, PowerPoint, or Text). This will be added to my knowledge base for future questions!")
            return
            
        if lower_text == "generate question paper":
            # Clear any existing state completely
            clear_user_session(sender)
            
            # Start fresh exam type selection
            user_states[sender] = "awaiting_exam_type"
            user_temp_data[sender] = {}  # Initialize empty temp data
            
            send_interactive_message(
                sender,
                "🎓 SELECT EXAM TYPE\n\nWhich type of question paper would you like to generate?",
                ["Semester Exam (100 marks)", "CAT Exam (Coming Soon)", "Unit Test (Coming Soon)"]
            )
            return
        
        if lower_text == "view uploaded files":
            # Get files that ACTUALLY exist in vector DB
            vector_db_files = get_uploaded_files_list()
            
            if vector_db_files:
                files_list = []
                for file_path in vector_db_files:
                    filename = os.path.basename(file_path)
                    # Try to get upload time from Firebase if available
                    upload_time = "Recently uploaded"
                    if db:
                        try:
                            docs = db.collection('uploaded_files').where('filename', '==', filename).stream()
                            for doc in docs:
                                file_data = doc.to_dict()
                                if 'uploaded_at' in file_data:
                                    upload_time = file_data['uploaded_at'].strftime('%Y-%m-%d')
                                    break
                        except:
                            pass
                    
                    files_list.append(f"• {filename} (Uploaded: {upload_time})")
                
                files_text = "\n".join(files_list[:10])  # Show first 10 files
                send_whatsapp_message(sender, f"📚 Files in your knowledge base:\n{files_text}")
                
                if len(vector_db_files) > 10:
                    send_whatsapp_message(sender, f"📖 ... and {len(vector_db_files) - 10} more files")
            else:
                send_whatsapp_message(sender, "❌ No files in knowledge base. Use 'Upload Materials' to add files.")
            
            end_conversation_and_show_menu(sender, None)
            return

            
        if lower_text == "podcast from image":
            update_user_state(sender, state="awaiting_podcast_image")
            send_whatsapp_message(sender, "Please send me an image of the text you'd like to convert to a podcast.")
            return
            
        if lower_text == "summary from image":
            update_user_state(sender, state="awaiting_summary_image")
            send_whatsapp_message(sender, "Please send me an image of the text you'd like me to summarize.")
            return

        if lower_text == "create video":
            send_whatsapp_message(sender, "This feature is coming soon!")
            end_conversation_and_show_menu(sender, None)
            return

    # --- State-Based Conversation Flow ---
    if current_state == "awaiting_question":
        response = query_subject_knowledge(text)
        end_conversation_and_show_menu(sender, response["result"])
        return
    
    if current_state == "awaiting_exam_type_selection":
        if  text.lower() == "Semester":
            start_question_paper_flow(sender)  # Start the semester flow we implemented
            return
        elif "cat" in text.lower():
            send_whatsapp_message(sender,
                "📅 CAT EXAM QUESTION PAPERS\n\n"
                "🚧 **Coming Soon!** 🚧\n\n"
                "CAT (Continuous Assessment Test) question paper generation "
                "is currently under development.\n\n"
                "Features coming:\n"
                "• Shorter format (20-30 marks)\n"
                "• Unit-wise question papers\n"
                "• Quick assessment patterns\n"
                "• Automatic difficulty scaling\n\n"
                "Please use 'Semester Exam' for now, or check back later!"
            )
            end_conversation_and_show_menu(sender, None)
            return
    elif "unit" in text.lower() or "test" in text.lower():
        send_whatsapp_message(sender,
            "📝 UNIT TEST QUESTION PAPERS\n\n"
            "🚧 **Coming Soon!** 🚧\n\n"
            "Unit test question paper generation is currently under development.\n\n"
            "Please use 'Semester Exam' for now!"
        )
        end_conversation_and_show_menu(sender, None)
        return
    else:
        send_whatsapp_message(sender, "Please select either 'Semester Exam', 'CAT Exam', or 'Unit Test'")
        return
    
    if current_state == "awaiting_ppt_topic":
        try:
            topic = text.strip()
            update_user_state(sender, temp_data={'ppt_topic': topic})
            send_whatsapp_message(sender, f"Okay, generating a 10-slide presentation on '{topic}'. This may take a moment...")
            ppt_content = generate_ppt_content(topic)
            if ppt_content:
                ppt_bytes = create_ppt_file(topic, ppt_content)
                if ppt_bytes:
                    status = send_whatsapp_ppt(sender, ppt_bytes, f"{topic.replace(' ', '_')}.pptx")
                    if status == "success":
                        end_conversation_and_show_menu(sender, "I've sent the presentation.")
                    else:
                        end_conversation_and_show_menu(sender, "I created the PPT, but there was an error sending it.")
                else:
                    end_conversation_and_show_menu(sender, "Sorry, I generated the content but failed to create the PPT file.")
            else:
                end_conversation_and_show_menu(sender, "Sorry, I couldn't generate content for that topic.")
        except Exception as e:
             print(f"Error in PPT generation flow: {e}")
             end_conversation_and_show_menu(sender, "An unexpected error occurred.")
        return
    
    if current_state == "awaiting_material_file":
    # This will be handled by the webhook when documents are received
        send_whatsapp_message(sender, 
            "📎 Please send the file you'd like to upload:\n"
            "• PDF, Word, PowerPoint, or Text files\n"
            "• I'll add it to my knowledge base\n"
            "• You can send multiple files"
        )
        return

    if current_state == "awaiting_qp_chapters":
        progress = get_qp_progress(sender)
        chapters_uploaded = progress.get('chapters_uploaded', 0)
        
        send_whatsapp_message(sender,
            f"Step 2/4: CHAPTER NOTES\n\n"
            f"📚 Progress: {chapters_uploaded}/5 chapters\n\n"
            "Please upload notes for all 5 chapters:\n"
            "• PDF/Word documents\n"
            "• Presentation slides\n"
            "• Study materials\n\n"
            "You can upload multiple files at once."
        )
        return

    if current_state == "awaiting_qp_sample_papers":
        send_whatsapp_message(sender,
            f"Step 3/4: SAMPLE PAPERS\n\n"
            "Please upload sample question papers:\n"
            "• Previous year question papers\n"
            "• Model question papers\n"
            "• Sample exam papers\n\n"
            "These will be used as reference for difficulty and pattern."
        )
        return
    

    # In process_message function, update the question paper states:

    if current_state == "ready_qp_generation":
        if "yes" in text.lower() or "generate" in text.lower():
            title = user_temp_data[sender].get('title', 'Unknown Title')
            progress = get_qp_progress(sender)
            
            send_whatsapp_message(sender,
                f"🎯 Generating Question Paper: {title}\n\n"
                f"📊 Materials Ready:\n"
                f"• Syllabus: {'✅' if progress.get('syllabus') else '❌'}\n"
                f"• Chapters: {progress.get('chapters_uploaded', 0)}/5\n"
                f"• Sample Papers: {'✅' if progress.get('sample_papers') else '❌'}\n\n"
                "⏱️ Generating high-quality question paper..."
            )
            
            # Generate question paper
            qp_content, status = generate_enhanced_qp_from_materials(title)
            
            if status == "success":
                pdf_bytes = create_question_paper_pdf(title, qp_content)
                if pdf_bytes:
                    filename = f"{title.replace(' ', '_')}.pdf"
                    send_whatsapp_document(sender, pdf_bytes, filename)
                    
                    send_whatsapp_message(sender,
                        f"✅ QUESTION PAPER GENERATED!\n\n"
                        f"📄 {title}\n"
                        f"📚 Based on your uploaded materials\n"
                        f"🎯 Professional quality with mixed difficulty"
                    )
                else:
                    send_whatsapp_message(sender, "❌ Generated content but couldn't create PDF.")
            else:
                send_whatsapp_message(sender, f"❌ {qp_content}")
            
            # Clean up and show menu
            clear_user_session(sender)
            end_conversation_and_show_menu(sender, None)
        else:
            # User wants to upload more materials
            user_states[sender] = "awaiting_qp_more_materials"
            send_whatsapp_message(sender, "What additional materials would you like to upload?")
        return

    if current_state == "awaiting_exam_type":
        debug_user_state(sender, f"Exam type selection: '{text}'")
        
        if "semester" in text.lower():
            # Clear current state and start semester flow
            clear_user_session(sender)
            start_question_paper_flow(sender)
        elif "cat" in text.lower():
            send_whatsapp_message(sender,
                "📅 CAT EXAM QUESTION PAPERS\n\n"
                "🚧 **Coming Soon!** 🚧\n\n"
                "CAT exam generation is under active development.\n"
                "Please use 'Semester Exam' for now!"
            )
            end_conversation_and_show_menu(sender, None)
        elif "unit" in text.lower() or "test" in text.lower():
            send_whatsapp_message(sender,
                "📝 UNIT TEST PAPERS\n\n"
                "🚧 **Coming Soon!** 🚧\n\n"
                "Unit test generation is under active development.\n"
                "Please use 'Semester Exam' for now!"
            )
            end_conversation_and_show_menu(sender, None)
        else:
            send_whatsapp_message(sender, "Please select 'Semester Exam', 'CAT Exam', or 'Unit Test'")
        return

        
    # ADD THIS: Handle the title input for semester flow

    if current_state == "awaiting_qp_title":
        debug_user_state(sender, f"Received title: '{text}'")
        
        # Initialize temp data if not exists
        if sender not in user_temp_data:
            user_temp_data[sender] = {}
        
        user_temp_data[sender]['title'] = text.strip()
        user_temp_data[sender]['current_step'] = 1
        user_temp_data[sender]['total_steps'] = 4
        user_temp_data[sender]['qp_materials'] = {}
        
        user_states[sender] = "awaiting_qp_syllabus"
        
        send_whatsapp_message(sender,
            f"📝 Title: {text}\n\n"
            f"Step 1/4: SYLLABUS UPLOAD\n\n"
            "Please upload the syllabus document:\n"
            "• Course outline PDF\n"
            "• Syllabus with topics\n"
            "• Curriculum document\n\n"
            "Send the file now..."
        )
        return

    if current_state == "awaiting_question_paper_topic":
        topic = text.strip()
        user_temp_data[sender] = {'qp_topic': topic}
        
        # Step 1: Check for syllabus
        send_whatsapp_message(sender, f"🔍 Checking for syllabus document: {topic}...")
        has_syllabus, syllabus_content = check_syllabus_document(topic)
        
        if not has_syllabus:
            user_states[sender] = "awaiting_syllabus_upload"
            send_whatsapp_message(sender,
                f"❌ No syllabus document found for '{topic}'\n\n"
                "Please upload the syllabus document (PDF preferred):\n"
                "• Course outline\n"
                "• Syllabus PDF\n"
                "• Curriculum document\n\n"
                "This helps me understand the course structure."
            )
            return end_conversation_and_show_menu(sender, None)
        
        # Step 2: Check chapter notes
        send_whatsapp_message(sender, "📚 Checking chapter notes for all 5 chapters...")
        chapter_status = check_chapter_notes_enhanced(topic)
        coverage = get_material_coverage_summary(chapter_status)
        
        user_temp_data[sender]['chapter_status'] = chapter_status
        user_temp_data[sender]['coverage'] = coverage
        
        if not coverage['sufficient_for_paper']:
            user_states[sender] = "awaiting_chapter_notes"
            missing_chapters = ", ".join(map(str, coverage['chapters_without_material']))
            
            send_whatsapp_message(sender,
                f"❌ Insufficient chapter notes for '{topic}'\n\n"
                f"📊 Current coverage: {len(coverage['chapters_with_material'])}/5 chapters\n"
                f"📋 Missing chapters: {missing_chapters}\n\n"
                "Please upload notes for the missing chapters:\n"
                "• PDF/Word documents\n"
                "• Presentation slides\n"
                "• Study materials\n\n"
                "I need at least 3 chapters to generate a proper question paper."
            )
            return end_conversation_and_show_menu(sender, None)
        
        # Step 3: Ready to generate
        user_states[sender] = "ready_to_generate_qp"
        send_whatsapp_message(sender,
            f"✅ All materials ready for '{topic}'\n\n"
            f"📋 Syllabus: Found ✓\n"
            f"📚 Chapters: {len(coverage['chapters_with_material'])}/5 covered ✓\n"
            f"📝 Pattern: 10×2 marks + 5×16 marks\n\n"
            "Ready to generate the question paper?"
        )
        send_interactive_message(sender, "Generate question paper?", ["Yes, Generate Now", "No, Upload More"])
        return

    if current_state == "awaiting_syllabus_upload":
        # This will be handled by the document upload webhook
        send_whatsapp_message(sender, "I'm ready for the syllabus document. Please upload it now.")
        # The webhook will process the upload and then re-check materials
        return

    if current_state == "awaiting_chapter_notes":
        # This will be handled by the document upload webhook
        missing_chapters = user_temp_data[sender]['coverage']['chapters_without_material']
        missing_text = ", ".join(map(str, missing_chapters))
        send_whatsapp_message(sender, f"Please upload notes for chapters: {missing_text}")
        return
    
        # In the webhook document handler, add this:
    # In the webhook document handler, update the bulk upload logic:

    if current_state in ["awaiting_bulk_upload", "awaiting_syllabus_upload", "awaiting_chapter_notes"]:
        send_whatsapp_message(sender, "I'm ready for your files! Please send the documents now.")
        return

    if current_state == "ready_to_generate_qp":
        if "yes" in text.lower() or "generate" in text.lower():
            topic = user_temp_data[sender]['qp_topic']
            
            send_whatsapp_message(sender,
                "📄 Generating semester question paper...\n"
                "🎯 Source: Your uploaded materials\n"
                "⏱️ This may take 1-2 minutes"
            )
            
            # Generate from the materials we've verified
            qp_content, status = generate_enhanced_semester_question_paper(topic)
            
            if status == "success" and qp_content:
                pdf_bytes = create_question_paper_pdf(f"Semester - {topic}", qp_content)
                
                if pdf_bytes:
                    filename = f"Semester_QP_{topic.replace(' ', '_')}.pdf"
                    send_whatsapp_document(sender, pdf_bytes, filename)
                    
                    coverage = user_temp_data[sender]['coverage']
                    success_msg = (
                        f"✅ Semester question paper generated!\n\n"
                        f"📚 Materials used:\n"
                        f"• Syllabus document ✓\n"
                        f"• {len(coverage['chapters_with_material'])}/5 chapters ✓\n"
                        f"📝 Format: 10×2 marks + 5×16 marks\n"
                        f"🎯 All questions from your uploaded content"
                    )
                    send_whatsapp_message(sender, success_msg)
                else:
                    send_whatsapp_message(sender, "Generated questions but couldn't create PDF.")
            else:
                send_whatsapp_message(sender, qp_content)  # Error message
            
            end_conversation_and_show_menu(sender, None)
        else:
            # User chose to upload more materials
            user_states[sender] = "awaiting_bulk_upload"
            send_whatsapp_message(sender, "Okay, please upload the additional materials you'd like to include.")
        return
        
    if current_state == "awaiting_question_paper_confirmation":
        topic = user_temp_data[sender]['qp_topic']
        
        if "yes" in text.lower() or "generate" in text.lower():
            send_whatsapp_message(sender, 
                "📄 Generating enhanced semester question paper...\n"
                "🎯 Validating materials thoroughly...\n"
                "📊 Checking content quality...\n"
                "⏱️ This may take 2-3 minutes"
            )
            
            # Use the ENHANCED generator
            qp_content, status = generate_enhanced_semester_question_paper(topic)
            
            if status == "success" and qp_content:
                pdf_bytes = create_question_paper_pdf(f"Semester - {topic}", qp_content)
                
                if pdf_bytes:
                    filename = f"Semester_QP_{topic.replace(' ', '_')}.pdf"
                    send_whatsapp_document(sender, pdf_bytes, filename)
                    
                    # Get coverage info
                    chapter_status = check_chapter_notes_enhanced(topic)
                    substantial_chapters = [chap for chap, status in chapter_status.items() 
                                        if status['coverage_quality'] in ['good', 'basic']]
                    
                    success_msg = (
                        f"✅ Enhanced Question Paper Generated!\n\n"
                        f"📚 Materials Used:\n"
                        f"• Syllabus: ✓\n"
                        f"• Substantial chapters: {len(substantial_chapters)}/5 ✓\n"
                        f"📝 Quality: Mixed difficulty levels\n"
                        f"🎯 All questions from YOUR uploaded content"
                    )
                    send_whatsapp_message(sender, success_msg)
                else:
                    send_whatsapp_message(sender, "Generated questions but couldn't create PDF.")
            else:
                send_whatsapp_message(sender, qp_content)  # Error message
            
            end_conversation_and_show_menu(sender, None)
        

    if current_state == "awaiting_bulk_upload":
        # This will be handled by the webhook when documents are received
        send_whatsapp_message(sender, "I'm ready for your files! Please send the PDFs now.")
        return

    if current_state == "awaiting_worksheet_topic":
        if text.lower() == "other topics":
            user_states[sender] = "awaiting_custom_topic"
            send_whatsapp_message(sender, "Please type the custom topic for your worksheet.")
        else:
            user_temp_data[sender]['topic'] = text
            user_states[sender] = "awaiting_worksheet_format"
            send_interactive_message(sender, f"Great! Topic is '{text}'.\nWhat format would you like?", ["PDF Worksheet", "Google Form Quiz"])
        return
        
    if current_state == "awaiting_custom_topic":
        user_temp_data[sender]['topic'] = text
        user_states[sender] = "awaiting_worksheet_format"
        send_interactive_message(sender, f"Great! Topic is '{text}'.\nWhat format would you like?", ["PDF Worksheet", "Google Form Quiz"])
        return
        
    if current_state == "awaiting_worksheet_format":
        user_temp_data[sender]['format'] = text
        user_states[sender] = "awaiting_worksheet_quantity"
        send_interactive_message(sender, f"Perfect, a {text}.\nHow many questions?", ["5", "10", "15"])
        return

    if current_state == "awaiting_worksheet_quantity":
        try:
            user_temp_data[sender]['quantity'] = int(text.strip())
            if "Google Form" in user_temp_data[sender]['format']:
                user_temp_data[sender]['type'] = "mcq"
                topic = user_temp_data[sender]['topic']
                quantity = user_temp_data[sender]['quantity']
                send_whatsapp_message(sender, f"Okay! Generating a {quantity}-question Google Form quiz on '{topic}'. Please wait...")
                
                worksheet_content_result = generate_worksheet_content_text(topic, quantity, "mcq", user_memory[sender])
                if worksheet_content_result.get("source") == "generated_worksheet_text":
                    form_result = create_google_form_mcq(f"Quiz: {topic}", worksheet_content_result["result"])
                    user_temp_data[sender]['form_url'] = form_result.get("result", "").split(' ')[-1]
                    send_whatsapp_message(sender, form_result["result"])
                    
                    user_states[sender] = "awaiting_assignment_title"
                    send_whatsapp_message(sender, "The quiz is ready! Now, what should the title of the assignment be for this quiz?")
                else:
                    end_conversation_and_show_menu(sender, "Sorry, I couldn't generate the quiz content.")
                return
            else: # PDF Flow
                user_states[sender] = "awaiting_worksheet_type"
                send_interactive_message(sender, f"Perfect, {text} questions.\nNow, what type of questions?", ["MCQ", "Short Answer", "Numerical"])
        except ValueError:
            send_whatsapp_message(sender, "Please select a valid number from the buttons.")
        return

    if current_state == "awaiting_worksheet_type":
        try:
            topic = user_temp_data[sender]['topic']
            quantity = user_temp_data[sender]['quantity']
            q_type = text.lower().replace(" answer", "")
            valid_types = ["mcq", "short", "numerical"]
            if q_type not in valid_types:
                send_whatsapp_message(sender, "That's not a valid type."); return

            send_whatsapp_message(sender, f"Okay! Generating {quantity} {q_type} questions on '{topic}'. Please wait...")
            
            worksheet_content_result = generate_worksheet_content_text(topic, quantity, q_type, user_memory[sender])
            
            if worksheet_content_result and worksheet_content_result.get("source") == "generated_worksheet_text":
                full_content = worksheet_content_result["result"]
                questions_text, answers_text = (full_content.split("--- ANSWERS ---", 1) + ["No answer key generated."])[:2]
                user_temp_data[sender]['questions_text'] = questions_text.strip()
                user_temp_data[sender]['answers_text'] = answers_text.strip()
                
                if worksheet_pdf_bytes := create_pdf_locally(f"Worksheet: {topic.title()}", questions_text.strip()):
                    send_whatsapp_document(sender, worksheet_pdf_bytes, f"{topic.replace(' ', '_')}_worksheet.pdf")
                
                if answer_key_pdf_bytes := create_pdf_locally(f"Answer Key: {topic.title()}", answers_text.strip()):
                    send_whatsapp_document(sender, answer_key_pdf_bytes, f"{topic.replace(' ', '_')}_answers.pdf")

                user_states[sender] = "awaiting_classroom_post_type"
                send_interactive_message(sender, "I've sent the PDFs. How should I post this to Classroom?", ["Post only questions", "Post with answers", "Don't post"])
            else:
                end_conversation_and_show_menu(sender, "Sorry, I couldn't generate the worksheet content.")
        except (ValueError, KeyError) as e:
            print(f"Error in worksheet type state: {e}")
            end_conversation_and_show_menu(sender, "An error occurred. Please start over.")
        return

    if current_state == "awaiting_classroom_post_type":
        user_temp_data[sender]['post_choice'] = text
        if text.lower() == "post only questions" or text.lower() == "post with answers":
            user_states[sender] = "awaiting_assignment_title"
            # SIMPLIFIED: Skip due dates, just ask for title
            send_whatsapp_message(sender, "Great! What should the assignment title be?")
        else: # Don't Post
            end_conversation_and_show_menu(sender, "Okay, I won't post to Classroom.")
        return

    if current_state == "awaiting_assignment_title":
        user_temp_data[sender]['assignment_title'] = text.strip()
        # SIMPLIFIED: Skip due date question and go straight to posting
        title = user_temp_data[sender]['assignment_title']
        handle_final_classroom_post(sender, title)
        return
    
    if text.lower() in ['reset', 'start over', 'restart', 'clear']:
        clear_user_session(sender)
        send_whatsapp_message(sender, "🔄 Session reset completely. Starting fresh...")
        menu_text = "What would you like to do next?"
        options = [
            "Ask Question", 
            "Create Worksheet", 
            "Create PPT", 
            "Generate Question Paper",
            "Upload Materials",
            "View Uploaded Files", 
            "Podcast from Image", 
            "Summary from Image"
        ]
        send_menu_message(sender, menu_text, options)
        return

    ## --- Fallback to a General Query ---
    print(f"Handling as a general query: '{text}'")
    response = query_subject_knowledge(text)  # Changed from query_os_subject to query_subject_knowledge
    send_whatsapp_message(sender, response["result"])
    end_conversation_and_show_menu(sender, None)

def cleanup_orphaned_file_metadata():
    """Remove Firebase entries for files that no longer exist in vector DB"""
    if not db:
        return
    
    try:
        # Get all files from Firebase
        docs = db.collection('uploaded_files').stream()
        vector_db_files = get_uploaded_files_list()
        vector_db_filenames = [os.path.basename(f) for f in vector_db_files]
        
        deleted_count = 0
        for doc in docs:
            file_data = doc.to_dict()
            filename = file_data.get('filename')
            
            # If file doesn't exist in vector DB, delete from Firebase
            if filename not in vector_db_filenames:
                db.collection('uploaded_files').document(doc.id).delete()
                deleted_count += 1
                print(f"Cleaned up orphaned file metadata: {filename}")
        
        if deleted_count > 0:
            print(f"Cleaned up {deleted_count} orphaned file metadata entries")
            
    except Exception as e:
        print(f"Error cleaning up orphaned file metadata: {e}")

def parse_deadline(text):
    text = text.lower()
    today = datetime.now()
    if "tomorrow" in text:
        return today + timedelta(days=1)
    if "next week" in text or "7 days" in text:
        return today + timedelta(days=7)
   
    return today + timedelta(days=2)

def send_interactive_message(to, text, buttons):
    """Sends an interactive message with up to 3 reply buttons with robust error handling."""
    url = f"https://graph.facebook.com/v19.0/{PHONE_NUMBER_ID}/messages"
    headers = {
        "Authorization": f"Bearer {ACCESS_TOKEN}", 
        "Content-Type": "application/json"
    }
    
    try:
        # Validate and format buttons properly
        button_payload = []
        for i, button_title in enumerate(buttons[:3]):  # Max 3 buttons
            # Clean button title - remove emojis and limit length
            clean_title = re.sub(r'[^\w\s\-()]', '', button_title)  # Remove special chars
            clean_title = clean_title[:20].strip()  # Limit to 20 chars
            
            if clean_title:  # Only add if we have valid text
                button_payload.append({
                    "type": "reply",
                    "reply": {
                        "id": f"btn_{i+1}",
                        "title": clean_title
                    }
                })
        
        # If no valid buttons, send as text instead
        if not button_payload:
            send_whatsapp_message(to, f"{text}\n\nPlease respond with your choice.")
            return
        
        payload = {
            "messaging_product": "whatsapp",
            "to": to,
            "type": "interactive",
            "interactive": {
                "type": "button",
                "body": {
                    "text": text[:1024]  # Limit text length
                },
                "action": {
                    "buttons": button_payload
                }
            }
        }
        
        print(f"Sending interactive message with payload: {json.dumps(payload, indent=2)}")
        
        response = requests.post(url, headers=headers, json=payload, timeout=10)
        response.raise_for_status()
        
        print(f"Interactive message sent successfully to {to}")
        
    except requests.exceptions.RequestException as e:
        print(f"Failed to send interactive message: {e}")
        if hasattr(e, 'response') and e.response is not None:
            print(f"Response content: {e.response.text}")
        
        # Fallback: Send as text message with options
        fallback_text = f"{text}\n\nPlease reply with:\n"
        for i, button in enumerate(buttons[:5], 1):  # Show up to 5 options in fallback
            fallback_text += f"{i}. {button}\n"
        
        send_whatsapp_message(to, fallback_text)
        
    except Exception as e:
        print(f"Unexpected error in interactive message: {e}")
        # Fallback to simple text
        send_whatsapp_message(to, f"{text}\n\nPlease respond with your choice.")

def send_text_options_message(to, text, options):
    """Alternative: Send options as numbered text message"""
    message = f"{text}\n\nPlease reply with the number:\n"
    for i, option in enumerate(options, 1):
        message += f"{i}. {option}\n"   
    
    send_whatsapp_message(to, message)
#
# ---------------------------------
# --- FLASK WEBHOOK ENDPOINTS (Updated for Document Uploads)
# ---------------------------------
@app.route("/webhook", methods=["GET", "POST"])
def webhook():
    if request.method == "GET":
        if request.args.get("hub.verify_token") == VERIFY_TOKEN:
            return request.args.get("hub.challenge")
        return "Verification token mismatch", 403
    
    data = request.get_json()
    try:
        if data and "entry" in data:
            for entry in data.get("entry", []):
                for change in entry.get("changes", []):
                    value = change.get("value", {})
                    if "messages" in value:
                        for msg in value.get("messages", []):
                            sender = msg["from"]
                            text = ""
                            msg_type = msg.get("type")

                            # Handle document messages for dynamic RAG
                            if msg_type == "document":
                                current_state = user_states.get(sender)
                                
                                # Handle different document upload states
                                if current_state in ["awaiting_bulk_upload", "awaiting_syllabus_upload", "awaiting_chapter_notes", "awaiting_material_file"]:
                                    media_id = msg['document']['id']  # ✅ msg is available here
                                    filename = msg['document'].get('filename', 'uploaded_file')
                                    
                                    # Initialize bulk upload items if not exists
                                    if 'bulk_upload_items' not in user_temp_data[sender]:
                                        user_temp_data[sender]['bulk_upload_items'] = []
                                    
                                    user_temp_data[sender]['bulk_upload_items'].append({
                                        'id': media_id,
                                        'filename': filename
                                    })
                                    
                                    # Wait for more files then process ALL files
                                    threading.Timer(8.0, process_bulk_upload, args=[sender]).start()
                                    
                                    send_whatsapp_message(sender, f"✅ Added {filename} to upload queue...")
                                else:
                                    # If not in upload state, process as single document
                                    media_id = msg['document']['id']
                                    filename = msg['document'].get('filename', 'uploaded_file')
                                    thread = threading.Thread(target=handle_document_upload, args=(sender, media_id, filename))
                                    thread.start()
                                return "ok", 200

                            # ✅ CORRECTED: Added handler for template button clicks
                            if msg_type == "button":
                                text = msg["button"]["text"]
                            elif msg_type == "text":
                                text = msg["text"]["body"]
                            elif msg_type == "interactive" and msg.get("interactive", {}).get("type") == "button_reply":
                                text = msg["interactive"]["button_reply"]["title"]
                            elif msg_type == "interactive" and msg.get("interactive", {}).get("type") == "list_reply":
                                text = msg["interactive"]["list_reply"]["title"]
                            elif msg_type == "audio":
                                media_id = msg['audio']['id']
                                thread = threading.Thread(target=handle_audio_task, args=(sender, media_id))
                                thread.start()
                            elif msg_type == "image":
                                current_state = user_states.get(sender)
                                media_id = msg['image']['id']
                                
                                if current_state == "awaiting_podcast_image":
                                    thread = threading.Thread(target=handle_podcast_task, args=(sender, media_id))
                                    thread.start()
                                elif current_state == "awaiting_summary_image":
                                    thread = threading.Thread(target=handle_summary_task, args=(sender, media_id))
                                    thread.start()
                                else:
                                    prompt = msg.get("image", {}).get("caption", "Explain this image.")
                                    thread = threading.Thread(target=handle_image_task, args=(sender, media_id, prompt))
                                    thread.start()
                                return "ok", 200
                            
                            if text:
                                process_message(sender, text)
    except Exception as e:
        print(f"Webhook processing error: {e}")
        traceback.print_exc()
    return "ok", 200

def process_bulk_upload(sender):
    """Process accumulated bulk upload items"""
    if sender in user_temp_data and 'bulk_upload_items' in user_temp_data[sender]:
        items = user_temp_data[sender]['bulk_upload_items']
        if items:
            print(f"🔄 Processing {len(items)} files for {sender}")
            
            # Get current state to determine what to do after upload
            current_state = user_states.get(sender)
            
            # Clear the items to avoid reprocessing
            user_temp_data[sender]['bulk_upload_items'] = []
            
            # Process the bulk upload
            handle_bulk_document_upload(sender, items)
            
            # After bulk upload, check if we need to continue question paper flow
            if current_state in ["awaiting_syllabus_upload", "awaiting_chapter_notes"]:
                threading.Timer(5.0, lambda: continue_question_paper_flow(sender)).start()
        else:
            print(f"ℹ️ No items to process for {sender}")
    else:
        print(f"ℹ️ No bulk upload items found for {sender}")

def continue_question_paper_flow(sender):
    """Continue question paper flow after uploads are processed"""
    if sender in user_temp_data and 'qp_topic' in user_temp_data[sender]:
        topic = user_temp_data[sender]['qp_topic']
        send_whatsapp_message(sender, "🔄 Uploads processed! Continuing with question paper generation...")
        
        # Re-trigger the question paper flow
        user_states[sender] = "awaiting_question_paper_topic"
        process_message(sender, topic)

@app.route("/status", methods=["GET"])
def status():
    return {"status": "running"}, 200

if __name__ == "__main__":
    os.makedirs("data", exist_ok=True)
    os.makedirs("vector_index", exist_ok=True)
    os.makedirs(DYNAMIC_VECTOR_INDEX_PATH, exist_ok=True)
    
    # Initialize APIs
    init_google_apis()
    
    # Check Firebase connection
    if db:
        print("Firebase connected successfully")
        cleanup_orphaned_file_metadata()
    else:
        print("Firebase not connected - running in local mode only")
    
    # --- Send Startup Template Message ---
    print("Sending startup template to users...")
    for number in student_phone_numbers:
        send_start_template(number)
    print("Finished sending templates.")

    # ✅ ONLY Dynamic RAG - No static embeddings or databases
    print("Dynamic RAG system ready - using only user-uploaded documents")
    
    print("Starting Flask app...")
    app.run(port=5000, debug=False)