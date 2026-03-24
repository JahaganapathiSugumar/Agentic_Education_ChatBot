# ==============================================================================
# --- 1. IMPORTS ---
# ==============================================================================
import os
import re
import json
import base64
import tempfile
import threading
import traceback
import uuid
import warnings
import pickle  # Added for OAuth token storage
from abc import ABC, abstractmethod
from collections import defaultdict
from datetime import datetime, timedelta
from io import BytesIO
from functools import wraps
import time
import gc  # Added for garbage collection

import requests
import numpy as np
from flask import Flask, request, render_template, jsonify, session, redirect, url_for, make_response
from functools import wraps

# --- Media / File Generation ---
from fpdf import FPDF
from fpdf.enums import XPos, YPos
from pydub import AudioSegment
from gtts import gTTS
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt
import fitz
from docx import Document as DocxDocument
from docx.shared import Inches, Pt
import moviepy
from moviepy import ImageClip, concatenate_videoclips, AudioFileClip

# --- AI & Machine Learning ---
import google.generativeai as genai
from sklearn.metrics.pairwise import cosine_similarity
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.documents import Document

# --- Google Cloud APIs ---
from googleapiclient.discovery import build
from googleapiclient.errors import HttpError
from googleapiclient.http import MediaFileUpload, MediaIoBaseUpload

# --- OAuth 2.0 for Google APIs ---
from google_auth_oauthlib.flow import InstalledAppFlow
from google.auth.transport.requests import Request
from google.oauth2.credentials import Credentials

# --- Firebase ---
import firebase_admin
from firebase_admin import credentials, firestore, auth as fb_auth

# --- Config ---
from dotenv import load_dotenv

# ==============================================================================
# --- 2. INITIAL SETUP ---
# ==============================================================================
load_dotenv()

VERIFY_TOKEN = os.getenv("VERIFY_TOKEN")
ACCESS_TOKEN = os.getenv("ACCESS_TOKEN")
PHONE_NUMBER_ID = os.getenv("PHONE_NUMBER_ID")
TEMPLATE_NAME = os.getenv("TEMPLATE_NAME")
GENAI_API_KEY = os.getenv("GENAI_API_KEY")
GOOGLE_FORMS_TARGET_FOLDER_ID = os.getenv("GOOGLE_FORMS_TARGET_FOLDER_ID")
CLASSROOM_COURSE_ID = os.getenv("CLASSROOM_COURSE_ID")
FIREBASE_WEB_API_KEY = os.getenv("FIREBASE_WEB_API_KEY", "")

student_phone_numbers = ["916379613654", "918870420449", "917530086388"]

genai.configure(api_key=GENAI_API_KEY)
GOOGLE_FORMS_TARGET_FOLDER_ID = '141ThSl8rtgjmBsrUz41xac6uThjkI0mQ'
CLASSROOM_COURSE_ID = '791255382357'

# --- State & Memory ---
user_states = defaultdict(str)
user_temp_data = defaultdict(dict)
user_memory = defaultdict(list)
MAX_HISTORY = 50

# --- Dynamic RAG ---
DYNAMIC_VECTOR_INDEX_PATH = "vector_index/dynamic_uploads"
os.makedirs(DYNAMIC_VECTOR_INDEX_PATH, exist_ok=True)

embeddings_model_for_classification = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
dynamic_embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")

# --- Flask App Initialization ---
app = Flask(__name__)
app.secret_key = os.getenv("SECRET_KEY", "your-secret-key-change-in-production")

# --- API Services ---
forms_service = None
drive_service = None
classroom_service = None
speech_client = None
creds = None
calendar_service = None

# ==============================================================================
# --- 3. AUTH DECORATORS ---
# ==============================================================================

def login_required(f):
    """Decorator to require login for a route"""
    @wraps(f)
    def decorated_function(*args, **kwargs):
        if 'user_id' not in session:
            return jsonify({"error": "Unauthorized"}), 401
        return f(*args, **kwargs)
    return decorated_function

def teacher_required(f):
    """Decorator to require teacher role for a route"""
    @wraps(f)
    def decorated_function(*args, **kwargs):
        if 'user_id' not in session:
            return jsonify({"error": "Unauthorized"}), 401
        if session.get('role') != 'teacher':
            return jsonify({"error": "Forbidden: Teacher access required"}), 403
        return f(*args, **kwargs)
    return decorated_function

# ==============================================================================
# --- 4. OAUTH 2.0 INITIALIZATION (FIXED TOKEN REFRESH) ---
# ==============================================================================

def init_google_apis_oauth():
    """Initialize Google APIs using OAuth 2.0 for user account with improved token handling"""
    global forms_service, drive_service, classroom_service, creds, calendar_service
    
    try:
        SCOPES = [
            "https://www.googleapis.com/auth/drive",              # Full Drive access
            "https://www.googleapis.com/auth/drive.file",         # Drive file access
            "https://www.googleapis.com/auth/forms.body",         # Forms creation
            "https://www.googleapis.com/auth/forms.responses.readonly",  # Read form responses
            "https://www.googleapis.com/auth/classroom.coursework.students",  # Coursework
            "https://www.googleapis.com/auth/classroom.announcements",  # Announcements
            "https://www.googleapis.com/auth/classroom.courses",  # Course management
            "https://www.googleapis.com/auth/calendar",           # Calendar access
            "https://www.googleapis.com/auth/calendar.events",    # Calendar events
        ]
        
        creds = None
        token_file = 'token.pickle'
        token_backup_file = 'token_backup.pickle'
        
        # Load existing token if available
        if os.path.exists(token_file):
            try:
                with open(token_file, 'rb') as token:
                    creds = pickle.load(token)
                    print("✅ Loaded existing OAuth token from token.pickle")
                    
                    # Check if token is valid and not expired
                    if creds and creds.valid:
                        print("✅ Token is valid")
                    elif creds and creds.expired and creds.refresh_token:
                        print("🔄 Token expired, attempting refresh...")
                        try:
                            creds.refresh(Request())
                            print("✅ Token refreshed successfully")
                            # Save refreshed token
                            with open(token_file, 'wb') as token:
                                pickle.dump(creds, token)
                            print("💾 Refreshed token saved")
                        except Exception as refresh_error:
                            print(f"❌ Token refresh failed: {refresh_error}")
                            # Try to load backup token if available
                            if os.path.exists(token_backup_file):
                                try:
                                    with open(token_backup_file, 'rb') as backup_token:
                                        creds = pickle.load(backup_token)
                                    print("✅ Loaded backup token")
                                except:
                                    creds = None
                            else:
                                creds = None
                    elif creds and creds.expired and not creds.refresh_token:
                        print("⚠️ Token expired with no refresh token - needs re-authentication")
                        creds = None
            except Exception as load_error:
                print(f"❌ Error loading token: {load_error}")
                # Try backup
                if os.path.exists(token_backup_file):
                    try:
                        with open(token_backup_file, 'rb') as backup_token:
                            creds = pickle.load(backup_token)
                        print("✅ Loaded backup token")
                    except:
                        creds = None
                else:
                    creds = None
        
        # If credentials are invalid or don't exist, start OAuth flow
        if not creds or not creds.valid:
            print("🔑 Starting OAuth 2.0 authorization flow...")
            print("📢 A browser window will open for you to log in and grant permissions.")
            
            # Check if credentials.json exists
            if not os.path.exists('credentials.json'):
                print("❌ credentials.json not found!")
                print("📥 Download OAuth 2.0 credentials from Google Cloud Console:")
                print("   1. Go to https://console.cloud.google.com/")
                print("   2. Select your project: sahayak-465916")
                print("   3. Go to APIs & Services → Credentials")
                print("   4. Create Credentials → OAuth 2.0 Client IDs")
                print("   5. Choose 'Desktop Application' as type")
                print("   6. Download the JSON and save as 'credentials.json' in project root")
                print("   7. Also ensure the OAuth consent screen is properly configured")
                return
            
            try:
                flow = InstalledAppFlow.from_client_secrets_file(
                    'credentials.json',
                    SCOPES
                )
                # Add a timeout to prevent hanging
                creds = flow.run_local_server(port=0, open_browser=True)
                print("✅ OAuth authorization complete!")
                
                # Save credentials with backup
                with open(token_file, 'wb') as token:
                    pickle.dump(creds, token)
                with open(token_backup_file, 'wb') as backup_token:
                    pickle.dump(creds, backup_token)
                print("💾 OAuth token saved to token.pickle and backup")
                
            except Exception as auth_error:
                print(f"❌ OAuth authorization failed: {auth_error}")
                print("💡 Troubleshooting steps:")
                print("   1. Delete token.pickle and try again")
                print("   2. Check that your credentials.json is valid and not expired")
                print("   3. Verify OAuth consent screen is properly configured")
                print("   4. Check that the correct scopes are enabled in Google Cloud Console")
                return
        
        # Build services with user credentials
        drive_service = build('drive', 'v3', credentials=creds)
        classroom_service = build('classroom', 'v1', credentials=creds)
        forms_service = build('forms', 'v1', credentials=creds)
        
        # Test the connection
        try:
            about = drive_service.about().get(fields="user").execute()
            user_email = about.get('user', {}).get('emailAddress', 'Unknown')
            print(f"✅ Google APIs initialized successfully with OAuth 2.0")
            print(f"👤 Authenticated as: {user_email}")
            
            # Verify drive access
            test_files = drive_service.files().list(pageSize=1).execute()
            print(f"✅ Drive access verified: found {len(test_files.get('files', []))} files")
            
        except HttpError as api_error:
            print(f"⚠️ API test failed: {api_error}")
            print("This might indicate insufficient permissions or API not enabled")
            print("Please ensure all required APIs are enabled in Google Cloud Console:")
            print("  - Google Drive API")
            print("  - Google Classroom API")
            print("  - Google Forms API")
            print("  - Google Calendar API")
            
        except Exception as e:
            print(f"⚠️ Connection test failed: {e}")
        
    except ImportError as e:
        print(f"❌ Missing required library: {e}")
        print("💡 Run: pip install google-auth-oauthlib google-auth-httplib2")
    except FileNotFoundError as e:
        print(f"❌ credentials.json not found: {e}")
        print("📥 Download OAuth 2.0 credentials from Google Cloud Console and save as 'credentials.json'")
    except Exception as e:
        print(f"❌ Failed to initialize Google APIs with OAuth: {e}")
        traceback.print_exc()

# ==============================================================================
# --- 5. RESPONDER CLASSES ---
# ==============================================================================

class Responder(ABC):
    @abstractmethod
    def text(self, to, message):
        pass
    @abstractmethod
    def menu(self, to, text, options):
        pass
    @abstractmethod
    def document(self, to, file_bytes, filename):
        pass
    @abstractmethod
    def audio(self, to, audio_bytes, filename):
        pass
    @abstractmethod
    def video(self, to, video_bytes, filename):
        pass
    @abstractmethod
    def interactive(self, to, text, buttons):
        pass

class WhatsAppResponder(Responder):
    """Uses existing WhatsApp send functions."""
    def __init__(self, access_token, phone_number_id):
        self.access_token = access_token
        self.phone_number_id = phone_number_id

    def _send_api(self, payload):
        url = f"https://graph.facebook.com/v19.0/{self.phone_number_id}/messages"
        headers = {"Authorization": f"Bearer {self.access_token}", "Content-Type": "application/json"}
        return requests.post(url, headers=headers, json=payload)

    def text(self, to, message):
        payload = {"messaging_product": "whatsapp", "to": to, "text": {"body": message}}
        self._send_api(payload)

    def menu(self, to, text, options):
        rows = [{"id": f"option_{i}", "title": opt} for i, opt in enumerate(options)]
        payload = {
            "messaging_product": "whatsapp",
            "to": to,
            "type": "interactive",
            "interactive": {
                "type": "list",
                "body": {"text": text},
                "action": {
                    "button": "Choose",
                    "sections": [{"title": "Options", "rows": rows}]
                }
            }
        }
        self._send_api(payload)

    def document(self, to, file_bytes, filename):
        send_whatsapp_document(to, file_bytes, filename)

    def audio(self, to, audio_bytes, filename):
        send_whatsapp_audio(to, audio_bytes, filename)

    def video(self, to, video_bytes, filename):
        send_whatsapp_video(to, video_bytes, filename)

    def interactive(self, to, text, buttons):
        send_interactive_message(to, text, buttons)

class WebResponder(Responder):
    """Collects actions for JSON response."""
    def __init__(self):
        self.actions = []

    def text(self, to, message):
        self.actions.append({"type": "text", "content": message})

    def menu(self, to, text, options):
        self.actions.append({"type": "menu", "text": text, "options": options})

    def document(self, to, file_bytes, filename):
        self.actions.append({
            "type": "document",
            "filename": filename,
            "data": base64.b64encode(file_bytes).decode('utf-8')
        })

    def audio(self, to, audio_bytes, filename):
        self.actions.append({
            "type": "audio",
            "filename": filename,
            "data": base64.b64encode(audio_bytes).decode('utf-8')
        })

    def video(self, to, video_bytes, filename):
        self.actions.append({
            "type": "video",
            "filename": filename,
            "data": base64.b64encode(video_bytes).decode('utf-8')
        })

    def interactive(self, to, text, buttons):
        self.actions.append({"type": "interactive", "text": text, "buttons": buttons})

# ==============================================================================
# --- 6. WHATSAPP FUNCTIONS ---
# ==============================================================================


def send_whatsapp_message(to, message):
    """Send a plain text WhatsApp message."""
    url = f"https://graph.facebook.com/v19.0/{PHONE_NUMBER_ID}/messages"
    headers = {"Authorization": f"Bearer {ACCESS_TOKEN}", "Content-Type": "application/json"}
    payload = {"messaging_product": "whatsapp", "to": to, "text": {"body": message}}
    try:
        requests.post(url, headers=headers, json=payload)
    except Exception as e:
        print(f"Error sending WhatsApp message: {e}")

def send_whatsapp_document(to, file_bytes, filename):
    """Upload and send a document via WhatsApp Cloud API."""
    url = f"https://graph.facebook.com/v19.0/{PHONE_NUMBER_ID}/media"
    headers = {"Authorization": f"Bearer {ACCESS_TOKEN}"}
    mime_map = {'.pdf': 'application/pdf', '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
                '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document', '.txt': 'text/plain'}
    ext = os.path.splitext(filename)[1].lower()
    mime = mime_map.get(ext, 'application/octet-stream')
    try:
        upload_resp = requests.post(url, headers=headers,
                                    files={"file": (filename, file_bytes, mime)},
                                    data={"messaging_product": "whatsapp"})
        media_id = upload_resp.json().get("id")
        if media_id:
            msg_url = f"https://graph.facebook.com/v19.0/{PHONE_NUMBER_ID}/messages"
            payload = {"messaging_product": "whatsapp", "to": to, "type": "document",
                       "document": {"id": media_id, "filename": filename}}
            requests.post(msg_url, headers={**headers, "Content-Type": "application/json"}, json=payload)
    except Exception as e:
        print(f"Error sending WhatsApp document: {e}")

def send_whatsapp_audio(to, audio_bytes, filename):
    """Upload and send an audio file via WhatsApp Cloud API."""
    url = f"https://graph.facebook.com/v19.0/{PHONE_NUMBER_ID}/media"
    headers = {"Authorization": f"Bearer {ACCESS_TOKEN}"}
    try:
        upload_resp = requests.post(url, headers=headers,
                                    files={"file": (filename, audio_bytes, "audio/mpeg")},
                                    data={"messaging_product": "whatsapp"})
        media_id = upload_resp.json().get("id")
        if media_id:
            msg_url = f"https://graph.facebook.com/v19.0/{PHONE_NUMBER_ID}/messages"
            payload = {"messaging_product": "whatsapp", "to": to, "type": "audio",
                       "audio": {"id": media_id}}
            requests.post(msg_url, headers={**headers, "Content-Type": "application/json"}, json=payload)
    except Exception as e:
        print(f"Error sending WhatsApp audio: {e}")

def send_whatsapp_video(to, video_bytes, filename):
    """Upload and send a video file via WhatsApp Cloud API."""
    url = f"https://graph.facebook.com/v19.0/{PHONE_NUMBER_ID}/media"
    headers = {"Authorization": f"Bearer {ACCESS_TOKEN}"}
    try:
        upload_resp = requests.post(url, headers=headers,
                                    files={"file": (filename, video_bytes, "video/mp4")},
                                    data={"messaging_product": "whatsapp"})
        media_id = upload_resp.json().get("id")
        if media_id:
            msg_url = f"https://graph.facebook.com/v19.0/{PHONE_NUMBER_ID}/messages"
            payload = {"messaging_product": "whatsapp", "to": to, "type": "video",
                       "video": {"id": media_id}}
            requests.post(msg_url, headers={**headers, "Content-Type": "application/json"}, json=payload)
    except Exception as e:
        print(f"Error sending WhatsApp video: {e}")

def send_interactive_message(to, text, buttons):
    """Send an interactive button message via WhatsApp Cloud API."""
    url = f"https://graph.facebook.com/v19.0/{PHONE_NUMBER_ID}/messages"
    headers = {"Authorization": f"Bearer {ACCESS_TOKEN}", "Content-Type": "application/json"}
    btn_list = [{"type": "reply", "reply": {"id": f"btn_{i}", "title": b[:20]}} for i, b in enumerate(buttons[:3])]
    payload = {"messaging_product": "whatsapp", "to": to, "type": "interactive",
               "interactive": {"type": "button", "body": {"text": text}, "action": {"buttons": btn_list}}}
    try:
        requests.post(url, headers=headers, json=payload)
    except Exception as e:
        print(f"Error sending interactive message: {e}")

def send_start_template(to):
    """Send the startup template message to a user."""
    url = f"https://graph.facebook.com/v19.0/{PHONE_NUMBER_ID}/messages"
    headers = {"Authorization": f"Bearer {ACCESS_TOKEN}", "Content-Type": "application/json"}
    payload = {"messaging_product": "whatsapp", "to": to, "type": "template",
               "template": {"name": TEMPLATE_NAME or "hello_world", "language": {"code": "en_US"}}}
    try:
        requests.post(url, headers=headers, json=payload)
    except Exception as e:
        print(f"Error sending start template to {to}: {e}")

def get_media_url(media_id):
    """Retrieve the download URL for a WhatsApp media object."""
    url = f"https://graph.facebook.com/v19.0/{media_id}"
    headers = {"Authorization": f"Bearer {ACCESS_TOKEN}"}
    try:
        resp = requests.get(url, headers=headers)
        return resp.json().get("url")
    except Exception as e:
        print(f"Error getting media URL: {e}")
        return None

def download_media_file(media_url):
    """Download a WhatsApp media file and return its bytes."""
    headers = {"Authorization": f"Bearer {ACCESS_TOKEN}"}
    try:
        resp = requests.get(media_url, headers=headers)
        if resp.status_code == 200:
            return resp.content
    except Exception as e:
        print(f"Error downloading media file: {e}")
    return None

def transcribe_audio(audio_bytes):
    """Transcribe audio bytes using OpenAI Whisper."""
    try:
        import whisper as _whisper
        with tempfile.NamedTemporaryFile(suffix=".ogg", delete=False) as tmp:
            tmp.write(audio_bytes)
            tmp_path = tmp.name
        model = _whisper.load_model("base")
        result = model.transcribe(tmp_path)
        os.unlink(tmp_path)
        return result.get("text", "").strip()
    except Exception as e:
        print(f"Error transcribing audio: {e}")
        return None

# ==============================================================================
# --- 7. DOCUMENT PROCESSING FUNCTIONS ---
# ==============================================================================


def extract_text_from_file(file_bytes, filename):
    """Extract plain text from PDF, DOCX, PPTX, or TXT file bytes."""
    ext = os.path.splitext(filename)[1].lower()
    text = ""
    try:
        if ext == ".pdf":
            import fitz as _fitz
            doc = _fitz.open(stream=file_bytes, filetype="pdf")
            text = "\n".join(page.get_text() for page in doc)
        elif ext == ".docx":
            from docx import Document as _DocxDoc
            from io import BytesIO
            doc = _DocxDoc(BytesIO(file_bytes))
            text = "\n".join(p.text for p in doc.paragraphs)
        elif ext == ".pptx":
            from pptx import Presentation as _Prs
            from io import BytesIO
            prs = _Prs(BytesIO(file_bytes))
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        elif ext in (".txt", ".md"):
            text = file_bytes.decode("utf-8", errors="ignore")
        else:
            text = file_bytes.decode("utf-8", errors="ignore")
    except Exception as e:
        print(f"Error extracting text from {filename}: {e}")
    return text.strip()

def add_documents_to_vector_index(file_bytes, filename):
    """Add a document to the dynamic FAISS vector index. Returns (success, message)."""
    try:
        text = extract_text_from_file(file_bytes, filename)
        if not text:
            return False, "Could not extract text from the file."
        splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
        chunks = splitter.split_text(text)
        if not chunks:
            return False, "No text chunks could be created."
        docs = [Document(page_content=chunk, metadata={"source": filename}) for chunk in chunks]
        index_path = DYNAMIC_VECTOR_INDEX_PATH
        try:
            existing = FAISS.load_local(index_path, dynamic_embeddings, allow_dangerous_deserialization=True)
            existing.add_documents(docs)
            existing.save_local(index_path)
        except Exception:
            new_index = FAISS.from_documents(docs, dynamic_embeddings)
            new_index.save_local(index_path)
        return True, f"Successfully added {len(chunks)} chunks from {filename}."
    except Exception as e:
        print(f"Error adding document to vector index: {e}")
        return False, str(e)

def get_uploaded_files_list():
    """Return a list of source file names currently in the vector index."""
    try:
        index_path = DYNAMIC_VECTOR_INDEX_PATH
        index = FAISS.load_local(index_path, dynamic_embeddings, allow_dangerous_deserialization=True)
        sources = list({doc.metadata.get("source", "") for doc in index.docstore._dict.values()})
        return [s for s in sources if s]
    except Exception:
        return []

def check_similarity_only(query, top_k=5):
    """
    Utility function to check similarity scores without generating answers.
    Useful for debugging and monitoring.
    """
    try:
        index_path = DYNAMIC_VECTOR_INDEX_PATH
        if not os.path.exists(index_path):
            return {"error": "No vector index found"}
        
        index = FAISS.load_local(index_path, dynamic_embeddings, allow_dangerous_deserialization=True)
        docs_with_scores = index.similarity_search_with_score(query, k=top_k)
        
        results = []
        for doc, score in docs_with_scores:
            similarity = 1 / (1 + score)
            results.append({
                "content": doc.page_content[:200] + "...",
                "source": doc.metadata.get("source", "Unknown"),
                "similarity": similarity,
                "distance": score
            })
        
        return results
    except Exception as e:
        return {"error": str(e)}

# ==============================================================================
# --- 7. AI FUNCTIONS ---
# ==============================================================================

def fallback_to_gemini(query):
    """Fallback function to use Gemini directly"""
    try:
        print(f"Falling back to Gemini for query: {query[:50]}...")
        model = genai.GenerativeModel('gemini-3-flash-preview')
        response = model.generate_content(query)
        return {"result": response.text, "source": "gemini_direct", "similarity": None}
    except Exception as e2:
        return {"result": f"Sorry, I couldn't find an answer. Error: {e2}", 
                "source": "error", "similarity": None}

def query_dynamic_rag(query, similarity_threshold=0.7):
    """
    Query the dynamic FAISS vector index and only use Gemini if similarity is low.
    
    Args:
        query: The user's question
        similarity_threshold: Minimum similarity score (0-1) to consider relevant
                             (lower threshold = more likely to use RAG)
    
    Returns:
        Dictionary with result and source information
    """
    try:
        index_path = DYNAMIC_VECTOR_INDEX_PATH
        
        # Check if vector store exists
        if not os.path.exists(index_path):
            print("No vector index found, falling back to Gemini")
            return fallback_to_gemini(query)
        
        # Load the index
        index = FAISS.load_local(index_path, dynamic_embeddings, allow_dangerous_deserialization=True)
        
        # Get documents with similarity scores
        docs_with_scores = index.similarity_search_with_score(query, k=4)
        
        # FAISS returns L2 distance (lower = more similar)
        # Convert to similarity score (0-1, higher = more similar)
        docs_with_similarity = []
        for doc, score in docs_with_scores:
            # Convert L2 distance to similarity score (approximate)
            similarity = 1 / (1 + score)  # Simple conversion
            docs_with_similarity.append((doc, similarity))
        
        # Filter documents above threshold
        relevant_docs = [(doc, sim) for doc, sim in docs_with_similarity if sim >= similarity_threshold]
        
        if not relevant_docs:
            print(f"No documents found above similarity threshold {similarity_threshold}. Best score: {docs_with_similarity[0][1] if docs_with_similarity else 'N/A'}")
            return fallback_to_gemini(query)
        
        # Use only relevant docs for context
        context = "\n\n".join(doc.page_content for doc, _ in relevant_docs)
        sources = list(dict.fromkeys(doc.metadata.get("source", "") for doc, _ in relevant_docs))
        
        # Calculate average similarity for logging
        avg_similarity = sum(sim for _, sim in relevant_docs) / len(relevant_docs)
        print(f"Using RAG with {len(relevant_docs)} docs, avg similarity: {avg_similarity:.3f}")
        
        # Generate answer using Gemini with context
        model = genai.GenerativeModel('gemini-3-flash-preview')
        prompt = f"""Based on the following context, answer the question. If the answer cannot be found in the context, say so.

Context:
{context}

Question: {query}

Answer based on the context above:"""
        
        response = model.generate_content(prompt)
        answer = response.text
        
        # Add source information
        if sources:
            source_label = ", ".join(sources)
            answer = f"{answer}\n\n📚 Source: {source_label}"
        
        return {"result": answer, "source": "dynamic_rag", "similarity": avg_similarity}
        
    except Exception as e:
        print(f"Error in query_dynamic_rag: {e}")
        return fallback_to_gemini(query)

def query_dynamic_rag_strict(query, use_only_vector_db=True, similarity_threshold=0.35):
    """
    Query ONLY the dynamic FAISS vector index. No fallback to Gemini.
    
    Args:
        query: The user's question
        use_only_vector_db: If True, only return results from vector DB (no Gemini)
        similarity_threshold: Minimum similarity score to consider relevant (lower = more inclusive)
    
    Returns:
        Dictionary with result and source information
    """
    try:
        index_path = DYNAMIC_VECTOR_INDEX_PATH
        
        # Check if vector store exists
        if not os.path.exists(index_path):
            print("No vector index found")
            return {"result": None, "source": "no_vector_db", "similarity": None}
        
        # Load the index
        index = FAISS.load_local(index_path, dynamic_embeddings, allow_dangerous_deserialization=True)
        
        # Get documents with similarity scores - increase k to get more candidates
        docs_with_scores = index.similarity_search_with_score(query, k=10)
        
        # Convert L2 distance to similarity score
        docs_with_similarity = []
        for doc, score in docs_with_scores:
            similarity = 1 / (1 + score)  # Simple conversion
            docs_with_similarity.append((doc, similarity))
        
        # Filter documents above threshold
        relevant_docs = [(doc, sim) for doc, sim in docs_with_similarity if sim >= similarity_threshold]
        
        if not relevant_docs:
            print(f"No documents found above similarity threshold {similarity_threshold}")
            # Try with even lower threshold as last resort
            if similarity_threshold > 0.2:
                return query_dynamic_rag_strict(query, use_only_vector_db, similarity_threshold=0.2)
            return {"result": None, "source": "no_relevant_docs", "similarity": None}
        
        # Use only relevant docs for context
        context = "\n\n".join(doc.page_content for doc, _ in relevant_docs)
        sources = list(dict.fromkeys(doc.metadata.get("source", "") for doc, _ in relevant_docs))
        
        # Calculate average similarity
        avg_similarity = sum(sim for _, sim in relevant_docs) / len(relevant_docs)
        print(f"Using RAG with {len(relevant_docs)} docs, avg similarity: {avg_similarity:.3f}")
        
        # Generate answer using Gemini with context (but ONLY from vector DB)
        model = genai.GenerativeModel('gemini-3-flash-preview')
        prompt = f"""Based ONLY on the following context from uploaded materials, answer the question. 
Do not use any external knowledge or general information. If the answer cannot be found in the context, say "I cannot find this information in your uploaded materials."

Context:
{context}

Question: {query}

Answer based ONLY on the context above:"""
        
        response = model.generate_content(prompt)
        answer = response.text
        
        # Add source information
        if sources:
            source_label = ", ".join(sources)
            answer = f"{answer}\n\n📚 Source: {source_label}"
        
        return {"result": answer, "source": "dynamic_rag", "similarity": avg_similarity}
        
    except Exception as e:
        print(f"Error in query_dynamic_rag_strict: {e}")
        return {"result": None, "source": "error", "similarity": None}

def query_vector_db_for_topic(topic, similarity_threshold=0.3):
    """
    Query vector DB for content related to a topic.
    Returns concatenated relevant content or None if no relevant content found.
    """
    try:
        index_path = DYNAMIC_VECTOR_INDEX_PATH
        
        if not os.path.exists(index_path):
            return None
        
        index = FAISS.load_local(index_path, dynamic_embeddings, allow_dangerous_deserialization=True)
        docs_with_scores = index.similarity_search_with_score(topic, k=15)  # Increased k
        
        # Convert scores to similarity and filter
        relevant_docs = []
        for doc, score in docs_with_scores:
            similarity = 1 / (1 + score)
            if similarity >= similarity_threshold:
                relevant_docs.append(doc)
        
        if not relevant_docs:
            return None
        
        # Combine all relevant content
        combined_content = "\n\n".join(doc.page_content for doc in relevant_docs)
        return combined_content
        
    except Exception as e:
        print(f"Error querying vector DB for topic: {e}")
        return None

def query_gemini_vision(prompt, image_bytes):
    """Send an image + prompt to Gemini Vision and return the response."""
    try:
        model = genai.GenerativeModel('gemini-3-flash-preview')
        image_part = {"mime_type": "image/jpeg", "data": image_bytes}
        response = model.generate_content([prompt, image_part])
        return {"result": response.text, "source": "gemini_vision"}
    except Exception as e:
        print(f"Error in query_gemini_vision: {e}")
        return {"result": f"Sorry, I couldn't analyze the image. Error: {e}", "source": "error"}

def generate_voiceover(text):
    """Convert text to speech using gTTS and return audio bytes (MP3)."""
    try:
        from gtts import gTTS
        from io import BytesIO
        tts = gTTS(text=text, lang='en', slow=False)
        buf = BytesIO()
        tts.write_to_fp(buf)
        buf.seek(0)
        return buf.read()
    except Exception as e:
        print(f"Error generating voiceover: {e}")
        return None

def generate_ppt_content(topic):
    """Use Gemini to generate structured content for a 10-slide PPT."""
    try:
        model = genai.GenerativeModel('gemini-3-flash-preview')
        prompt = (
            f"Create a structured 10-slide PowerPoint presentation about '{topic}'. "
            "For each slide provide: SLIDE_TITLE: <title> and SLIDE_CONTENT: <3-5 bullet points>. "
            "Separate slides with '---'."
        )
        response = model.generate_content(prompt)
        return response.text
    except Exception as e:
        print(f"Error generating PPT content: {e}")
        return None

def create_ppt_file(topic, content_text):
    """Create a PPTX file from generated content text and return bytes."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from io import BytesIO
        prs = Presentation()
        slide_layout = prs.slide_layouts[1]
        slides_raw = content_text.split("---")
        for slide_raw in slides_raw:
            if not slide_raw.strip():
                continue
            title_line = ""
            content_lines = []
            for line in slide_raw.strip().splitlines():
                if line.startswith("SLIDE_TITLE:"):
                    title_line = line.replace("SLIDE_TITLE:", "").strip()
                elif line.startswith("SLIDE_CONTENT:"):
                    content_lines.append(line.replace("SLIDE_CONTENT:", "").strip())
                elif line.startswith("-") or line.startswith("•"):
                    content_lines.append(line.strip("- •").strip())
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = title_line or topic
            tf = slide.placeholders[1].text_frame
            tf.clear()
            for bullet in content_lines:
                p = tf.add_paragraph()
                p.text = bullet
        buf = BytesIO()
        prs.save(buf)
        buf.seek(0)
        return buf.read()
    except Exception as e:
        print(f"Error creating PPT file: {e}")
        return None

def generate_ppt_from_context(topic, context):
    """Generate PowerPoint content using ONLY the provided context."""
    try:
        model = genai.GenerativeModel('gemini-3-flash-preview')
        prompt = f"""Based ONLY on the following context from uploaded materials, create a structured 10-slide PowerPoint presentation about '{topic}'.

Context from uploaded materials:
{context}

For each slide provide:
SLIDE_TITLE: <title>
SLIDE_CONTENT: <3-5 bullet points based ONLY on the context>

Separate slides with '---'."""
        
        response = model.generate_content(prompt)
        return response.text
    except Exception as e:
        print(f"Error generating PPT content: {e}")
        return None

# ==============================================================================
# --- 8. WORKSHEET GENERATION FUNCTIONS ---
# ==============================================================================

def ensure_worksheet_format(content):
    """Ensure the worksheet content has proper QUESTIONS and ANSWERS sections."""
    if not content:
        return content
        
    # If markers already exist, return as is
    if "---QUESTIONS---" in content and "---ANSWERS---" in content:
        return content
    
    # Try to intelligently split if markers are missing
    lines = content.split('\n')
    questions = []
    answers = []
    in_answers = False
    found_answer_section = False
    
    # Common answer section indicators
    answer_indicators = ['answer key', 'answers:', 'solutions:', 'answer:', 
                         'answers', 'solutions', 'answer key:', 'ANSWER KEY']
    
    for line in lines:
        # Check if this line indicates start of answers section
        if any(indicator.lower() in line.lower() for indicator in answer_indicators):
            in_answers = True
            found_answer_section = True
            # Don't add the indicator line to answers
            continue
        
        if not in_answers:
            if line.strip():  # Only add non-empty lines
                questions.append(line)
        else:
            if line.strip():  # Only add non-empty lines
                answers.append(line)
    
    # If we found an answer section, format properly
    if found_answer_section and questions and answers:
        return f"---QUESTIONS---\n{''.join(questions)}\n\n---ANSWERS---\n{''.join(answers)}"
    
    # If no clear separation, try to split by common patterns
    if not found_answer_section:
        # Look for patterns like "1." and then later "Answers:" 
        question_pattern = r'^\d+\.'
        answer_pattern = r'answers|solutions'
        
        temp_questions = []
        temp_answers = []
        current_section = "questions"
        
        for line in lines:
            if re.search(answer_pattern, line.lower()) and len(temp_questions) > 5:
                current_section = "answers"
                continue
            
            if current_section == "questions" and line.strip():
                temp_questions.append(line)
            elif current_section == "answers" and line.strip():
                temp_answers.append(line)
        
        if temp_questions and temp_answers:
            return f"---QUESTIONS---\n{''.join(temp_questions)}\n\n---ANSWERS---\n{''.join(temp_answers)}"
    
    # If all else fails, put everything in questions and create a simple answer section
    return f"---QUESTIONS---\n{content}\n\n---ANSWERS---\nPlease refer to the source materials for answers."

def generate_worksheet_parts(topic, context, worksheet_type='descriptive', question_count=10):
    """Generate worksheet content with separate questions and answers using ONLY the provided context.
    
    Args:
        topic: The worksheet topic
        context: The context from uploaded materials
        worksheet_type: 'mcq' for multiple choice or 'descriptive' for long answer questions
        question_count: Number of questions to generate (5, 10, or 15)
    """
    try:
        model = genai.GenerativeModel('gemini-3-flash-preview')
        
        if worksheet_type == 'mcq':
            question_instruction = f"""Create {question_count} MULTIPLE CHOICE QUESTIONS with 4 options (A, B, C, D) each.
Format each question as:
Q1. [Question text]
A) [Option A]
B) [Option B]
C) [Option C]
D) [Option D]
Answer: [Correct option letter]"""
        else:
            question_instruction = f"""Create {question_count} DESCRIPTIVE/LONG ANSWER QUESTIONS that require detailed explanations.
Format each question as:
Q1. [Question text]
(Expect 2-3 paragraph answers)"""
        
        prompt = f"""Based ONLY on the following context from uploaded materials, create a worksheet about '{topic}'.

Context from uploaded materials:
{context}

{question_instruction}

IMPORTANT: Format your response EXACTLY as follows with these exact markers:

---QUESTIONS---
Q1. [First question]
Q2. [Second question]
... (continue for all {question_count} questions)

---ANSWERS---
Answer 1: [Answer to Q1 with explanation]
Answer 2: [Answer to Q2 with explanation]
... (continue for all {question_count} answers)

Make sure:
- Exactly {question_count} questions are generated
- Each question and answer is numbered
- Questions and answers are clearly separated with the markers above
- For MCQ: Include all 4 options and indicate the correct answer
- For Descriptive: Provide comprehensive answers with explanations
- Use ONLY information from the context provided
- Do not include any other text outside these sections"""
        
        response = model.generate_content(prompt)
        content = response.text
        
        # Log the content for debugging
        print(f"Generated {worksheet_type.upper()} worksheet ({question_count} questions): {content[:200]}...")
        
        # Ensure proper formatting
        content = ensure_worksheet_format(content)
        
        return content
    except Exception as e:
        print(f"Error generating worksheet parts: {e}")
        return None

def create_worksheet_pdfs(topic, worksheet_text):
    """Create two PDF files (questions and answers) from generated content."""
    try:
        from fpdf import FPDF
        from fpdf.enums import XPos, YPos
        import textwrap
        
        # Split into questions and answers
        questions_text = ""
        answers_text = ""
        
        if "---QUESTIONS---" in worksheet_text and "---ANSWERS---" in worksheet_text:
            parts = worksheet_text.split("---QUESTIONS---")
            if len(parts) > 1:
                qa_parts = parts[1].split("---ANSWERS---")
                if len(qa_parts) > 1:
                    questions_text = qa_parts[0].strip()
                    answers_text = qa_parts[1].strip()
        
        # If parsing fails, use the whole text as questions
        if not questions_text:
            questions_text = worksheet_text
            answers_text = "Answer key not available. Please refer to the source materials."
        
        # Create Questions PDF with better formatting
        questions_pdf = FPDF()
        questions_pdf.add_page()
        
        # Set margins (left, top, right) - give more space
        questions_pdf.set_left_margin(15)
        questions_pdf.set_right_margin(15)
        questions_pdf.set_auto_page_break(auto=True, margin=15)
        
        # Add title - using modern FPDF2 syntax with new_x/new_y instead of ln
        questions_pdf.set_font("Helvetica", style='B', size=16)
        questions_pdf.cell(0, 10, text=f"Worksheet: {topic}", align='C', new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        questions_pdf.ln(5)
        
        # Add subtitle
        questions_pdf.set_font("Helvetica", style='B', size=14)
        questions_pdf.cell(0, 10, text="QUESTIONS", align='C', new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        questions_pdf.ln(10)
        
        # Add questions content with proper text wrapping
        questions_pdf.set_font("Helvetica", size=11)
        
        # Process each line and handle long text
        for line in questions_text.split('\n'):
            if not line.strip():
                questions_pdf.ln(5)
                continue
            
            # Remove any problematic characters
            safe_line = ''.join(char if ord(char) < 128 else '?' for char in line)
            
            # Handle numbered questions (like "Q1." or "1.")
            if re.match(r'^Q?\d+[\.\)]', safe_line.strip()) or re.match(r'^\d+[\.\)]', safe_line.strip()):
                questions_pdf.set_font("Helvetica", style='B', size=11)
                # Split long lines if needed
                wrapped_lines = textwrap.wrap(safe_line, width=80)
                for wrapped_line in wrapped_lines:
                    questions_pdf.cell(0, 6, text=wrapped_line, new_x=XPos.LMARGIN, new_y=YPos.NEXT)
                questions_pdf.set_font("Helvetica", size=11)
            else:
                # For options in MCQ (A) B) C) D))
                if safe_line.strip() and len(safe_line.strip()) > 0 and safe_line.strip()[0] in ['A', 'B', 'C', 'D'] and ')' in safe_line[:3]:
                    questions_pdf.set_font("Helvetica", style='I', size=11)
                    # Indent options with spaces
                    wrapped_lines = textwrap.wrap("    " + safe_line, width=80)
                    for wrapped_line in wrapped_lines:
                        questions_pdf.cell(0, 5, text=wrapped_line, new_x=XPos.LMARGIN, new_y=YPos.NEXT)
                    questions_pdf.set_font("Helvetica", size=11)
                else:
                    # Regular text
                    wrapped_lines = textwrap.wrap(safe_line, width=80)
                    for wrapped_line in wrapped_lines:
                        questions_pdf.cell(0, 5, text=wrapped_line, new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        
        # Create Answers PDF with better formatting
        answers_pdf = FPDF()
        answers_pdf.add_page()
        
        # Set margins
        answers_pdf.set_left_margin(15)
        answers_pdf.set_right_margin(15)
        answers_pdf.set_auto_page_break(auto=True, margin=15)
        
        # Add title
        answers_pdf.set_font("Helvetica", style='B', size=16)
        answers_pdf.cell(0, 10, text=f"Worksheet: {topic}", align='C', new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        answers_pdf.ln(5)
        
        # Add subtitle
        answers_pdf.set_font("Helvetica", style='B', size=14)
        answers_pdf.cell(0, 10, text="ANSWER KEY", align='C', new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        answers_pdf.ln(10)
        
        # Add answers content
        answers_pdf.set_font("Helvetica", size=11)
        
        for line in answers_text.split('\n'):
            if not line.strip():
                answers_pdf.ln(5)
                continue
            
            safe_line = ''.join(char if ord(char) < 128 else '?' for char in line)
            
            # Handle answer numbering
            if re.match(r'^(Answer|answer|Ans|ans)\s*\d+', safe_line) or re.match(r'^\d+\.', safe_line):
                answers_pdf.set_font("Helvetica", style='B', size=11)
                wrapped_lines = textwrap.wrap(safe_line, width=80)
                for wrapped_line in wrapped_lines:
                    answers_pdf.cell(0, 6, text=wrapped_line, new_x=XPos.LMARGIN, new_y=YPos.NEXT)
                answers_pdf.set_font("Helvetica", size=11)
            else:
                wrapped_lines = textwrap.wrap(safe_line, width=80)
                for wrapped_line in wrapped_lines:
                    answers_pdf.cell(0, 5, text=wrapped_line, new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        
        # Get PDF bytes using output()
        questions_bytes = questions_pdf.output()
        answers_bytes = answers_pdf.output()
        
        # Convert to bytes if needed
        if isinstance(questions_bytes, str):
            questions_bytes = questions_bytes.encode('latin-1')
        if isinstance(answers_bytes, str):
            answers_bytes = answers_bytes.encode('latin-1')
        
        return questions_bytes, answers_bytes
        
    except Exception as e:
        print(f"Error creating worksheet PDFs: {e}")
        traceback.print_exc()
        return None, None

def create_worksheet_text_files(topic, worksheet_text):
    """Create text files as fallback if PDF generation fails"""
    try:
        # Split into questions and answers
        questions_text = ""
        answers_text = ""
        
        if "---QUESTIONS---" in worksheet_text and "---ANSWERS---" in worksheet_text:
            parts = worksheet_text.split("---QUESTIONS---")
            if len(parts) > 1:
                qa_parts = parts[1].split("---ANSWERS---")
                if len(qa_parts) > 1:
                    questions_text = qa_parts[0].strip()
                    answers_text = qa_parts[1].strip()
        
        # If parsing failed, create a simple format
        if not questions_text:
            questions_text = "Questions could not be parsed from the generated content."
            answers_text = "Answer key not available."
        
        # Create questions text file
        questions_content = f"WORKSHEET: {topic}\n"
        questions_content += "=" * 50 + "\n"
        questions_content += "QUESTIONS\n"
        questions_content += "=" * 50 + "\n\n"
        questions_content += questions_text
        
        # Create answers text file
        answers_content = f"WORKSHEET: {topic}\n"
        answers_content += "=" * 50 + "\n"
        answers_content += "ANSWER KEY\n"
        answers_content += "=" * 50 + "\n\n"
        answers_content += answers_text
        
        # Convert to bytes
        questions_bytes = questions_content.encode('utf-8')
        answers_bytes = answers_content.encode('utf-8')
        
        return questions_bytes, answers_bytes
        
    except Exception as e:
        print(f"Error creating text files: {e}")
        return None, None

def generate_worksheet_from_context(topic, context):
    """Generate worksheet content using ONLY the provided context."""
    try:
        model = genai.GenerativeModel('gemini-3-flash-preview')
        prompt = f"""Based ONLY on the following context from uploaded materials, create a comprehensive worksheet about '{topic}'.

Context from uploaded materials:
{context}

Create a worksheet with:
1. 5-10 questions (mix of multiple choice, short answer, and true/false)
2. An answer key at the end
3. Use ONLY information from the context provided

Worksheet:"""
        
        response = model.generate_content(prompt)
        return response.text
    except Exception as e:
        print(f"Error generating worksheet: {e}")
        return None

def generate_quiz_5_questions(topic, context=None):
    """Generate a 5-question quiz based on topic and uploaded materials if available."""
    try:
        model = genai.GenerativeModel('gemini-2.0-flash')
        
        if context:
            prompt = f"""Based ONLY on the following context from uploaded materials, create a 5-question quiz about '{topic}'.

Context from uploaded materials:
{context}

Create exactly 5 quiz questions with:
1. Question format: Mix of multiple choice, short answer, and true/false
2. Each question clearly numbered (1-5)
3. For multiple choice questions, provide options (A, B, C, D)
4. Include the correct answer at the end for each question
5. Use ONLY information from the context provided

Quiz:"""
        else:
            prompt = f"""Create a 5-question quiz about '{topic}' as a learning assessment.

Create exactly 5 quiz questions with:
1. Question format: Mix of multiple choice, short answer, and true/false
2. Each question clearly numbered (1-5)
3. For multiple choice questions, provide options (A, B, C, D)
4. Include the correct answer at the end for each question
5. Make questions appropriate for students learning this topic

Quiz:"""
        
        response = model.generate_content(prompt)
        return response.text
    except Exception as e:
        print(f"Error generating quiz: {e}")
        return None

def create_worksheet_file(topic, content):
    """Create a PDF worksheet file from content."""
    try:
        from fpdf import FPDF
        from fpdf.enums import XPos, YPos
        from io import BytesIO
        
        pdf = FPDF()
        pdf.set_margins(left=10, top=10, right=10)
        pdf.add_page()
        pdf.set_font("Arial", size=12)
        
        # Add title - using new syntax
        pdf.set_font("Arial", 'B', 16)
        pdf.cell(0, 10, txt=f"Worksheet: {topic}", new_x=XPos.LMARGIN, new_y=YPos.NEXT, align='C')
        pdf.ln(10)
        
        # Add content
        pdf.set_font("Arial", size=12)
        
        # Split content into lines and add to PDF with proper wrapping
        for line in content.split('\n'):
            # Handle encoding issues
            try:
                # Encode to latin-1 to handle special characters
                line = line.encode('latin-1', 'replace').decode('latin-1')
                pdf.multi_cell(0, 10, txt=line, align='L')
            except:
                # If encoding fails, try a simpler approach
                safe_line = ''.join(char if ord(char) < 128 else '?' for char in line)
                pdf.multi_cell(0, 10, txt=safe_line, align='L')
        
        # Get PDF bytes
        pdf_bytes = pdf.output()
        if isinstance(pdf_bytes, str):
            pdf_bytes = pdf_bytes.encode('latin-1')
        return pdf_bytes
        
    except Exception as e:
        print(f"Error creating worksheet file: {e}")
        return None

# ==============================================================================
# --- 9. GOOGLE DRIVE UPLOAD FUNCTION (UPDATED WITH OAUTH) ---
# ==============================================================================

def upload_file_to_drive_with_folder_check(file_bytes, filename, target_folder_id=None):
    """
    Upload a file to Google Drive using OAuth 2.0 authentication.
    Returns (file_id, success_message, error_message)
    """
    if not drive_service:
        return None, None, "Google Drive API not initialized"
    
    temp_file_path = None
    try:
        # Create temporary file
        with tempfile.NamedTemporaryFile(suffix=os.path.splitext(filename)[1], delete=False) as temp_file:
            temp_file.write(file_bytes)
            temp_file.flush()
            os.fsync(temp_file.fileno())
            temp_file_path = temp_file.name
        
        # Determine MIME type
        mime_types = {
            '.pdf': 'application/pdf',
            '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
            '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            '.txt': 'text/plain',
            '.jpg': 'image/jpeg',
            '.jpeg': 'image/jpeg',
            '.png': 'image/png',
        }
        
        file_ext = os.path.splitext(filename)[1].lower()
        mime_type = mime_types.get(file_ext, 'application/octet-stream')
        
        # Prepare file metadata
        file_metadata = {'name': filename}
        
        # Add to Shared Drive if folder ID provided
        if target_folder_id and target_folder_id.strip():
            file_metadata['parents'] = [target_folder_id]
            print(f"📁 Uploading '{filename}' to folder: {target_folder_id}")
        
        # Upload file
        media = MediaFileUpload(temp_file_path, mimetype=mime_type, resumable=True)
        
        try:
            file = drive_service.files().create(
                body=file_metadata,
                media_body=media,
                fields='id, name, webViewLink',
                supportsAllDrives=True  # Still needed for Shared Drives
            ).execute()
            
            file_id = file.get('id')
            web_link = file.get('webViewLink', f"https://drive.google.com/file/d/{file_id}/view")
            
            print(f"✅ File uploaded successfully! ID: {file_id}")
            
            # Make file readable by anyone with link (optional)
            try:
                permission = {
                    'type': 'anyone',
                    'role': 'reader'
                }
                drive_service.permissions().create(
                    fileId=file_id,
                    body=permission,
                    supportsAllDrives=True
                ).execute()
                print(f"🔓 File made publicly accessible")
            except Exception as perm_error:
                print(f"⚠️ Could not set public permission: {perm_error}")
            
            return file_id, f"✅ File '{filename}' uploaded successfully", None
            
        except HttpError as e:
            error_message = str(e)
            print(f"❌ Drive upload error: {error_message}")
            return None, None, f"Drive API error: {error_message}"
        
    except Exception as e:
        error_message = f"Unexpected error: {str(e)}"
        print(f"❌ {error_message}")
        return None, None, error_message
        
    finally:
        # Clean up temporary file with retry mechanism
        if temp_file_path and os.path.exists(temp_file_path):
            try:
                # Close any potential file handles
                if 'media' in locals():
                    if hasattr(media, '_fd') and media._fd:
                        media._fd.close()
                
                # Force garbage collection
                gc.collect()
                
                # Small delay to ensure file is released
                time.sleep(0.5)
                
                os.remove(temp_file_path)
                print(f"✅ Temp file deleted: {temp_file_path}")
            except Exception as e:
                print(f"⚠️ Could not delete temp file: {e}")

# ==============================================================================
# --- 10. GOOGLE CLASSROOM FUNCTIONS ---
# ==============================================================================

def post_announcement_with_attachments(course_id, text, file_ids=None):
    """
    Post an announcement to Google Classroom with optional file attachments.
    Note: Announcements don't have titles, only text.
    """
    if not classroom_service:
        return {"result": "Google Classroom API not initialized.", "source": "api_error"}
    
    try:
        # Correct format for announcement (no title field)
        announcement_body = {
            'text': text,  # Only 'text' field, no 'title'
            'state': 'PUBLISHED'
        }
        
        # Add materials if file_ids provided
        if file_ids and len(file_ids) > 0:
            materials = []
            for file_id in file_ids:
                materials.append({
                    'driveFile': {
                        'driveFile': {
                            'id': file_id
                        }
                    }
                })
            announcement_body['materials'] = materials
        
        # Post announcement
        announcement = classroom_service.courses().announcements().create(
            courseId=course_id,
            body=announcement_body
        ).execute()
        
        announcement_id = announcement.get('id')
        print(f"✓ Announcement posted successfully: {announcement_id}")
        
        return {
            "result": f"✅ Successfully posted announcement in Google Classroom!",
            "source": "classroom_success",
            "announcement_id": announcement_id
        }
        
    except HttpError as api_error:
        error_content = api_error.content if hasattr(api_error, 'content') else str(api_error)
        print(f"Google Classroom API Error: {error_content}")
        return {
            "result": "❌ Couldn't post announcement. Please check permissions and Course ID.",
            "source": "api_error",
            "error": str(api_error)
        }
    except Exception as e:
        print(f"Unexpected error posting announcement: {e}")
        return {
            "result": "❌ An unexpected error occurred while posting the announcement.",
            "source": "error",
            "error": str(e)
        }

def post_assignment_with_attachments(course_id, title, description, file_ids=None, due_date=None):
    """
    Create a classroom assignment with optional file attachments
    Note: Assignments DO have titles
    """
    if not classroom_service:
        return {"result": "Google Classroom API not initialized.", "source": "api_error"}
    
    try:
        coursework = {
            'title': title,  # Title is valid for assignments
            'description': description,
            'workType': 'ASSIGNMENT',
            'state': 'PUBLISHED'
        }
        
        # Add due date if provided
        if due_date:
            coursework['dueDate'] = {
                'year': due_date.year,
                'month': due_date.month,
                'day': due_date.day
            }
        
        # Add materials if file_ids provided
        if file_ids and len(file_ids) > 0:
            materials = []
            for file_id in file_ids:
                materials.append({
                    'driveFile': {
                        'driveFile': {
                            'id': file_id
                        }
                    }
                })
            coursework['materials'] = materials
        
        # Create assignment
        assignment = classroom_service.courses().courseWork().create(
            courseId=course_id,
            body=coursework
        ).execute()
        
        assignment_id = assignment.get('id')
        assignment_title = assignment.get('title')
        print(f"✓ Assignment created successfully: {assignment_title} (ID: {assignment_id})")
        
        return {
            "result": f"✅ Successfully created assignment '{assignment_title}' in Google Classroom!",
            "source": "classroom_success",
            "assignment_id": assignment_id,
            "assignment_url": f"https://classroom.google.com/c/{course_id}/a/{assignment_id}"
        }
        
    except HttpError as api_error:
        error_content = api_error.content if hasattr(api_error, 'content') else str(api_error)
        print(f"Google Classroom API Error: {error_content}")
        return {
            "result": "❌ Couldn't create Classroom assignment. Please check permissions and Course ID.",
            "source": "api_error",
            "error": str(api_error)
        }
    except Exception as e:
        print(f"General error creating Classroom assignment: {e}")
        return {
            "result": "❌ An unexpected error occurred while creating the Classroom assignment.",
            "source": "error",
            "error": str(e)
        }

def post_to_classroom(announcement_title, announcement_text, file_bytes=None, filename=None):
    """Post an announcement to Google Classroom with optional file attachment."""
    try:
        if not classroom_service:
            return False, "Classroom service not initialized"
        
        # Create announcement with text
        announcement = {
            "text": announcement_text,
            "state": "PUBLISHED"
        }
        
        # If file provided, upload to Drive and attach
        if file_bytes and filename:
            try:
                from googleapiclient.http import MediaInMemoryUpload
                media_upload = MediaInMemoryUpload(file_bytes, mimetype='application/pdf', resumable=True)
                drive_file = drive_service.files().create(
                    body={
                        'name': filename,
                        'parents': [GOOGLE_FORMS_TARGET_FOLDER_ID]
                    },
                    media_body=media_upload,
                    supportsAllDrives=True  # Critical for Shared Drives
                ).execute()
                
                file_id = drive_file.get('id')
                
                # Add attachment to announcement
                announcement['materials'] = [{
                    'driveFile': {
                        'driveFile': {
                            'id': file_id
                        }
                    }
                }]
            except Exception as e:
                print(f"Error uploading file to Drive: {e}")
                # Continue without file attachment
                pass
        
        # Post to classroom
        result = classroom_service.courses().announcements().create(
            courseId=CLASSROOM_COURSE_ID,
            body={
                'title': announcement_title,
                'text': announcement_text,
                'state': 'PUBLISHED',
                'materials': announcement.get('materials', [])
            }
        ).execute()
        
        return True, f"Posted to classroom: {result.get('id', 'Unknown')}"
        
    except Exception as e:
        print(f"Error posting to classroom: {e}")
        return False, str(e)

# ==============================================================================
# --- 12. MESSAGE HANDLING FUNCTIONS ---
# ==============================================================================


def end_conversation_and_show_menu(sender, final_message, responder):
    if final_message:
        responder.text(sender, final_message)
    clear_user_session(sender)
    menu_text = "What would you like to do next?"
    options = [
        "Ask Question", "Create Worksheet", "Create PPT", "Upload Materials",
        "View Uploaded Files", "Podcast from Image", "Summary from Image"
    ]
    responder.menu(sender, menu_text, options)

def handle_worksheet_with_drive_upload(sender, responder, topic, questions_bytes, answers_bytes):
    """Upload worksheet files to Google Drive using the common upload function"""
    try:
        responder.text(sender, "📤 Uploading worksheet files to Google Drive...")
        
        # Create filenames
        questions_filename = f"{topic.replace(' ', '_')}_worksheet_questions.pdf"
        answers_filename = f"{topic.replace(' ', '_')}_worksheet_answers.pdf"
        
        # Upload files
        q_id, q_msg, q_err = upload_file_to_drive_with_folder_check(
            questions_bytes, questions_filename, GOOGLE_FORMS_TARGET_FOLDER_ID
        )
        
        a_id, a_msg, a_err = upload_file_to_drive_with_folder_check(
            answers_bytes, answers_filename, GOOGLE_FORMS_TARGET_FOLDER_ID
        )
        
        if q_id and a_id:
            responder.text(sender, f"✅ Questions uploaded to Drive\n✅ Answer key uploaded to Drive")
        elif q_id or a_id:
            responder.text(sender, "⚠️ Partially uploaded to Drive\n✅ Files are available for download from your device")
        else:
            responder.text(sender, "⚠️ Could not upload to Drive\n✅ Files are available for download from your device")
        
        end_conversation_and_show_menu(sender, "✅ Worksheet generated successfully!", responder)
        
    except Exception as e:
        print(f"Error in handle_worksheet_with_drive_upload: {e}")
        responder.text(sender, "✅ Worksheet generated successfully! Download from your device.")
        end_conversation_and_show_menu(sender, None, responder)

def generate_and_send_worksheet(sender, topic, worksheet_type, question_count, relevant_content, responder):
    """Helper function to generate and send worksheet with proper error handling"""
    try:
        # Generate worksheet with type and count parameters
        worksheet_full = generate_worksheet_parts(topic, relevant_content, worksheet_type, question_count)
        
        if not worksheet_full:
            responder.text(sender, f"❌ Couldn't generate worksheet for '{topic}'. Please try again.")
            end_conversation_and_show_menu(sender, None, responder)
            return
        
        # Try to create PDFs first
        questions_bytes, answers_bytes = create_worksheet_pdfs(topic, worksheet_full)
        
        if questions_bytes and answers_bytes:
            # Success with PDFs
            responder.text(sender, f"📝 Your {worksheet_type.upper()} worksheet ({question_count} questions) is ready:")
            
            # Send questions PDF
            questions_filename = f"{topic.replace(' ', '_')}_worksheet_questions.pdf"
            responder.document(sender, questions_bytes, questions_filename)
            
            # Small delay
            time.sleep(1)
            
            # Send answers PDF
            answers_filename = f"{topic.replace(' ', '_')}_worksheet_answers.pdf"
            responder.document(sender, answers_bytes, answers_filename)
            
            # Store PDFs in temp_data for potential classroom upload
            update_user_state(sender, state="awaiting_classroom_upload_decision", temp_data={
                'questions_bytes': base64.b64encode(questions_bytes).decode('utf-8'),
                'answers_bytes': base64.b64encode(answers_bytes).decode('utf-8'),
                'questions_filename': questions_filename,
                'answers_filename': answers_filename,
                'topic': topic,
                'worksheet_type': worksheet_type,
                'question_count': question_count
            })
            
            # Ask about classroom upload
            responder.menu(sender, "Would you like to upload this worksheet to Google Classroom?", ["Yes, post to Classroom", "No, just download"])
        else:
            # Fallback to text files if PDF fails
            responder.text(sender, "⚠️ PDF generation had issues. Sending as text files instead.")
            
            questions_bytes, answers_bytes = create_worksheet_text_files(topic, worksheet_full)
            
            if questions_bytes and answers_bytes:
                questions_filename = f"{topic.replace(' ', '_')}_worksheet_questions.txt"
                responder.document(sender, questions_bytes, questions_filename)
                
                time.sleep(1)
                
                answers_filename = f"{topic.replace(' ', '_')}_worksheet_answers.txt"
                responder.document(sender, answers_bytes, answers_filename)
                
                end_conversation_and_show_menu(sender, "✅ Worksheet sent as text files.", responder)
            else:
                # Last resort: send as text message
                responder.text(sender, f"📝 Here's your worksheet:\n\n{worksheet_full}")
                end_conversation_and_show_menu(sender, None, responder)
                
    except Exception as e:
        print(f"Error in worksheet generation: {e}")
        responder.text(sender, "❌ An error occurred while generating the worksheet.")
        end_conversation_and_show_menu(sender, None, responder)

def process_message(sender, text, responder=None):
    # If no responder provided, use WhatsApp responder
    if responder is None:
        responder = WhatsAppResponder(ACCESS_TOKEN, PHONE_NUMBER_ID)

    initialize_user_session(sender)
    append_to_memory(sender, "user", text)
    current_state = user_states.get(sender)

    # --- Universal Cancel/Menu ---
    if text.lower().strip() in ['cancel', 'stop', 'menu', 'start', 'exit']:
        if current_state:
            clear_user_session(sender)
            responder.text(sender, "Okay, I've canceled the current operation.")
        menu_text = "Hello! I'm your AI teaching assistant. Please choose an option:"
        options = [
            "Ask Question", "Create Worksheet", "Create PPT", "Upload Materials",
            "View Uploaded Files", "Podcast from Image", "Summary from Image"
        ]
        responder.menu(sender, menu_text, options)
        append_to_memory(sender, "assistant", "Displayed main menu.")
        return

    # --- Initial Triggers ---
    if not current_state:
        lower_text = text.lower().strip()
        if lower_text == "ask question":
            update_user_state(sender, state="awaiting_question")
            responder.text(sender, "Of course! What is your question? I'll search through your uploaded documents ONLY (no external knowledge will be used).")
            return
        if lower_text == "take quiz":
            update_user_state(sender, state="awaiting_quiz_topic")
            responder.text(sender, "Great! What topic would you like to be quizzed on? I'll create 5 questions for you.")
            return
        if lower_text == "create worksheet":
            update_user_state(sender, state="awaiting_worksheet_type")
            responder.menu(sender, "What type of questions would you like?", ["MCQ (Multiple Choice)", "Descriptive (Long Answer)"])
            return
        if lower_text == "create ppt":
            update_user_state(sender, state="awaiting_ppt_topic")
            responder.text(sender, "Excellent! What topic would you like the presentation to be about? I'll use ONLY your uploaded materials.")
            return
        if lower_text == "upload materials" or lower_text == "upload material":
            update_user_state(sender, state="awaiting_material_file")
            responder.text(sender, "Please send the file you'd like to upload (PDF, Word, PowerPoint, or Text). This will be added to my knowledge base for future questions!")
            return
        if lower_text == "view uploaded files":
            vector_db_files = get_uploaded_files_list()
            if vector_db_files:
                files_list = [f"• {os.path.basename(f)}" for f in vector_db_files[:10]]
                responder.text(sender, "Files in your knowledge base:\n" + "\n".join(files_list))
                if len(vector_db_files) > 10:
                    responder.text(sender, f"... and {len(vector_db_files) - 10} more files")
            else:
                responder.text(sender, "No files in knowledge base. Use 'Upload Materials' to add files.")
            end_conversation_and_show_menu(sender, None, responder)
            return
        if lower_text == "podcast from image":
            update_user_state(sender, state="awaiting_podcast_image")
            responder.text(sender, "Please send me an image of the text you'd like to convert to a podcast.")
            return
        if lower_text == "summary from image":
            update_user_state(sender, state="awaiting_summary_image")
            responder.text(sender, "Please send me an image of the text you'd like me to summarize.")
            return

    # --- State-Based Flows (ALL USING ONLY UPLOADED DOCUMENTS) ---
    
    # Question Answering Flow - USING ONLY UPLOADED DOCUMENTS
    if current_state == "awaiting_question":
        # Use strict RAG that ONLY uses vector DB - NO GEMINI FALLBACK
        response = query_dynamic_rag_strict(text, similarity_threshold=0.4)
        
        if response.get("result"):
            # Found relevant content in uploaded documents
            end_conversation_and_show_menu(sender, response["result"], responder)
        else:
            # No relevant content found in uploaded documents
            no_content_msg = "I couldn't find any relevant information in your uploaded materials to answer this question. Please upload relevant documents first using the 'Upload Materials' option, or try a different question based on your existing materials."
            end_conversation_and_show_menu(sender, no_content_msg, responder)
        return

    # Quiz Flow - Generate 5 Questions
    if current_state == "awaiting_quiz_topic":
        try:
            topic = text.strip()
            update_user_state(sender, temp_data={'quiz_topic': topic})
            
            # Try to find context from uploaded materials
            context = None
            try:
                rag_response = query_dynamic_rag(topic, similarity_threshold=0.4)
                if rag_response and rag_response.get("result"):
                    context = rag_response["result"]
            except:
                pass
            
            # Generate 5-question quiz
            quiz_content = generate_quiz_5_questions(topic, context=context)
            
            if quiz_content:
                responder.text(sender, f"📝 Here's your 5-question quiz on '{topic}':\n\n{quiz_content}")
                end_conversation_and_show_menu(sender, None, responder)
            else:
                responder.text(sender, "❌ Could not generate quiz. Please try again.")
                end_conversation_and_show_menu(sender, None, responder)
        except Exception as e:
            print(f"Error in quiz generation: {e}")
            responder.text(sender, "❌ An error occurred while generating the quiz.")
            end_conversation_and_show_menu(sender, None, responder)
        return

    # PPT Creation Flow - USING ONLY UPLOADED DOCUMENTS
    if current_state == "awaiting_ppt_topic":
        try:
            topic = text.strip()
            update_user_state(sender, temp_data={'ppt_topic': topic})
            
            # First check if we have relevant content in vector DB
            relevant_content = query_vector_db_for_topic(topic, similarity_threshold=0.4)
            
            if relevant_content:
                responder.text(sender, f"Found relevant content in your uploaded materials. Creating a 10-slide presentation on '{topic}' using ONLY your documents...")
                
                # Generate PPT content using ONLY the context from uploaded documents
                ppt_content = generate_ppt_from_context(topic, relevant_content)
                
                if ppt_content:
                    ppt_bytes = create_ppt_file(topic, ppt_content)
                    if ppt_bytes:
                        responder.document(sender, ppt_bytes, f"{topic.replace(' ', '_')}.pptx")
                        end_conversation_and_show_menu(sender, "✅ I've sent the presentation based on your uploaded materials.", responder)
                    else:
                        end_conversation_and_show_menu(sender, "❌ Sorry, I generated the content but failed to create the PPT file.", responder)
                else:
                    end_conversation_and_show_menu(sender, f"❌ Couldn't generate presentation content for '{topic}' from your uploaded materials.", responder)
            else:
                # No relevant content found
                no_content_msg = f"I couldn't find any relevant information about '{topic}' in your uploaded materials to create a presentation. Please upload relevant documents first using the 'Upload Materials' option, or try a different topic that exists in your materials."
                end_conversation_and_show_menu(sender, no_content_msg, responder)
                
        except Exception as e:
            print(f"Error in PPT generation flow: {e}")
            end_conversation_and_show_menu(sender, "❌ An unexpected error occurred while creating the presentation.", responder)
        return

    # Worksheet Creation Flow - USING ONLY UPLOADED DOCUMENTS (with two PDFs)
    
    # Step 1: Await worksheet type selection (MCQ or Descriptive)
    if current_state == "awaiting_worksheet_type":
        question_type = text.strip().lower()
        
        # Determine if MCQ or Descriptive
        if "mcq" in question_type or "multiple" in question_type or "choice" in question_type:
            worksheet_type = "mcq"
        elif "descriptive" in question_type or "long" in question_type or "answer" in question_type:
            worksheet_type = "descriptive"
        else:
            responder.text(sender, "Please select either MCQ or Descriptive type.")
            responder.menu(sender, "What type of questions would you like?", ["MCQ (Multiple Choice)", "Descriptive (Long Answer)"])
            return
        
        # Move to next state: ask for number of questions
        update_user_state(sender, state="awaiting_worksheet_count", temp_data={'worksheet_type': worksheet_type})
        responder.menu(sender, "How many questions would you like?", ["5 Questions", "10 Questions", "15 Questions"])
        return
    
    # Step 2: Await worksheet count selection
    if current_state == "awaiting_worksheet_count":
        question_count_text = text.strip().lower()
        
        # Extract number
        if "5" in question_count_text:
            question_count = 5
        elif "10" in question_count_text:
            question_count = 10
        elif "15" in question_count_text:
            question_count = 15
        else:
            responder.text(sender, "Please select either 5, 10, or 15 questions.")
            responder.menu(sender, "How many questions would you like?", ["5 Questions", "10 Questions", "15 Questions"])
            return
        
        # Move to next state: ask for topic
        update_user_state(sender, state="awaiting_worksheet_topic_with_type", temp_data={'question_count': question_count})
        responder.text(sender, f"Great! Now, what topic would you like the worksheet to be about? I'll create {question_count} {user_temp_data[sender]['worksheet_type'].upper()} questions using ONLY your uploaded materials.")
        return
    
    # Step 3: Await topic and generate worksheet
    if current_state == "awaiting_worksheet_topic_with_type":
        try:
            topic = text.strip()
            worksheet_type = user_temp_data.get(sender, {}).get('worksheet_type', 'descriptive')
            question_count = user_temp_data.get(sender, {}).get('question_count', 10)
            
            # First check if we have relevant content in vector DB
            relevant_content = query_vector_db_for_topic(topic, similarity_threshold=0.4)
            
            if relevant_content:
                responder.text(sender, f"Found relevant content! Creating {question_count} {worksheet_type.upper()} questions on '{topic}'...")
                # Use the helper function to handle generation and sending with proper error handling
                generate_and_send_worksheet(sender, topic, worksheet_type, question_count, relevant_content, responder)
            else:
                # No relevant content found
                no_content_msg = f"I couldn't find relevant information about '{topic}' in your uploaded materials. Please upload relevant documents or try a different topic."
                end_conversation_and_show_menu(sender, no_content_msg, responder)
                
        except Exception as e:
            print(f"Error in worksheet generation flow: {e}")
            end_conversation_and_show_menu(sender, "❌ An error occurred while creating the worksheet.", responder)
        return

    # Classroom Upload Decision Flow (from generate_and_send_worksheet)
    if current_state == "awaiting_classroom_upload_decision":
        decision = text.strip().lower()
        print(f"DEBUG: Classroom decision flow - decision text: '{decision}'")
        
        if "classroom" in decision or "post" in decision:
            # Direct to Classroom post flow - CHECK THIS FIRST!
            print(f"DEBUG: Routing to Classroom posting (matched 'classroom' or 'post')")
            update_user_state(sender, state="awaiting_classroom_post_choice", temp_data={
                'questions_bytes': user_temp_data[sender].get('questions_bytes'),
                'answers_bytes': user_temp_data[sender].get('answers_bytes'),
                'questions_filename': user_temp_data[sender].get('questions_filename'),
                'answers_filename': user_temp_data[sender].get('answers_filename'),
                'topic': user_temp_data[sender].get('topic')
            })
            responder.menu(sender, "What would you like to post to Classroom?", 
                          ["Post Questions Only", "Post Answers Only", "Post Both"])
        elif "yes" in decision or "drive" in decision or "upload" in decision:
            # Upload to Drive
            print(f"DEBUG: Routing to Drive upload (matched 'yes', 'drive', or 'upload')")
            questions_bytes = base64.b64decode(user_temp_data.get(sender, {}).get('questions_bytes', ''))
            answers_bytes = base64.b64decode(user_temp_data.get(sender, {}).get('answers_bytes', ''))
            topic = user_temp_data.get(sender, {}).get('topic', 'Worksheet')
            
            handle_worksheet_with_drive_upload(sender, responder, topic, questions_bytes, answers_bytes)
        else:
            # Just download
            print(f"DEBUG: Just download (no match)")
            end_conversation_and_show_menu(sender, "✅ Worksheet generated successfully! Download from your device.", responder)
        return

    # Classroom Post Choice (what content to post)
    if current_state == "awaiting_classroom_post_choice":
        choice = text.strip().lower()
        print(f"DEBUG: User chose classroom post type: '{choice}'")
        
        questions_bytes = base64.b64decode(user_temp_data.get(sender, {}).get('questions_bytes', ''))
        answers_bytes = base64.b64decode(user_temp_data.get(sender, {}).get('answers_bytes', ''))
        questions_filename = user_temp_data.get(sender, {}).get('questions_filename', 'questions.pdf')
        answers_filename = user_temp_data.get(sender, {}).get('answers_filename', 'answers.pdf')
        topic = user_temp_data.get(sender, {}).get('topic', 'Worksheet')
        
        responder.text(sender, "📤 Posting to Google Classroom...")
        
        # Upload files to Drive first
        file_ids = []
        
        if "questions" in choice or "both" in choice:
            print(f"DEBUG: Uploading questions to Drive")
            q_id, q_msg, q_err = upload_file_to_drive_with_folder_check(
                questions_bytes, questions_filename, GOOGLE_FORMS_TARGET_FOLDER_ID
            )
            if q_id:
                file_ids.append(q_id)
                responder.text(sender, f"✅ Questions uploaded to Drive")
                print(f"DEBUG: Questions uploaded, ID: {q_id}")
        
        if "answers" in choice or "both" in choice:
            print(f"DEBUG: Uploading answers to Drive")
            a_id, a_msg, a_err = upload_file_to_drive_with_folder_check(
                answers_bytes, answers_filename, GOOGLE_FORMS_TARGET_FOLDER_ID
            )
            if a_id:
                file_ids.append(a_id)
                responder.text(sender, f"✅ Answers uploaded to Drive")
                print(f"DEBUG: Answers uploaded, ID: {a_id}")
        
        if file_ids:
            # Create announcement (NO TITLE, just text)
            if "questions" in choice and "answers" in choice:
                announcement_text = f"📚 New Worksheet: {topic}\n\nBoth questions and answer key are attached below."
            elif "questions" in choice:
                announcement_text = f"📚 New Worksheet: {topic}\n\nQuestions are attached below."
            else:
                announcement_text = f"📚 Answer Key: {topic}\n\nAnswer key is attached below."
            
            print(f"DEBUG: Posting announcement with {len(file_ids)} file(s) to Classroom")
            # Post to Classroom (announcements don't have titles)
            result = post_announcement_with_attachments(
                CLASSROOM_COURSE_ID, 
                announcement_text,
                file_ids
            )
            responder.text(sender, result['result'])
            print(f"DEBUG: Classroom posting result: {result}")
        else:
            responder.text(sender, "❌ Failed to upload files to Drive. Cannot post to Classroom.")
            print(f"DEBUG: No file IDs received from Drive upload")
        
        end_conversation_and_show_menu(sender, None, responder)
        return

    # Material Upload Flow
    if current_state == "awaiting_material_file":
        responder.text(sender, "Please send me the file you'd like to upload (PDF, Word, PowerPoint, or Text).")
        return

    # Podcast from Image Flow
    if current_state == "awaiting_podcast_image":
        responder.text(sender, "Please send me the image you'd like to convert to a podcast.")
        return

    # Summary from Image Flow
    if current_state == "awaiting_summary_image":
        responder.text(sender, "Please send me the image you'd like me to summarize.")
        return

    # Interactive confirmation states (if needed)
    if current_state == "awaiting_topic_confirmation":
        if text.lower().strip() in ['yes', 'y', 'ok', 'sure', 'confirm']:
            # User confirmed, proceed with generation
            topic = user_temp_data.get(sender, {}).get('pending_topic', '')
            if topic:
                update_user_state(sender, state="awaiting_ppt_topic")
                responder.text(sender, f"Great! Let's proceed with '{topic}'. I'll generate the presentation using your uploaded materials.")
                # Reprocess with the topic
                process_message(sender, topic, responder)
            else:
                end_conversation_and_show_menu(sender, "Sorry, I lost track of the topic. Please start over.", responder)
        else:
            # User wants a different topic
            update_user_state(sender, state="awaiting_ppt_topic")
            responder.text(sender, "Okay, please provide the topic you'd like for the presentation (must be from your uploaded materials).")
        return

    # Bulk upload confirmation
    if current_state == "awaiting_bulk_upload_confirmation":
        if text.lower().strip() in ['yes', 'y', 'ok', 'sure', 'confirm', 'upload']:
            if sender in user_temp_data and 'bulk_upload_items' in user_temp_data[sender]:
                items = user_temp_data[sender]['bulk_upload_items']
                if items:
                    responder.text(sender, f"📦 Processing {len(items)} files...")
                    handle_bulk_document_upload(sender, items, responder)
                else:
                    end_conversation_and_show_menu(sender, "No files to upload.", responder)
            else:
                end_conversation_and_show_menu(sender, "No files in upload queue.", responder)
        else:
            # Cancel bulk upload
            if sender in user_temp_data:
                user_temp_data[sender].pop('bulk_upload_items', None)
            end_conversation_and_show_menu(sender, "Bulk upload cancelled.", responder)
        return

    # --- Fallback for unhandled states or general queries ---
    print(f"Handling as a general query: '{text}' (current_state: {current_state})")
    
    # If we're in some state we don't recognize, clear it and treat as general query
    if current_state:
        print(f"Clearing unknown state: {current_state}")
        clear_user_session(sender)
    
    # For general queries, also use ONLY uploaded documents
    response = query_dynamic_rag_strict(text, similarity_threshold=0.4)
    
    # Check if we got a meaningful response from uploaded documents
    if response and response.get("result"):
        responder.text(sender, response["result"])
    else:
        responder.text(sender, "I can only answer questions based on your uploaded materials. Please upload relevant documents using the 'Upload Materials' option, or select an option from the menu.")
    
    end_conversation_and_show_menu(sender, None, responder)

# ==============================================================================
# --- 12. BACKGROUND TASK HANDLERS ---
# ==============================================================================

def handle_audio_task(sender, media_id, responder):
    print(f"Starting audio task for {sender}")
    responder.text(sender, "Got your voice message, listening... 🎤")
    if media_url := get_media_url(media_id):
        if audio_bytes := download_media_file(media_url):
            if transcribed_text := transcribe_audio(audio_bytes):
                process_message(sender, transcribed_text, responder)
            else:
                responder.text(sender, "Sorry, I couldn't understand your audio.")
        else:
            responder.text(sender, "Sorry, I couldn't download your voice message.")
    print(f"Audio task for {sender} finished.")

def handle_image_task(sender, media_id, prompt, responder):
    print(f"Starting image task for {sender}")
    responder.text(sender, "🖼️ Analyzing your image...")
    if media_url := get_media_url(media_id):
        if image_bytes := download_media_file(media_url):
            if len(image_bytes) > 20 * 1024 * 1024:
                responder.text(sender, "❌ Image is too large (over 20MB). Please send a smaller image.")
                end_conversation_and_show_menu(sender, None, responder)
                return
            if 'podcast' in prompt.lower() or 'audio' in prompt.lower():
                ocr_prompt = "Extract all the text from this image. Do not summarize or explain, just provide the raw text."
                ocr_response = query_gemini_vision(ocr_prompt, image_bytes)
                if ocr_response.get("source") == "gemini_vision":
                    extracted_text = ocr_response["result"]
                    if len(extracted_text.strip()) > 50:
                        if audio_bytes := generate_voiceover(extracted_text):
                            responder.audio(sender, audio_bytes, "podcast.mp3")
                            end_conversation_and_show_menu(sender, "✅ Here is your podcast from the image text!", responder)
                        else:
                            responder.text(sender, "❌ I extracted the text, but couldn't create the audio.")
                            end_conversation_and_show_menu(sender, None, responder)
                    else:
                        responder.text(sender, "❌ Not enough text found in the image to create a podcast.")
                        end_conversation_and_show_menu(sender, None, responder)
                else:
                    responder.text(sender, ocr_response["result"])
                    end_conversation_and_show_menu(sender, None, responder)
            else:
                response = query_gemini_vision(prompt, image_bytes)
                responder.text(sender, response["result"])
                end_conversation_and_show_menu(sender, None, responder)
        else:
            responder.text(sender, "❌ Sorry, I couldn't download the image. Please try again.")
            end_conversation_and_show_menu(sender, None, responder)
    else:
        responder.text(sender, "❌ Sorry, I couldn't access the image. Please try again.")
        end_conversation_and_show_menu(sender, None, responder)
    print(f"Image task for {sender} finished.")

def handle_document_upload(sender, media_id, filename, responder):
    print(f"Starting document upload task for {sender}, file: {filename}")
    responder.text(sender, f"Received {filename}. Processing and adding to knowledge base...")
    if media_url := get_media_url(media_id):
        if file_bytes := download_media_file(media_url):
            success, message = add_documents_to_vector_index(file_bytes, filename)
            if success:
                if db:
                    save_uploaded_file_metadata(sender, filename, len(file_bytes),
                                                 message.split(' ')[-2] if 'chunks' in message else 0)
                responder.text(sender, f"✅ Success! {message}\n\nI can now answer questions based on this content.")
            else:
                responder.text(sender, f"❌ Failed to process file: {message}")
        else:
            responder.text(sender, "❌ Sorry, I couldn't download the file.")
    else:
        responder.text(sender, "❌ Sorry, I couldn't access the file.")
    if sender in user_states:
        del user_states[sender]
        clear_user_state_from_firebase(sender)
    end_conversation_and_show_menu(sender, None, responder)
    print(f"Document upload task for {sender} finished.")

def handle_bulk_document_upload(sender, media_items, responder):
    print(f"Starting bulk upload for {sender}, {len(media_items)} files")
    successful_uploads = 0
    failed_uploads = 0
    processed_files = []
    responder.text(sender, f"📦 Processing {len(media_items)} files... This may take a few minutes.")
    for i, media_item in enumerate(media_items, 1):
        media_id = media_item.get('id')
        filename = media_item.get('filename', f'document_{i}')
        responder.text(sender, f"📄 Processing ({i}/{len(media_items)}): {filename}")
        if media_url := get_media_url(media_id):
            if file_bytes := download_media_file(media_url):
                success, message = add_documents_to_vector_index(file_bytes, filename)
                if success:
                    successful_uploads += 1
                    processed_files.append(filename)
                    if db:
                        save_uploaded_file_metadata(sender, filename, len(file_bytes),
                                                     message.split(' ')[-2] if 'chunks' in message else 0)
                else:
                    failed_uploads += 1
                    responder.text(sender, f"❌ Failed to process {filename}: {message}")
            else:
                failed_uploads += 1
                responder.text(sender, f"❌ Couldn't download {filename}")
        else:
            failed_uploads += 1
            responder.text(sender, f"❌ Couldn't access {filename}")
    summary_msg = f"✅ Bulk Upload Complete!\n\n"
    summary_msg += f"📊 Success: {successful_uploads} files\n"
    if failed_uploads > 0:
        summary_msg += f"❌ Failed: {failed_uploads} files\n"
    if successful_uploads > 0:
        summary_msg += f"\n📚 Added to knowledge base:\n"
        for file in processed_files[:5]:
            summary_msg += f"• {file}\n"
        if len(processed_files) > 5:
            summary_msg += f"• ... and {len(processed_files) - 5} more\n"
        summary_msg += f"\nYou can now ask questions based on this content!"
    responder.text(sender, summary_msg)
    if sender in user_states:
        del user_states[sender]
    end_conversation_and_show_menu(sender, None, responder)

# New byte-based handlers for web uploads
def handle_document_upload_bytes(sender, file_bytes, filename, responder):
    responder.text(sender, f"Received {filename}. Processing and adding to knowledge base...")
    success, message = add_documents_to_vector_index(file_bytes, filename)
    if success:
        if db:
            save_uploaded_file_metadata(sender, filename, len(file_bytes),
                                         message.split(' ')[-2] if 'chunks' in message else 0)
        responder.text(sender, f"✅ Success! {message}\n\nI can now answer questions based on this content.")
    else:
        responder.text(sender, f"❌ Failed to process file: {message}")
    end_conversation_and_show_menu(sender, None, responder)

def handle_image_task_bytes(sender, image_bytes, prompt, responder):
    responder.text(sender, "🖼️ Analyzing your image...")
    if len(image_bytes) > 20 * 1024 * 1024:
        responder.text(sender, "❌ Image is too large (over 20MB).")
        end_conversation_and_show_menu(sender, None, responder)
        return
    if 'podcast' in prompt.lower() or 'audio' in prompt.lower():
        ocr_prompt = "Extract all the text from this image. Do not summarize or explain, just provide the raw text."
        ocr_response = query_gemini_vision(ocr_prompt, image_bytes)
        if ocr_response.get("source") == "gemini_vision":
            extracted_text = ocr_response["result"]
            if len(extracted_text.strip()) > 50:
                audio_bytes = generate_voiceover(extracted_text)
                if audio_bytes:
                    responder.audio(sender, audio_bytes, "podcast.mp3")
                    end_conversation_and_show_menu(sender, "✅ Here is your podcast from the image text!", responder)
                else:
                    responder.text(sender, "❌ I extracted the text, but couldn't create the audio.")
                    end_conversation_and_show_menu(sender, None, responder)
            else:
                responder.text(sender, "❌ Not enough text found in the image to create a podcast.")
                end_conversation_and_show_menu(sender, None, responder)
        else:
            responder.text(sender, ocr_response["result"])
            end_conversation_and_show_menu(sender, None, responder)
    else:
        response = query_gemini_vision(prompt, image_bytes)
        responder.text(sender, response["result"])
        end_conversation_and_show_menu(sender, None, responder)

def handle_podcast_task_bytes(sender, image_bytes, responder):
    handle_image_task_bytes(sender, image_bytes, "podcast", responder)

def handle_summary_task_bytes(sender, image_bytes, responder):
    handle_image_task_bytes(sender, image_bytes, "summary", responder)

# ==============================================================================
# --- 14. FIREBASE FUNCTIONS ---
# ==============================================================================


db = None

try:
    if not firebase_admin._apps:
        # Try to use firebase_config.json first, fallback to service account
        if os.path.exists("firebase_config.json"):
            _fb_cred = credentials.Certificate("firebase_config.json")
        else:
            # Note: We're no longer using the service account for Google APIs,
            # but Firebase can still use it if needed
            _fb_cred = credentials.Certificate("sahayak-465916-8bf5ddce5515.json")
        firebase_admin.initialize_app(_fb_cred)
    db = firestore.client()
    print("Firebase initialized successfully.")
except Exception as _fb_err:
    print(f"Warning: Firebase not available – running without Firestore. ({_fb_err})")

def get_user_memory_from_firebase(sender):
    """Load conversation memory for a user from Firestore."""
    if not db:
        return []
    try:
        doc = db.collection("user_sessions").document(sender).get()
        if doc.exists:
            return doc.to_dict().get("memory", [])
    except Exception as e:
        print(f"Error loading memory from Firebase: {e}")
    return []

def save_user_memory_to_firebase(sender, memory):
    """Persist conversation memory for a user to Firestore."""
    if not db:
        return
    try:
        db.collection("user_sessions").document(sender).set({"memory": memory}, merge=True)
    except Exception as e:
        print(f"Error saving memory to Firebase: {e}")

def initialize_user_session(sender):
    """Load user memory from Firebase into in-memory store if not already loaded."""
    if sender not in user_memory or not user_memory[sender]:
        user_memory[sender] = get_user_memory_from_firebase(sender)
    if db:
        try:
            doc = db.collection("user_sessions").document(sender).get()
            if doc.exists:
                data = doc.to_dict()
                if data.get("state"):
                    user_states[sender] = data["state"]
                if data.get("temp_data"):
                    user_temp_data[sender] = data["temp_data"]
        except Exception as e:
            print(f"Error initializing session from Firebase: {e}")

def update_user_state(sender, state=None, temp_data=None):
    """Update in-memory and Firestore state for a user."""
    if state is not None:
        user_states[sender] = state
    if temp_data is not None:
        user_temp_data[sender].update(temp_data)
    if db:
        try:
            update = {}
            if state is not None:
                update["state"] = state
            if temp_data is not None:
                update["temp_data"] = dict(user_temp_data[sender])
            if update:
                db.collection("user_sessions").document(sender).set(update, merge=True)
        except Exception as e:
            print(f"Error updating user state in Firebase: {e}")

def append_to_memory(sender, role, content):
    """Append a message to the user's conversation memory."""
    user_memory[sender].append({"role": role, "content": content, "timestamp": datetime.now().isoformat()})
    if len(user_memory[sender]) > MAX_HISTORY:
        user_memory[sender] = user_memory[sender][-MAX_HISTORY:]
    save_user_memory_to_firebase(sender, user_memory[sender])

def clear_user_session(sender):
    """Clear in-memory state for a user and remove from Firestore."""
    user_states.pop(sender, None)
    user_temp_data.pop(sender, None)
    clear_user_state_from_firebase(sender)

def clear_user_state_from_firebase(sender):
    """Delete the user's session document from Firestore."""
    if not db:
        return
    try:
        db.collection("user_sessions").document(sender).delete()
    except Exception as e:
        print(f"Error clearing user state from Firebase: {e}")

def save_uploaded_file_metadata(sender, filename, file_size, chunks_count):
    """Save metadata about an uploaded file to Firestore."""
    if not db:
        return
    try:
        db.collection("uploaded_files").add({
            "sender": sender,
            "filename": filename,
            "file_size": file_size,
            "chunks_count": chunks_count,
            "uploaded_at": datetime.now().isoformat(),
        })
    except Exception as e:
        print(f"Error saving file metadata: {e}")

def get_user_uploaded_files(sender):
    """Retrieve metadata for all files uploaded by a specific user."""
    if not db:
        return []
    try:
        docs = db.collection("uploaded_files").where("sender", "==", sender).stream()
        return [doc.to_dict() for doc in docs]
    except Exception as e:
        print(f"Error getting uploaded files: {e}")
        return []

def cleanup_orphaned_file_metadata():
    """Remove Firestore file metadata records whose vector index entries no longer exist."""
    if not db:
        return
    try:
        active_files = set(os.path.basename(f) for f in get_uploaded_files_list())
        docs = db.collection("uploaded_files").stream()
        for doc in docs:
            data = doc.to_dict()
            if data.get("filename") not in active_files:
                doc.reference.delete()
        print("Orphaned file metadata cleanup complete.")
    except Exception as e:
        print(f"Error during cleanup: {e}")

# ── Web chat persistence ───────────────────────────────────────────────────

def _serialize_chat_value(value):
    """Convert Firestore values into JSON-safe primitives."""
    if isinstance(value, dict):
        return {key: _serialize_chat_value(item) for key, item in value.items()}
    if isinstance(value, list):
        return [_serialize_chat_value(item) for item in value]
    if isinstance(value, datetime):
        return value.isoformat()

    isoformat = getattr(value, "isoformat", None)
    if callable(isoformat):
        try:
            return isoformat()
        except Exception:
            pass

    return value

def _build_chat_title(chat_data):
    """Return a stable chat title, even for older documents missing title."""
    title = (chat_data or {}).get("title")
    if isinstance(title, str) and title.strip():
        return title.strip()

    for message in chat_data.get("messages", []):
        if message.get("role") == "user":
            content = (message.get("content") or "").strip()
            if content:
                return content[:50]

    return "New conversation"

def _normalize_chat_message(message):
    """Normalize legacy and Firestore-backed messages for the web client."""
    if not isinstance(message, dict):
        return None

    normalized = _serialize_chat_value(message)

    if normalized.get("role") == "assistant" and "actions" not in normalized:
        content = normalized.get("content")
        if content:
            normalized["actions"] = [{"type": "text", "content": str(content)}]

    return normalized

def get_web_chats_list(uid):
    """Return chat summaries [{chat_id, title, updated_at}] for a user, newest first."""
    if not db:
        return []
    try:
        docs = db.collection("web_chats").document(uid).collection("sessions").stream()
        result = []
        for d in docs:
            data = _serialize_chat_value(d.to_dict() or {})
            result.append({
                "chat_id": d.id,
                "title": _build_chat_title(data),
                "updated_at": data.get("updated_at", "") or "",
            })
        result.sort(key=lambda x: x["updated_at"], reverse=True)
        return result[:100]
    except Exception as e:
        print(f"Error fetching web chat list: {e}")
        return []

def get_web_chat_detail(uid, chat_id):
    """Return full chat document {title, messages, ...} or None."""
    if not db:
        return None
    try:
        doc = db.collection("web_chats").document(uid) \
                .collection("sessions").document(chat_id).get()
        if doc.exists:
            data = _serialize_chat_value(doc.to_dict() or {})
            data["title"] = _build_chat_title(data)
            messages = []
            for message in data.get("messages", []):
                normalized = _normalize_chat_message(message)
                if normalized:
                    messages.append(normalized)
            data["messages"] = messages
            return data
    except Exception as e:
        print(f"Error fetching web chat: {e}")
    return None

def _strip_binary_from_actions(actions):
    """Remove large base64 payloads from actions before saving to Firestore."""
    safe = []
    for action in (actions or []):
        if action.get("type") in ("document", "audio", "video") and "data" in action:
            stripped = {k: v for k, v in action.items() if k != "data"}
            stripped["data_stripped"] = True
            safe.append(stripped)
        else:
            safe.append(action)
    return safe

def save_web_chat_exchange(uid, chat_id, user_content, bot_actions):
    """Append a user→bot exchange to the chat session in Firestore."""
    if not db:
        return
    try:
        now = datetime.now().isoformat()
        ref = db.collection("web_chats").document(uid) \
                .collection("sessions").document(chat_id)
        doc = ref.get()
        if doc.exists:
            data = doc.to_dict()
            messages = data.get("messages", [])
            title = data.get("title", "New conversation")
            created_at = data.get("created_at", now)
        else:
            messages = []
            title = (user_content or "New conversation")[:50]
            created_at = now
        if user_content:
            messages.append({"role": "user", "content": user_content, "timestamp": now})
        safe_actions = _strip_binary_from_actions(bot_actions)
        messages.append({"role": "assistant", "actions": safe_actions, "timestamp": now})
        if len(messages) > MAX_HISTORY * 2:
            messages = messages[-(MAX_HISTORY * 2):]
        ref.set({
            "title": title,
            "created_at": created_at,
            "updated_at": now,
            "messages": messages,
        })
    except Exception as e:
        print(f"Error saving web chat: {e}")

def delete_web_chat_session(uid, chat_id):
    """Delete a single chat session from Firestore."""
    if not db:
        return
    try:
        db.collection("web_chats").document(uid) \
          .collection("sessions").document(chat_id).delete()
    except Exception as e:
        print(f"Error deleting web chat: {e}")

def clear_all_web_chats(uid):
    """Delete all chat sessions for a user."""
    if not db:
        return
    try:
        sessions = db.collection("web_chats").document(uid) \
                     .collection("sessions").stream()
        for doc in sessions:
            doc.reference.delete()
    except Exception as e:
        print(f"Error clearing web chats: {e}")

# ==============================================================================
# --- 14. TEACHER-SPECIFIC CLASSES ---
# ==============================================================================

class CourseMaterialManager:
    """Efficient storage and retrieval of course materials"""
    
    def __init__(self):
        self.vector_index_path = DYNAMIC_VECTOR_INDEX_PATH
        self.metadata_collection = "course_materials" if db else None
    
    def store_unit_materials(self, course_id, unit_number, syllabus_text, ppt_bytes, ppt_filename):
        """
        Store syllabus and PPT for a unit with proper metadata
        """
        # Create unit identifier
        unit_id = f"course_{course_id}_unit_{unit_number}"
        
        # Store syllabus in vector DB with metadata
        syllabus_success, syllabus_msg = add_documents_to_vector_index(
            syllabus_text.encode('utf-8'), 
            f"{unit_id}_syllabus.txt"
        )
        
        # Store PPT in vector DB
        ppt_success, ppt_msg = add_documents_to_vector_index(
            ppt_bytes,
            f"{unit_id}_{ppt_filename}"
        )
        
        # Store metadata in Firestore
        if db and syllabus_success and ppt_success:
            doc_ref = db.collection("course_units").document(unit_id)
            doc_ref.set({
                "course_id": course_id,
                "unit_number": unit_number,
                "syllabus_file": f"{unit_id}_syllabus.txt",
                "ppt_file": f"{unit_id}_{ppt_filename}",
                "stored_at": datetime.now().isoformat(),
                "question_pattern": {
                    "two_marks_per_unit": 2,
                    "sixteen_marks_per_unit": 1
                }
            })
        
        return {
            "unit_id": unit_id,
            "syllabus_status": syllabus_msg,
            "ppt_status": ppt_msg
        }
    
    def get_unit_materials(self, course_id, unit_number):
        """Retrieve all materials for a specific unit"""
        unit_id = f"course_{course_id}_unit_{unit_number}"
        
        if not db:
            return None
            
        doc = db.collection("course_units").document(unit_id).get()
        if doc.exists:
            return doc.to_dict()
        return None
    
    def get_all_units(self, course_id):
        """Get all units for a course"""
        if not db:
            return []
            
        docs = db.collection("course_units")\
                 .where("course_id", "==", course_id)\
                 .order_by("unit_number")\
                 .stream()
        
        return [doc.to_dict() for doc in docs]
    
    def check_materials_exist(self, course_id, unit_number):
        """Check if materials exist for a specific unit"""
        unit_id = f"course_{course_id}_unit_{unit_number}"
        
        if not db:
            return False
            
        doc = db.collection("course_units").document(unit_id).get()
        return doc.exists

class QuestionPaperGenerator:
    """Generate question papers following specific pattern: 10x2 marks + 5x16 marks per unit"""
    
    def __init__(self, course_id):
        self.course_id = course_id
        self.manager = CourseMaterialManager()
        self.model = genai.GenerativeModel('gemini-3-flash-preview')
    
    def extract_unit_content(self, unit_number):
        """Extract all content from syllabus and PPT for a unit"""
        unit_data = self.manager.get_unit_materials(self.course_id, unit_number)
        if not unit_data:
            return None
        
        # Query vector DB for all content related to this unit
        unit_id = f"course_{self.course_id}_unit_{unit_number}"
        
        # Get syllabus content
        syllabus_query = f"unit {unit_number} syllabus"
        syllabus_docs = query_dynamic_rag_strict(syllabus_query, similarity_threshold=0.3)
        
        # Get PPT content
        ppt_query = f"unit {unit_number} presentation ppt"
        ppt_docs = query_dynamic_rag_strict(ppt_query, similarity_threshold=0.3)
        
        combined_content = ""
        if syllabus_docs.get("result"):
            combined_content += syllabus_docs["result"] + "\n\n"
        if ppt_docs.get("result"):
            combined_content += ppt_docs["result"]
        
        return combined_content if combined_content else None
    
    def generate_unit_questions(self, unit_number, content):
        """Generate exactly 10 two-mark questions and 5 sixteen-mark questions for a unit"""
        
        # Prompt for 2-mark questions (factual, definition-based)
        two_marks_prompt = f"""
        Based on the following content from Unit {unit_number}, generate EXACTLY 10 two-mark questions.
        
        Content:
        {content[:3000]}  # Limit context for focused generation
        
        Rules for 2-mark questions:
        - Questions should test basic understanding and recall
        - Each question should have a concise answer (2-3 sentences max)
        - Cover different topics from the unit
        - Format as: Q1. [question]
        
        Generate exactly 10 questions:
        """
        
        two_marks_response = self.model.generate_content(two_marks_prompt)
        two_marks_questions = self.parse_questions(two_marks_response.text, num_expected=10)
        
        # Prompt for 16-mark questions (analytical, comprehensive)
        sixteen_marks_prompt = f"""
        Based on the following content from Unit {unit_number}, generate EXACTLY 5 sixteen-mark questions.
        
        Content:
        {content}
        
        Rules for 16-mark questions:
        - Questions should require detailed explanation, analysis, or application
        - Should cover major topics and concepts
        - May have sub-parts (a, b, c)
        - Should test deeper understanding
        - Format as: Q1. [question] (with sub-parts if applicable)
        
        Generate exactly 5 questions:
        """
        
        sixteen_marks_response = self.model.generate_content(sixteen_marks_prompt)
        sixteen_marks_questions = self.parse_questions(sixteen_marks_response.text, num_expected=5)
        
        return {
            "unit": unit_number,
            "two_marks": two_marks_questions,
            "sixteen_marks": sixteen_marks_questions
        }
    
    def parse_questions(self, text, num_expected):
        """Parse generated questions into structured format"""
        questions = []
        lines = text.strip().split('\n')
        
        for line in lines:
            # Look for question patterns (Q1., 1., etc.)
            if re.match(r'^Q?\d+[\.\)]', line.strip()) or re.match(r'^\d+[\.\)]', line.strip()):
                questions.append(line.strip())
        
        # If we didn't get enough, try a different parsing approach
        if len(questions) < num_expected:
            # Fallback: take first N non-empty lines
            non_empty = [l.strip() for l in lines if l.strip()]
            questions = non_empty[:num_expected]
        
        return questions[:num_expected]  # Ensure we don't return more than expected
    
    def generate_full_question_paper(self, units=[1,2,3,4,5]):
        """Generate complete question paper for all units"""
        
        paper_structure = {
            "course_id": self.course_id,
            "generated_at": datetime.now().isoformat(),
            "total_marks": 0,
            "sections": []
        }
        
        all_questions = []
        total_marks = 0
        
        for unit in units:
            print(f"Generating questions for Unit {unit}...")
            content = self.extract_unit_content(unit)
            
            if not content:
                print(f"Warning: No content found for Unit {unit}")
                continue
            
            unit_questions = self.generate_unit_questions(unit, content)
            
            # Calculate marks
            unit_marks = (len(unit_questions["two_marks"]) * 2) + (len(unit_questions["sixteen_marks"]) * 16)
            total_marks += unit_marks
            
            all_questions.append({
                "unit": unit,
                "two_marks": unit_questions["two_marks"],
                "sixteen_marks": unit_questions["sixteen_marks"],
                "section_marks": unit_marks
            })
        
        paper_structure["sections"] = all_questions
        paper_structure["total_marks"] = total_marks
        
        return paper_structure

class QuestionPaperFormatter:
    """Format question paper in different output formats"""
    
    @staticmethod
    def format_as_text(paper_data):
        """Format as plain text"""
        lines = []
        lines.append("=" * 80)
        lines.append(f"QUESTION PAPER - {paper_data['course_id']}")
        lines.append(f"Generated: {paper_data['generated_at']}")
        lines.append(f"Total Marks: {paper_data['total_marks']}")
        lines.append("=" * 80)
        lines.append("")
        
        for section in paper_data['sections']:
            lines.append(f"\nUNIT {section['unit']}")
            lines.append("-" * 40)
            
            lines.append(f"\nPART A - (10 x 2 = 20 marks)")
            for i, q in enumerate(section['two_marks'], 1):
                lines.append(f"  {i}. {q}")
            
            lines.append(f"\nPART B - (5 x 16 = 80 marks)")
            for i, q in enumerate(section['sixteen_marks'], 1):
                lines.append(f"  {i}. {q}")
            
            lines.append("")
        
        return "\n".join(lines)
    
    @staticmethod
    def format_as_pdf(paper_data, filename="question_paper.pdf"):
        """Generate PDF with proper formatting"""
        from fpdf import FPDF
        from fpdf.enums import XPos, YPos
        
        pdf = FPDF()
        pdf.set_margins(left=10, top=10, right=10)
        pdf.add_page()
        
        # Title
        pdf.set_font("Arial", 'B', 16)
        pdf.cell(0, 10, txt="QUESTION PAPER", new_x=XPos.LMARGIN, new_y=YPos.NEXT, align='C')
        pdf.set_font("Arial", '', 12)
        pdf.cell(0, 10, txt=f"Course: {paper_data['course_id']}", new_x=XPos.LMARGIN, new_y=YPos.NEXT, align='C')
        pdf.cell(0, 10, txt=f"Total Marks: {paper_data['total_marks']}", new_x=XPos.LMARGIN, new_y=YPos.NEXT, align='C')
        pdf.ln(10)
        
        for section in paper_data['sections']:
            # Unit header
            pdf.set_font("Arial", 'B', 14)
            pdf.cell(0, 10, txt=f"UNIT {section['unit']}", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
            
            # Part A
            pdf.set_font("Arial", 'B', 12)
            pdf.cell(0, 10, txt="PART A - (10 x 2 = 20 marks)", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
            pdf.set_font("Arial", '', 11)
            for i, q in enumerate(section['two_marks'], 1):
                # Handle long questions
                safe_q = ''.join(char if ord(char) < 128 else '?' for char in q)
                pdf.multi_cell(0, 5, txt=f"{i}. {safe_q}", align='L')
            
            # Part B
            pdf.ln(5)
            pdf.set_font("Arial", 'B', 12)
            pdf.cell(0, 10, txt="PART B - (5 x 16 = 80 marks)", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
            pdf.set_font("Arial", '', 11)
            for i, q in enumerate(section['sixteen_marks'], 1):
                safe_q = ''.join(char if ord(char) < 128 else '?' for char in q)
                pdf.multi_cell(0, 5, txt=f"{i}. {safe_q}", align='L')
            
            pdf.ln(10)
        
        # Use output() for compatibility with all versions
        pdf_bytes = pdf.output()
        if isinstance(pdf_bytes, str):
            pdf_bytes = pdf_bytes.encode('latin-1')
        return pdf_bytes
    
    @staticmethod
    def format_as_docx(paper_data):
        """Generate Word document"""
        from docx import Document
        from io import BytesIO
        
        doc = Document()
        
        # Title
        title = doc.add_heading('QUESTION PAPER', 0)
        title.alignment = 1  # Center
        
        doc.add_paragraph(f"Course: {paper_data['course_id']}")
        doc.add_paragraph(f"Total Marks: {paper_data['total_marks']}")
        doc.add_paragraph(f"Generated: {paper_data['generated_at']}")
        
        for section in paper_data['sections']:
            # Unit header
            doc.add_heading(f'UNIT {section["unit"]}', level=1)
            
            # Part A
            doc.add_heading('PART A - (10 x 2 = 20 marks)', level=2)
            for i, q in enumerate(section['two_marks'], 1):
                doc.add_paragraph(f"{i}. {q}", style='List Number')
            
            # Part B
            doc.add_heading('PART B - (5 x 16 = 80 marks)', level=2)
            for i, q in enumerate(section['sixteen_marks'], 1):
                doc.add_paragraph(f"{i}. {q}", style='List Number')
        
        # Save to bytes
        bytes_io = BytesIO()
        doc.save(bytes_io)
        bytes_io.seek(0)
        return bytes_io.getvalue()

class StudentProgressTracker:
    """Track student progress and generate analytics"""
    
    def __init__(self):
        self.db = db
    
    def get_student_queries(self, student_id, days=30):
        """Get all queries from a student in the last X days"""
        if not self.db:
            return []
        
        cutoff = datetime.now() - timedelta(days=days)
        
        try:
            docs = self.db.collection("user_sessions")\
                         .document(student_id)\
                         .collection("queries")\
                         .where("timestamp", ">=", cutoff.isoformat())\
                         .stream()
            
            return [doc.to_dict() for doc in docs]
        except Exception as e:
            print(f"Error fetching student queries: {e}")
            return []
    
    def analyze_student_performance(self, student_id):
        """Analyze student's strengths and weaknesses"""
        queries = self.get_student_queries(student_id)
        
        if not queries:
            return None
        
        # Extract topics and analyze patterns
        topics = defaultdict(lambda: {"count": 0, "success_rate": 0})
        
        for query in queries:
            topic = self.extract_topic(query.get("content", ""))
            topics[topic]["count"] += 1
            # Check if query was successfully answered
            if query.get("answered", False):
                topics[topic]["success_rate"] += 1
        
        # Calculate success rates
        for topic in topics:
            if topics[topic]["count"] > 0:
                topics[topic]["success_rate"] = (topics[topic]["success_rate"] / topics[topic]["count"]) * 100
        
        # Identify strengths and weaknesses
        strengths = [t for t, data in topics.items() if data["success_rate"] >= 70 and data["count"] >= 3]
        weaknesses = [t for t, data in topics.items() if data["success_rate"] < 50 and data["count"] >= 3]
        
        return {
            "student_id": student_id,
            "total_queries": len(queries),
            "topics_covered": len(topics),
            "strengths": strengths[:5],
            "weaknesses": weaknesses[:5],
            "topic_details": dict(topics)
        }
    
    def extract_topic(self, query):
        """Extract main topic from query using simple NLP"""
        # Simple implementation - can be enhanced with ML
        common_topics = ["mathematics", "science", "history", "literature", 
                         "programming", "physics", "chemistry", "biology"]
        
        query_lower = query.lower()
        for topic in common_topics:
            if topic in query_lower:
                return topic
        
        # If no common topic found, return first few words
        words = query.split()[:3]
        return " ".join(words) if words else "general"
    
    def get_class_analytics(self, class_id):
        """Get analytics for entire class"""
        if not self.db:
            return None
        
        try:
            # Get all students in class
            students = self.db.collection("classes")\
                             .document(class_id)\
                             .collection("students")\
                             .stream()
            
            class_data = {
                "class_id": class_id,
                "total_students": 0,
                "active_students": 0,
                "student_performance": [],
                "common_topics": defaultdict(int)
            }
            
            for student in students:
                student_data = student.to_dict()
                class_data["total_students"] += 1
                
                # Check if student is active
                if student_data.get("last_active"):
                    last_active = datetime.fromisoformat(student_data["last_active"])
                    if datetime.now() - last_active < timedelta(days=7):
                        class_data["active_students"] += 1
                
                # Get performance
                performance = self.analyze_student_performance(student.id)
                if performance:
                    class_data["student_performance"].append(performance)
                    
                    # Track common topics
                    for topic in performance.get("weaknesses", []):
                        class_data["common_topics"][topic] += 1
            
            return class_data
            
        except Exception as e:
            print(f"Error getting class analytics: {e}")
            return None

# ==============================================================================
# --- 16. DIAGNOSTIC ENDPOINT ---
# ==============================================================================


@app.route("/api/diagnose-drive", methods=["GET"])
@login_required
def diagnose_drive():
    """Diagnose Google Drive OAuth setup and permissions"""
    if not drive_service:
        return jsonify({"error": "Drive service not initialized. Check OAuth setup."}), 500
    
    results = {
        "auth_type": "OAuth 2.0",
        "folder_id": GOOGLE_FORMS_TARGET_FOLDER_ID,
        "tests": []
    }
    
    try:
        # Get user info
        about = drive_service.about().get(fields="user").execute()
        results["user"] = about.get("user", {}).get("emailAddress", "Unknown")
        
        # Test 1: Check if we can list files
        try:
            files = drive_service.files().list(pageSize=5).execute()
            results["tests"].append({
                "test": "list_files",
                "success": True,
                "files_found": len(files.get("files", []))
            })
        except Exception as e:
            results["tests"].append({
                "test": "list_files",
                "success": False,
                "error": str(e)
            })
        
        # Test 2: Check folder access
        if GOOGLE_FORMS_TARGET_FOLDER_ID:
            try:
                folder = drive_service.files().get(
                    fileId=GOOGLE_FORMS_TARGET_FOLDER_ID,
                    fields="id,name,parents"
                ).execute()
                results["tests"].append({
                    "test": "access_target_folder",
                    "success": True,
                    "folder_name": folder.get("name")
                })
            except Exception as e:
                results["tests"].append({
                    "test": "access_target_folder",
                    "success": False,
                    "error": str(e)
                })
        
        # Test 3: Try to create a test file
        try:
            test_content = b"Test file to verify upload permissions"
            test_filename = "permission_test.txt"
            
            file_id, msg, err = upload_file_to_drive_with_folder_check(
                test_content, 
                test_filename
            )
            
            if file_id:
                # Clean up test file
                try:
                    drive_service.files().delete(fileId=file_id).execute()
                except:
                    pass
                
                results["tests"].append({
                    "test": "create_file",
                    "success": True
                })
            else:
                results["tests"].append({
                    "test": "create_file",
                    "success": False,
                    "error": err
                })
        except Exception as e:
            results["tests"].append({
                "test": "create_file",
                "success": False,
                "error": str(e)
            })
        
        return jsonify(results)
        
    except Exception as e:
        return jsonify({"error": str(e), "traceback": traceback.format_exc()}), 500

# ==============================================================================
# --- 17. WEBHOOK ROUTES ---
# ==============================================================================

@app.route("/webhook", methods=["GET", "POST"])
def webhook():
    if request.method == "GET":
        if request.args.get("hub.verify_token") == VERIFY_TOKEN:
            return request.args.get("hub.challenge")
        return "Verification token mismatch", 403

    data = request.get_json()
    try:
        if data and "entry" in data:
            for entry in data.get("entry", []):
                for change in entry.get("changes", []):
                    value = change.get("value", {})
                    if "messages" in value:
                        for msg in value.get("messages", []):
                            sender = msg["from"]
                            text = ""
                            msg_type = msg.get("type")
                            if msg_type == "document":
                                current_state = user_states.get(sender)
                                if current_state in ["awaiting_bulk_upload", "awaiting_material_file"]:
                                    media_id = msg['document']['id']
                                    filename = msg['document'].get('filename', 'uploaded_file')
                                    if 'bulk_upload_items' not in user_temp_data[sender]:
                                        user_temp_data[sender]['bulk_upload_items'] = []
                                    user_temp_data[sender]['bulk_upload_items'].append({
                                        'id': media_id,
                                        'filename': filename
                                    })
                                    threading.Timer(8.0, process_bulk_upload, args=[sender]).start()
                                    send_whatsapp_message(sender, f"✅ Added {filename} to upload queue...")
                                else:
                                    media_id = msg['document']['id']
                                    filename = msg['document'].get('filename', 'uploaded_file')
                                    thread = threading.Thread(target=handle_document_upload,
                                                              args=(sender, media_id, filename,
                                                                    WhatsAppResponder(ACCESS_TOKEN, PHONE_NUMBER_ID)))
                                    thread.start()
                                return "ok", 200
                            if msg_type == "button":
                                text = msg["button"]["text"]
                            elif msg_type == "text":
                                text = msg["text"]["body"]
                            elif msg_type == "interactive" and msg.get("interactive", {}).get("type") == "button_reply":
                                text = msg["interactive"]["button_reply"]["title"]
                            elif msg_type == "interactive" and msg.get("interactive", {}).get("type") == "list_reply":
                                text = msg["interactive"]["list_reply"]["title"]
                            elif msg_type == "audio":
                                media_id = msg['audio']['id']
                                thread = threading.Thread(target=handle_audio_task,
                                                          args=(sender, media_id,
                                                                WhatsAppResponder(ACCESS_TOKEN, PHONE_NUMBER_ID)))
                                thread.start()
                            elif msg_type == "image":
                                current_state = user_states.get(sender)
                                media_id = msg['image']['id']
                                if current_state == "awaiting_podcast_image":
                                    thread = threading.Thread(target=handle_podcast_task,
                                                              args=(sender, media_id,
                                                                    WhatsAppResponder(ACCESS_TOKEN, PHONE_NUMBER_ID)))
                                    thread.start()
                                elif current_state == "awaiting_summary_image":
                                    thread = threading.Thread(target=handle_summary_task,
                                                              args=(sender, media_id,
                                                                    WhatsAppResponder(ACCESS_TOKEN, PHONE_NUMBER_ID)))
                                    thread.start()
                                else:
                                    prompt = msg.get("image", {}).get("caption", "Explain this image.")
                                    thread = threading.Thread(target=handle_image_task,
                                                              args=(sender, media_id, prompt,
                                                                    WhatsAppResponder(ACCESS_TOKEN, PHONE_NUMBER_ID)))
                                    thread.start()
                                return "ok", 200
                            if text:
                                process_message(sender, text, WhatsAppResponder(ACCESS_TOKEN, PHONE_NUMBER_ID))
    except Exception as e:
        print(f"Webhook processing error: {e}")
        traceback.print_exc()
    return "ok", 200

def process_bulk_upload(sender):
    if sender in user_temp_data and 'bulk_upload_items' in user_temp_data[sender]:
        items = user_temp_data[sender]['bulk_upload_items']
        if items:
            print(f"🔄 Processing {len(items)} files for {sender}")
            user_temp_data[sender]['bulk_upload_items'] = []
            responder = WhatsAppResponder(ACCESS_TOKEN, PHONE_NUMBER_ID)
            handle_bulk_document_upload(sender, items, responder)

def handle_podcast_task(sender, media_id, responder):
    # This is a wrapper for handle_image_task with podcast prompt
    handle_image_task(sender, media_id, "podcast", responder)

def handle_summary_task(sender, media_id, responder):
    handle_image_task(sender, media_id, "summary", responder)

@app.route("/status", methods=["GET"])
def status():
    return {"status": "running"}, 200

@app.route("/api/check_similarity", methods=["POST"])
@login_required
def api_check_similarity():
    """Endpoint to check similarity scores without generating answers"""
    data = request.get_json()
    query = data.get("query", "")
    top_k = data.get("top_k", 5)
    
    if not query:
        return jsonify({"error": "No query provided"}), 400
    
    results = check_similarity_only(query, top_k)
    return jsonify(results)

# ==============================================================================
# --- 17. WEB API ROUTES ---
# ==============================================================================
def get_or_create_session(session_id):
    if not session_id:
        session_id = str(uuid.uuid4())
    return session_id

@app.route("/api/chat", methods=["POST"])
@login_required
def api_chat():
    data = request.get_json()
    message = data.get("message", "")
    chat_id = data.get("chat_id") or str(uuid.uuid4())
    uid = session['user_id']
    sender_key = f"web:{uid}:{chat_id}"

    if not message:
        return jsonify({"chat_id": chat_id, "actions": []})

    responder = WebResponder()
    process_message(sender_key, message, responder)
    save_web_chat_exchange(uid, chat_id, message, responder.actions)

    return jsonify({"chat_id": chat_id, "actions": responder.actions})

@app.route("/api/test-drive-upload", methods=["GET"])
@login_required
def test_drive_upload():
    """Test uploading a file to Shared Drive"""
    if not drive_service:
        return jsonify({"error": "Drive not initialized"})
    
    try:
        # Create a simple test file
        test_content = b"This is a test file to verify Shared Drive uploads work."
        test_filename = "test_upload.txt"
        
        file_id, success, error = upload_file_to_drive_with_folder_check(
            test_content,
            test_filename,
            GOOGLE_FORMS_TARGET_FOLDER_ID
        )
        
        if file_id:
            return jsonify({
                "success": True,
                "file_id": file_id,
                "file_url": f"https://drive.google.com/file/d/{file_id}/view",
                "message": "✅ Test file uploaded successfully!"
            })
        else:
            return jsonify({
                "success": False,
                "error": error
            })
            
    except Exception as e:
        return jsonify({
            "success": False,
            "error": str(e)
        })

@app.route("/api/upload", methods=["POST"])
@login_required
def api_upload():
    chat_id = request.form.get("chat_id") or str(uuid.uuid4())
    uid = session['user_id']
    sender_key = f"web:{uid}:{chat_id}"
    file = request.files.get("file")
    if not file:
        return jsonify({"error": "No file"}), 400

    filename = file.filename
    file_bytes = file.read()
    mime = file.mimetype or ""
    
    responder = WebResponder()
    
    # Determine if file is an image by checking MIME type and file extension
    image_extensions = {'.png', '.jpg', '.jpeg', '.gif', '.webp', '.bmp', '.svg', '.tiff', '.ico'}
    file_ext = os.path.splitext(filename)[1].lower()
    is_image = mime.lower().startswith("image/") or file_ext in image_extensions
    
    if is_image:
        responder.text(sender_key, f"✅ Image received: {filename}")
        current_state = user_states.get(sender_key)
        if current_state == "awaiting_podcast_image":
            handle_podcast_task_bytes(sender_key, file_bytes, responder)
        elif current_state == "awaiting_summary_image":
            handle_summary_task_bytes(sender_key, file_bytes, responder)
        else:
            caption = request.form.get("caption", "Explain this image.")
            handle_image_task_bytes(sender_key, file_bytes, caption, responder)
    else:
        handle_document_upload_bytes(sender_key, file_bytes, filename, responder)

    user_content = request.form.get("message") or f"\U0001f4ce {filename}"
    save_web_chat_exchange(uid, chat_id, user_content, responder.actions)

    return jsonify({"chat_id": chat_id, "actions": responder.actions})

@app.route("/api/transcribe", methods=["POST"])
@login_required
def api_transcribe():
    """Transcribe voice message to text using Whisper"""
    file = request.files.get("file")
    if not file:
        return jsonify({"error": "No audio file provided"}), 400
    
    try:
        audio_bytes = file.read()
        transcribed_text = transcribe_audio(audio_bytes)
        if transcribed_text:
            return jsonify({"text": transcribed_text, "success": True})
        else:
            return jsonify({"error": "Failed to transcribe audio"}), 400
    except Exception as e:
        print(f"Transcription error: {e}")
        return jsonify({"error": str(e)}), 500

@app.route("/api/chats/<chat_id>/rename", methods=["POST"])
@login_required
def api_rename_chat(chat_id):
    """Rename a chat session"""
    data = request.get_json()
    new_title = (data.get("title", "") or "").strip()
    if not new_title:
        return jsonify({"error": "Title cannot be empty"}), 400
    
    uid = session['user_id']
    if not db:
        return jsonify({"error": "Database not available"}), 500
    
    try:
        db.collection("web_chats").document(uid) \
          .collection("sessions").document(chat_id) \
          .update({"title": new_title})
        return jsonify({"success": True})
    except Exception as e:
        print(f"Error renaming chat: {e}")
        return jsonify({"error": str(e)}), 500

@app.route("/api/generate-worksheet-pdf", methods=["POST"])
@login_required
def api_generate_worksheet_pdf():
    """Generate worksheet PDFs (questions and answers) from uploaded materials"""
    data = request.get_json()
    topic = (data.get("topic", "") or "").strip()
    
    if not topic:
        return jsonify({"error": "Topic is required"}), 400
    
    try:
        # Query uploaded materials for the topic
        relevant_content = query_vector_db_for_topic(topic, similarity_threshold=0.4)
        
        if not relevant_content:
            return jsonify({
                "error": f"No relevant information found about '{topic}' in uploaded materials. Please upload relevant documents first."
            }), 400
        
        # Generate worksheet with separate questions and answers
        worksheet_full = generate_worksheet_parts(topic, relevant_content)
        
        if not worksheet_full:
            return jsonify({"error": "Failed to generate worksheet content"}), 500
        
        # Create separate PDFs for questions and answers
        questions_bytes, answers_bytes = create_worksheet_pdfs(topic, worksheet_full)
        
        if not questions_bytes or not answers_bytes:
            return jsonify({"error": "Failed to create PDF files"}), 500
        
        # Return both PDFs as base64-encoded data
        return jsonify({
            "success": True,
            "topic": topic,
            "questions_pdf": base64.b64encode(questions_bytes).decode('utf-8'),
            "questions_filename": f"{topic.replace(' ', '_')}_worksheet_questions.pdf",
            "answers_pdf": base64.b64encode(answers_bytes).decode('utf-8'),
            "answers_filename": f"{topic.replace(' ', '_')}_worksheet_answers.pdf"
        })
    
    except Exception as e:
        print(f"Error generating worksheet PDF: {e}")
        traceback.print_exc()
        return jsonify({"error": f"Error generating worksheet: {str(e)}"}), 500

# ==============================================================================
# --- 18. TEACHER API ROUTES ---
# ==============================================================================

@app.route("/api/teacher/upload_course_materials", methods=["POST"])
@teacher_required
def upload_course_materials():
    """Upload syllabus and PPT for all 5 units at once"""
    course_id = request.form.get("course_id", "default_course")
    files = request.files.getlist("files")
    
    if not files:
        return jsonify({"error": "No files provided"}), 400
    
    # Expected: syllabus_unit1.txt, ppt_unit1.pptx, syllabus_unit2.txt, etc.
    materials = defaultdict(dict)
    
    for file in files:
        filename = file.filename
        # Parse filename to identify unit and type
        unit_match = re.search(r'unit[_\s]*(\d+)', filename.lower())
        if not unit_match:
            continue
            
        unit = int(unit_match.group(1))
        if unit < 1 or unit > 5:
            continue
            
        if "syllabus" in filename.lower():
            materials[unit]["syllabus"] = file.read().decode('utf-8')
        elif "ppt" in filename.lower() or "presentation" in filename.lower():
            materials[unit]["ppt_bytes"] = file.read()
            materials[unit]["ppt_filename"] = filename
    
    # Store each unit
    manager = CourseMaterialManager()
    results = []
    
    for unit in range(1, 6):  # Units 1-5
        if unit in materials and "syllabus" in materials[unit] and "ppt_bytes" in materials[unit]:
            result = manager.store_unit_materials(
                course_id, 
                unit,
                materials[unit]["syllabus"],
                materials[unit]["ppt_bytes"],
                materials[unit]["ppt_filename"]
            )
            results.append(result)
    
    if not results:
        return jsonify({"error": "No valid files found. Please name files like 'syllabus_unit1.txt' and 'ppt_unit1.pptx'"}), 400
    
    return jsonify({
        "success": True,
        "units_stored": len(results),
        "details": results
    })

@app.route("/api/teacher/generate_question_paper", methods=["POST"])
@teacher_required
def generate_question_paper():
    """Generate question paper based on stored syllabus and PPTs"""
    data = request.get_json()
    course_id = data.get("course_id", "default_course")
    units = data.get("units", [1, 2, 3, 4, 5])
    format_type = data.get("format", "pdf")  # pdf, docx, text
    
    # Check if materials exist
    manager = CourseMaterialManager()
    existing_units = []
    missing_units = []
    
    for unit in units:
        if manager.check_materials_exist(course_id, unit):
            existing_units.append(unit)
        else:
            missing_units.append(unit)
    
    if not existing_units:
        return jsonify({
            "error": "No course materials found. Please upload syllabus and PPTs first.",
            "missing_units": units
        }), 404
    
    if missing_units:
        return jsonify({
            "warning": f"Materials not found for units: {missing_units}. Generating paper for units: {existing_units}",
            "missing_units": missing_units,
            "existing_units": existing_units
        }), 200
    
    # Generate question paper
    generator = QuestionPaperGenerator(course_id)
    formatter = QuestionPaperFormatter()
    
    try:
        paper_data = generator.generate_full_question_paper(existing_units)
        
        if not paper_data["sections"]:
            return jsonify({"error": "Failed to generate questions. Please check your materials."}), 500
        
        # Format based on request
        if format_type == "pdf":
            pdf_bytes = formatter.format_as_pdf(paper_data)
            filename = f"question_paper_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf"
            
            # Return as downloadable file
            response = make_response(pdf_bytes)
            response.headers['Content-Type'] = 'application/pdf'
            response.headers['Content-Disposition'] = f'attachment; filename={filename}'
            return response
            
        elif format_type == "docx":
            docx_bytes = formatter.format_as_docx(paper_data)
            filename = f"question_paper_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
            
            response = make_response(docx_bytes)
            response.headers['Content-Type'] = 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
            response.headers['Content-Disposition'] = f'attachment; filename={filename}'
            return response
            
        else:  # text format
            text = formatter.format_as_text(paper_data)
            return jsonify({
                "success": True,
                "paper": text,
                "total_marks": paper_data["total_marks"],
                "units_covered": existing_units
            })
            
    except Exception as e:
        print(f"Error generating question paper: {e}")
        traceback.print_exc()
        return jsonify({"error": str(e)}), 500

@app.route("/api/teacher/check_materials", methods=["POST"])
@teacher_required
def check_materials():
    """Check which units have materials uploaded"""
    data = request.get_json()
    course_id = data.get("course_id", "default_course")
    units = data.get("units", [1, 2, 3, 4, 5])
    
    manager = CourseMaterialManager()
    status = {}
    
    for unit in units:
        status[unit] = manager.check_materials_exist(course_id, unit)
    
    return jsonify({
        "course_id": course_id,
        "units": status
    })

@app.route("/api/teacher/get_uploaded_files", methods=["GET"])
@teacher_required
def get_teacher_uploaded_files():
    """Get list of all uploaded files in vector database"""
    files = get_uploaded_files_list()
    return jsonify({
        "files": files,
        "count": len(files)
    })

@app.route("/api/teacher/delete_materials", methods=["POST"])
@teacher_required
def delete_materials():
    """Delete materials for specific units"""
    data = request.get_json()
    course_id = data.get("course_id", "default_course")
    units = data.get("units", [])
    
    if not units:
        return jsonify({"error": "No units specified"}), 400
    
    if not db:
        return jsonify({"error": "Database not available"}), 500
    
    deleted = []
    failed = []
    
    for unit in units:
        unit_id = f"course_{course_id}_unit_{unit}"
        try:
            # Delete from Firestore
            db.collection("course_units").document(unit_id).delete()
            
            # Note: Vector DB cleanup is handled by cleanup_orphaned_file_metadata()
            deleted.append(unit)
        except Exception as e:
            print(f"Error deleting unit {unit}: {e}")
            failed.append(unit)
    
    # Trigger cleanup
    cleanup_orphaned_file_metadata()
    
    return jsonify({
        "success": True,
        "deleted": deleted,
        "failed": failed
    })

# ==============================================================================
# --- 19. ANALYTICS ROUTES ---
# ==============================================================================

@app.route("/api/teacher/student_analytics/<student_id>", methods=["GET"])
@teacher_required
def get_student_analytics(student_id):
    """Get analytics for a specific student"""
    tracker = StudentProgressTracker()
    analytics = tracker.analyze_student_performance(student_id)
    
    if not analytics:
        return jsonify({"error": "No data found for student"}), 404
    
    return jsonify(analytics)

@app.route("/api/teacher/class_analytics/<class_id>", methods=["GET"])
@teacher_required
def get_class_analytics(class_id):
    """Get analytics for an entire class"""
    tracker = StudentProgressTracker()
    analytics = tracker.get_class_analytics(class_id)
    
    if not analytics:
        return jsonify({"error": "No data found for class"}), 404
    
    return jsonify(analytics)

@app.route("/api/teacher/student_list", methods=["GET"])
@teacher_required
def get_student_list():
    """Get list of all students"""
    if not db:
        return jsonify({"error": "Database not available"}), 500
    
    try:
        students = db.collection("users")\
                     .where("role", "==", "student")\
                     .stream()
        
        student_list = []
        for student in students:
            data = student.to_dict()
            student_list.append({
                "id": student.id,
                "name": data.get("name", "Unknown"),
                "email": data.get("email", ""),
                "created_at": data.get("created_at", "")
            })
        
        return jsonify({
            "students": student_list,
            "count": len(student_list)
        })
        
    except Exception as e:
        print(f"Error fetching student list: {e}")
        return jsonify({"error": str(e)}), 500

# ==============================================================================
# --- 20. CHAT HISTORY ROUTES ---
# ==============================================================================

@app.route("/api/chats", methods=["GET"])
@login_required
def api_get_chats():
    return jsonify(get_web_chats_list(session['user_id']))

@app.route("/api/chats/<chat_id>", methods=["GET"])
@login_required
def api_get_chat(chat_id):
    chat = get_web_chat_detail(session['user_id'], chat_id)
    if not chat:
        return jsonify({"error": "Not found"}), 404
    return jsonify(chat)

@app.route("/api/chats/<chat_id>", methods=["DELETE"])
@login_required
def api_delete_chat(chat_id):
    delete_web_chat_session(session['user_id'], chat_id)
    return jsonify({"success": True})

@app.route("/api/chats", methods=["DELETE"])
@login_required
def api_clear_chats():
    clear_all_web_chats(session['user_id'])
    return jsonify({"success": True})

# ==============================================================================
# --- 21. AUTH ROUTES ---
# ==============================================================================

@app.route("/auth")
def auth_page():
    if 'user_id' in session:
        return redirect(url_for('index'))
    return render_template("auth.html")

@app.route("/api/auth/signup", methods=["POST"])
def api_signup():
    data     = request.get_json()
    name     = (data.get("name", "") or "").strip()
    email    = (data.get("email", "") or "").strip().lower()
    password = data.get("password", "") or ""
    role     = data.get("role", "student")

    if not name or not email or not password:
        return jsonify({"error": "All fields are required"}), 400
    if len(password) < 6:
        return jsonify({"error": "Password must be at least 6 characters"}), 400
    if role not in ("teacher", "student"):
        return jsonify({"error": "Invalid role"}), 400

    # Create user in Firebase Authentication
    try:
        user_record = fb_auth.create_user(
            email=email,
            password=password,
            display_name=name
        )
    except fb_auth.EmailAlreadyExistsError:
        return jsonify({"error": "Email already registered"}), 409
    except Exception as e:
        print(f"Firebase Auth signup error: {e}")
        return jsonify({"error": "Could not create account. Please try again."}), 500

    # Store role + name in Firestore keyed by Firebase UID
    if db:
        try:
            db.collection("users").document(user_record.uid).set({
                "name": name,
                "email": email,
                "role": role,
                "created_at": datetime.now().isoformat()
            })
        except Exception as e:
            print(f"Firestore write error after signup: {e}")

    session['user_id'] = user_record.uid
    session['role']    = role
    session['name']    = name
    return jsonify({"success": True, "role": role, "name": name})

@app.route("/api/auth/login", methods=["POST"])
def api_login():
    data     = request.get_json()
    email    = (data.get("email", "") or "").strip().lower()
    password = data.get("password", "") or ""

    if not email or not password:
        return jsonify({"error": "Email and password are required"}), 400
    if not FIREBASE_WEB_API_KEY:
        return jsonify({"error": "Server auth configuration missing"}), 500

    # Verify credentials via Firebase Auth REST API
    rest_resp = requests.post(
        f"https://identitytoolkit.googleapis.com/v1/accounts:signInWithPassword?key={FIREBASE_WEB_API_KEY}",
        json={"email": email, "password": password, "returnSecureToken": True},
        timeout=10
    )
    if not rest_resp.ok:
        err_msg = rest_resp.json().get("error", {}).get("message", "")
        if "EMAIL_NOT_FOUND" in err_msg or "INVALID_PASSWORD" in err_msg or "INVALID_LOGIN_CREDENTIALS" in err_msg:
            return jsonify({"error": "Invalid email or password"}), 401
        return jsonify({"error": "Login failed. Please try again."}), 401

    id_token = rest_resp.json().get("idToken")

    # Verify the ID token and extract UID (with retry for clock skew)
    decoded = None
    max_retries = 3
    for attempt in range(max_retries):
        try:
            decoded = fb_auth.verify_id_token(id_token)
            uid = decoded["uid"]
            break
        except Exception as e:
            error_str = str(e)
            print(f"Token verification attempt {attempt + 1}/{max_retries} error: {e}")
            
            # If it's a clock skew error, retry after a small delay
            if "Token used too early" in error_str and attempt < max_retries - 1:
                time.sleep(0.5)  # Small delay to allow clock sync
                continue
            elif attempt == max_retries - 1:
                # Final attempt failed
                print(f"Token verification failed after {max_retries} attempts")
                return jsonify({"error": "Authentication error. Please try again."}), 401
            else:
                # Non-timing error, fail immediately
                print(f"Token verification error: {e}")
                return jsonify({"error": "Authentication error. Please try again."}), 401
    
    if decoded is None:
        return jsonify({"error": "Authentication error. Please try again."}), 401

    # Fetch role from Firestore
    role, name = "student", email.split("@")[0]
    if db:
        try:
            doc = db.collection("users").document(uid).get()
            if doc.exists:
                d    = doc.to_dict()
                role = d.get("role", role)
                name = d.get("name", name)
        except Exception as e:
            print(f"Firestore role fetch error: {e}")

    session['user_id'] = uid
    session['role']    = role
    session['name']    = name
    return jsonify({"success": True, "role": role, "name": name})

@app.route("/api/auth/logout", methods=["POST"])
def api_logout():
    session.clear()
    return jsonify({"success": True})

@app.route("/coming-soon")
def coming_soon():
    if 'user_id' not in session:
        return redirect(url_for('auth_page'))
    return render_template("coming_soon.html", name=session.get('name', 'Student'))

# ==============================================================================
# --- 22. SERVE FRONTEND ---
# ==============================================================================
@app.route("/")
def index():
    if 'user_id' not in session:
        return redirect(url_for('auth_page'))
    return render_template("index.html", name=session.get('name', 'User'), role=session.get('role', 'student'))

# ==============================================================================
# --- 23. MAIN ---
# ==============================================================================
if __name__ == "__main__":
    os.makedirs("data", exist_ok=True)
    os.makedirs("vector_index", exist_ok=True)
    os.makedirs(DYNAMIC_VECTOR_INDEX_PATH, exist_ok=True)

    # Initialize Google APIs with OAuth 2.0 (replaces service account)
    print("🔄 Initializing Google APIs with OAuth 2.0...")
    init_google_apis_oauth()

    if db:
        print("Firebase connected successfully")
        cleanup_orphaned_file_metadata()
    else:
        print("Firebase not connected - running in local mode only")

    print("Sending startup template to users...")
    for number in student_phone_numbers:
        send_start_template(number)
    print("Finished sending templates.")

    print("Dynamic RAG system ready - using only user-uploaded documents")
    print("Teacher features enabled: Course Material Management, Question Paper Generation, Student Analytics")
    print("\n" + "="*60)
    print("📢 IMPORTANT: OAuth 2.0 Setup Required")
    print("="*60)
    print("1️⃣ Go to https://console.cloud.google.com/")
    print("2️⃣ Select project: sahayak-465916")
    print("3️⃣ Go to APIs & Services → Credentials")
    print("4️⃣ Create Credentials → OAuth 2.0 Client IDs")
    print("5️⃣ Choose 'Desktop Application' as type")
    print("6️⃣ Download JSON and save as 'credentials.json' in project root")
    print("7️⃣ Run this app - a browser will open for authentication")
    print("8️⃣ After authentication, 'token.pickle' will be created")
    print("="*60 + "\n")
    print("💡 Run /api/diagnose-drive to verify OAuth setup")
    print("Starting Flask app on port 5000...")
    app.run(port=5000, debug=False)