# ==============================================================================
# --- 1. IMPORTS (Cleaned Version)
# ==============================================================================

import os
from dotenv import load_dotenv
import torch
import requests
import traceback
from flask import Flask, request
from collections import defaultdict
import re
from datetime import datetime, timedelta
from io import BytesIO
import tempfile
import numpy as np
import whisper 
import base64
import moviepy
import moviepy.editor
from fpdf import FPDF
from pydub import AudioSegment
from moviepy.editor import ImageClip, concatenate_videoclips, AudioFileClip
from gtts import gTTS
from PIL import Image, ImageDraw
import json
import threading

# --- AI & Machine Learning ---
import google.generativeai as genai
from sklearn.metrics.pairwise import cosine_similarity
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_community.document_loaders import UnstructuredFileLoader
from langchain_community.vectorstores import FAISS
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.schema import Document

# --- Google Cloud APIs ---
import google.auth.transport.requests
from google.oauth2.credentials import Credentials
from google_auth_oauthlib.flow import InstalledAppFlow
from googleapiclient.discovery import build
from googleapiclient.errors import HttpError
from googleapiclient.http import MediaFileUpload

# --- Utility & File Generation ---
from fpdf import FPDF
from pydub import AudioSegment
from gtts import gTTS
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches

# --- PDF and Document Processing ---
import fitz  # PyMuPDF
from docx import Document as DocxDocument

# --- Firebase Imports ---
import firebase_admin
from firebase_admin import credentials, firestore
import uuid
from datetime import datetime

#
# ---------------------------------
# --- MODEL LOADING (at startup)
# ---------------------------------
#

print("Loading Whisper model...")
whisper_model = whisper.load_model("base")
print("Whisper model loaded successfully.")
os.environ["KMP_DUPLICATE_LIB_OK"] = "TRUE"

#
# ---------------------------------
# --- ENVIRONMENT & CONFIG
# ---------------------------------
#
app = Flask(__name__)

load_dotenv()
# --- WhatsApp & Meta Config ---
VERIFY_TOKEN = os.getenv("VERIFY_TOKEN")
ACCESS_TOKEN = os.getenv("ACCESS_TOKEN")
PHONE_NUMBER_ID = os.getenv("PHONE_NUMBER_ID")
TEMPLATE_NAME = os.getenv("TEMPLATE_NAME")

# Google & GenAI
GENAI_API_KEY = os.getenv("GENAI_API_KEY")
GOOGLE_FORMS_SERVICE_ACCOUNT_KEY_PATH = os.getenv("GOOGLE_FORMS_SERVICE_ACCOUNT_KEY_PATH")
GOOGLE_FORMS_TARGET_FOLDER_ID = os.getenv("GOOGLE_FORMS_TARGET_FOLDER_ID")
CLASSROOM_COURSE_ID = os.getenv("CLASSROOM_COURSE_ID")

print("Verify Token:", VERIFY_TOKEN)
print("Configs loaded successfully")

# --- List of users to notify on startup ---
student_phone_numbers = [
    "916379613654",
    "918870420449",
    "917530086388"
]

# --- Google AI & API Config ---
genai.configure(api_key=GENAI_API_KEY)
GOOGLE_FORMS_SERVICE_ACCOUNT_KEY_PATH = 'sahayak-465916-8bf5ddce5515.json'
GOOGLE_FORMS_TARGET_FOLDER_ID = '11vL5AgJiLYbX6fgAJFoMEjea9jiR1WMk'
CLASSROOM_COURSE_ID = '791255014049'

# --- State & Memory Management ---
user_states = defaultdict(str)
user_temp_data = defaultdict(dict)
user_memory = defaultdict(list)
MAX_HISTORY = 50

# --- Dynamic RAG Configuration ---
DYNAMIC_VECTOR_INDEX_PATH = "vector_index/dynamic_uploads"
os.makedirs(DYNAMIC_VECTOR_INDEX_PATH, exist_ok=True)

embeddings_model_for_classification = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
dynamic_embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")

# --- API Service Clients (Initialized later) ---
forms_service = None
drive_service = None
classroom_service = None
speech_client = None
creds = None
calendar_service = None

#
# ==============================================================================
# --- DYNAMIC RAG FUNCTIONS
# ==============================================================================
#

def extract_text_from_file(file_bytes, filename):
    """Extract text from various file formats."""
    try:
        file_extension = filename.lower().split('.')[-1]
        
        # Create temporary file
        with tempfile.NamedTemporaryFile(delete=False, suffix=f'.{file_extension}') as temp_file:
            temp_file.write(file_bytes)
            temp_file_path = temp_file.name
        
        text = ""
        
        if file_extension == 'pdf':
            # Using PyMuPDF for PDF extraction
            doc = fitz.open(temp_file_path)
            for  page in doc:
                text += page.get_text()
            doc.close()
            
        elif file_extension in ['docx', 'doc']:
            # Using python-docx for Word documents
            doc = DocxDocument(temp_file_path)
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
                
        elif file_extension in ['pptx', 'ppt']:
            # Using python-pptx for PowerPoint
            prs = Presentation(temp_file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                        
        elif file_extension == 'txt':
            with open(temp_file_path, 'r', encoding='utf-8') as f:
                text = f.read()
                
        else:
            # Fallback to unstructured for other formats
            try:
                loader = UnstructuredFileLoader(temp_file_path)
                documents = loader.load()
                text = "\n".join([doc.page_content for doc in documents])
            except Exception as e:
                print(f"Error with unstructured loader: {e}")
                return None
        
        # Clean up temporary file
        os.unlink(temp_file_path)
        
        return text.strip() if text else None
        
    except Exception as e:
        print(f"Error extracting text from {filename}: {e}")
        # Clean up temporary file if it exists
        if 'temp_file_path' in locals() and os.path.exists(temp_file_path):
            try:
                os.unlink(temp_file_path)
            except:
                pass
        return None

def add_documents_to_vector_index(file_bytes, filename, metadata=None):
    """Add new documents to the dynamic vector index."""
    try:
        # Extract text from uploaded file
        text_content = extract_text_from_file(file_bytes, filename)
        if not text_content:
            return False, "Could not extract text from the uploaded file."
        
        # Prepare metadata
        if metadata is None:
            metadata = {}
        
        base_metadata = {
            'source': filename,
            'uploaded_at': datetime.now().isoformat(),
            'type': 'user_upload'
        }
        base_metadata.update(metadata)
        
        # Create document object
        document = Document(
            page_content=text_content,
            metadata=base_metadata
        )
        
        # Split text into chunks
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=100
        )
        split_documents = text_splitter.split_documents([document])
        
        # Load existing vector store or create new one
        if os.path.exists(DYNAMIC_VECTOR_INDEX_PATH) and os.listdir(DYNAMIC_VECTOR_INDEX_PATH):
            try:
                vector_store = FAISS.load_local(
                    DYNAMIC_VECTOR_INDEX_PATH, 
                    dynamic_embeddings, 
                    allow_dangerous_deserialization=True
                )
                # Add new documents
                vector_store.add_documents(split_documents)
                print(f"Added {len(split_documents)} chunks to existing vector store")
            except Exception as e:
                print(f"Error loading existing index, creating new: {e}")
                vector_store = FAISS.from_documents(split_documents, dynamic_embeddings)
        else:
            # Create new vector store
            vector_store = FAISS.from_documents(split_documents, dynamic_embeddings)
            print(f"Created new vector store with {len(split_documents)} chunks")
        
        # Save the updated vector store
        vector_store.save_local(DYNAMIC_VECTOR_INDEX_PATH)
        
        return True, f"Successfully added {filename} to knowledge base with {len(split_documents)} chunks."
        
    except Exception as e:
        print(f"Error adding document to vector index: {e}")
        traceback.print_exc()
        return False, f"Error processing file: {str(e)}"

def query_dynamic_rag(question, k=3):
    """Query ONLY dynamic user uploads (no static database)"""
    try:
        # Check if dynamic database exists and has content
        if not os.path.exists(DYNAMIC_VECTOR_INDEX_PATH) or not os.listdir(DYNAMIC_VECTOR_INDEX_PATH):
            return {
                "result": "❌ No materials uploaded yet. Please upload documents first using 'Upload Materials' option.", 
                "source": "no_content"
            }
        
        # Load only dynamic database
        dynamic_db = FAISS.load_local(
            DYNAMIC_VECTOR_INDEX_PATH, 
            dynamic_embeddings, 
            allow_dangerous_deserialization=True
        )
        dynamic_retriever = dynamic_db.as_retriever(search_kwargs={"k": k})
        dynamic_docs = dynamic_retriever.invoke(question)
        
        if not dynamic_docs:
            return {
                "result": "❌ No relevant information found in your uploaded documents. Try uploading more relevant materials.", 
                "source": "no_matches"
            }
        
        # Combine context from dynamic docs only
        context = "\n\n".join([doc.page_content for doc in dynamic_docs])
        
        # Identify sources
        sources = []
        for doc in dynamic_docs:
            if 'source' in doc.metadata:
                source_name = os.path.basename(doc.metadata['source'])
                if source_name not in sources:
                    sources.append(source_name)
        
        source_info = "\n\n📚 Sources: " + ", ".join(sources) if sources else ""
        
        # Generate answer using Gemini
        prompt = f"""Based ONLY on the following context from user-uploaded documents, answer the question accurately:

Context:
{context}

Question: {question}

If the context doesn't contain enough information to answer the question, clearly state that and suggest uploading relevant documents."""

        model = genai.GenerativeModel('gemini-2.0-flash')
        response = model.generate_content(prompt)
        
        final_answer = response.text + source_info
        return {"result": final_answer, "source": "dynamic_rag"}
        
    except Exception as e:
        print(f"Error in dynamic RAG query: {e}")
        return {"result": "❌ Error during knowledge base lookup. Please try again.", "source": "error"}

def get_uploaded_files_list():
    """Get list of files that ACTUALLY exist in the dynamic vector database"""
    try:
        if not os.path.exists(DYNAMIC_VECTOR_INDEX_PATH) or not os.listdir(DYNAMIC_VECTOR_INDEX_PATH):
            return []
        
        vector_store = FAISS.load_local(
            DYNAMIC_VECTOR_INDEX_PATH, 
            dynamic_embeddings, 
            allow_dangerous_deserialization=True
        )
        
        # Extract unique sources from metadata that actually exist in the vector store
        sources = set()
        
        # Method 1: Check docstore
        if hasattr(vector_store, 'docstore') and hasattr(vector_store.docstore, '_dict'):
            for doc_id, doc in vector_store.docstore._dict.items():
                if hasattr(doc, 'metadata') and 'source' in doc.metadata:
                    sources.add(doc.metadata['source'])
        
        # Method 2: Check index if available
        if hasattr(vector_store, 'index_to_docstore_id'):
            for doc_id in vector_store.index_to_docstore_id.values():
                doc = vector_store.docstore.search(doc_id)
                if doc and hasattr(doc, 'metadata') and 'source' in doc.metadata:
                    sources.add(doc.metadata['source'])
        
        return list(sources)
        
    except Exception as e:
        print(f"Error getting uploaded files list from vector DB: {e}")
        return []

def handle_document_upload(sender, media_id, filename):
    """Worker function to handle single document upload and processing."""
    print(f"Starting document upload task for {sender}, file: {filename}")
    send_whatsapp_message(sender, f"Received {filename}. Processing and adding to knowledge base...")
    
    if media_url := get_media_url(media_id):
        if file_bytes := download_media_file(media_url):
            # Add document to vector index
            success, message = add_documents_to_vector_index(file_bytes, filename)
            
            if success:
                # Save file metadata to Firebase ONLY if successfully added to vector DB
                save_uploaded_file_metadata(sender, filename, len(file_bytes), message.split(' ')[-2] if 'chunks' in message else 0)
                send_whatsapp_message(sender, f"✅ Success! {message}\n\nI can now answer questions based on this content.")
            else:
                send_whatsapp_message(sender, f"❌ Failed to process file: {message}")
        else:
            send_whatsapp_message(sender, "❌ Sorry, I couldn't download the file.")
    else:
        send_whatsapp_message(sender, "❌ Sorry, I couldn't access the file.")
    
    # Clear state and show menu ONLY AFTER processing is complete
    if sender in user_states:
        del user_states[sender]
        clear_user_state_from_firebase(sender)
    end_conversation_and_show_menu(sender, None)
    print(f"Document upload task for {sender} finished.")

def get_upload_queue_status(sender):
    """Get current upload queue status"""
    if sender in user_temp_data and 'bulk_upload_items' in user_temp_data[sender]:
        items = user_temp_data[sender]['bulk_upload_items']
        return len(items), [item['filename'] for item in items]
    return 0, []

def check_upload_queue(sender):
    """Check current upload queue"""
    count, files = get_upload_queue_status(sender)
    if count > 0:
        file_list = "\n".join([f"• {f}" for f in files[:3]])
        if count > 3:
            file_list += f"\n• ... and {count - 3} more"
        send_whatsapp_message(sender, f"📋 Upload Queue: {count} files\n{file_list}")
    else:
        send_whatsapp_message(sender, "📋 No files in upload queue")

#
# ---------------------------------
# --- API INITIALIZATION
# ---------------------------------
#
SCOPES = [
    'https://www.googleapis.com/auth/classroom.courses',
    'https://www.googleapis.com/auth/classroom.coursework.students',
    'https://www.googleapis.com/auth/classroom.announcements',
    'https://www.googleapis.com/auth/drive',
    'https://www.googleapis.com/auth/forms.body',
    'https://www.googleapis.com/auth/calendar.events'
]

def init_google_apis():
    """Authenticates user and initializes Google API clients with proper scopes"""
    global classroom_service, drive_service, forms_service, calendar_service, creds
    
    if os.path.exists('token.json'):
        creds = Credentials.from_authorized_user_file('token.json', SCOPES)
    if not creds or not creds.valid:
        if creds and creds.expired and creds.refresh_token:
            creds.refresh(google.auth.transport.requests.Request())
        else:
            flow = InstalledAppFlow.from_client_secrets_file('credentials.json', SCOPES)
            creds = flow.run_local_server(port=0)
        with open('token.json', 'w') as token:
            token.write(creds.to_json())
    
    try:
        drive_service = build('drive', 'v3', credentials=creds)
        classroom_service = build('classroom', 'v1', credentials=creds)
        forms_service = build('forms', 'v1', credentials=creds)
        calendar_service = build('calendar', 'v3', credentials=creds)
        print("Google APIs initialized successfully with proper scopes.")
    except Exception as e:
        print(f"Error initializing Google APIs: {e}")

def send_start_template(to_number):
    """Sends the approved 'start' template message to a phone number."""
    url = f"https://graph.facebook.com/v19.0/{PHONE_NUMBER_ID}/messages"
    headers = {"Authorization": f"Bearer {ACCESS_TOKEN}", "Content-Type": "application/json"}
    payload = {
        "messaging_product": "whatsapp",
        "to": to_number,
        "type": "template",
        "template": {
            "name": "start_conversation_prompt",
            "language": {
                "code": "en"
            }
        }
    }
    try:
        response = requests.post(url, headers=headers, json=payload)
        response.raise_for_status()
        print(f"Successfully sent 'START' template to {to_number}.")
    except requests.exceptions.RequestException as e:
        print(f"Failed to send template to {to_number}: {e}")
        print(f"Response: {response.text}")

# ---------------------------------
# --- WHATSAPP & MEMORY FUNCTIONS
# ---------------------------------
#
def append_to_memory(user_id, role, content):
    """
    Appends a message to a user's conversation history in the correct
    format for the Gemini API.
    """
    # The Gemini API expects "user" and "model" as roles.
    api_role = "model" if role == "assistant" else "user"
    
    # This uses the 'parts' key required by the API instead of 'content'.
    user_memory[user_id].append({"role": api_role, "parts": [str(content)]})
    
    if len(user_memory[user_id]) > MAX_HISTORY:
        user_memory[user_id].pop(0)

def send_whatsapp_message(to, message):
    url = f"https://graph.facebook.com/v19.0/{PHONE_NUMBER_ID}/messages"
    headers = {"Authorization": f"Bearer {ACCESS_TOKEN}", "Content-Type": "application/json"}
    payload = {"messaging_product": "whatsapp", "to": to, "text": {"body": message}}
    try:
        res = requests.post(url, headers=headers, json=payload)
        res.raise_for_status()
        print(f"WhatsApp text message sent to {to}: {message[:75]}...")
    except Exception as e:
        print(f"Failed to send WhatsApp text message: {e}")

def send_whatsapp_document(to, file_bytes, filename):
    media_url = f"https://graph.facebook.com/v19.0/{PHONE_NUMBER_ID}/media"
    headers = {"Authorization": f"Bearer {ACCESS_TOKEN}"}
    files_and_data = {
        'file': (filename, file_bytes, 'application/pdf'),
        'messaging_product': (None, 'whatsapp'), 'type': (None, 'document')
    }
    try:
        media_res = requests.post(media_url, headers=headers, files=files_and_data)
        media_res.raise_for_status()
        media_id = media_res.json().get("id")
        if not media_id:
            raise Exception("Failed to get media ID from upload response.")
    except Exception as e:
        print(f"Failed to upload PDF to WhatsApp: {e}")
        return
    
    message_url = f"https://graph.facebook.com/v19.0/{PHONE_NUMBER_ID}/messages"
    message_headers = {"Authorization": f"Bearer {ACCESS_TOKEN}", "Content-Type": "application/json"}
    message_payload = {"messaging_product": "whatsapp", "to": to, "type": "document", "document": {"id": media_id, "filename": filename}}
    try:
        msg_res = requests.post(message_url, headers=message_headers, json=message_payload)
        msg_res.raise_for_status()
        print(f"WhatsApp document sent to {to}: {filename}")
    except Exception as e:
        print(f"Failed to send WhatsApp document message: {e}")

def send_menu_message(to, text, options):
    """Sends an interactive List Message."""
    url = f"https://graph.facebook.com/v19.0/{PHONE_NUMBER_ID}/messages"
    headers = {"Authorization": f"Bearer {ACCESS_TOKEN}", "Content-Type": "application/json"}
    
    rows = [{"id": f"option_{i}", "title": option} for i, option in enumerate(options)]
    
    payload = {
        "messaging_product": "whatsapp",
        "to": to,
        "type": "interactive",
        "interactive": {
            "type": "list",
            "header": {"type": "text", "text": "Main Menu"},
            "body": {"text": text},
            "footer": {"text": "Please choose an option"},
            "action": {
                "button": "Choose an Option",
                "sections": [{"title": "Available Actions", "rows": rows}]
            }
        }
    }
    try:
        requests.post(url, headers=headers, data=json.dumps(payload)).raise_for_status()
        print(f"List Menu sent to {to}.")
    except Exception as e:
        print(f"Failed to send List Menu: {e}")

def send_whatsapp_video(to, video_bytes, filename):
    media_url = f"https://graph.facebook.com/v19.0/{PHONE_NUMBER_ID}/media"
    headers = {"Authorization": f"Bearer {ACCESS_TOKEN}"}
    files = {'file': (filename, video_bytes, 'video/mp4'), 'messaging_product': (None, 'whatsapp')}
    try:
        media_res = requests.post(media_url, headers=headers, files=files)
        media_res.raise_for_status()
        media_id = media_res.json().get("id")
        if not media_id: return "failure"
        
        msg_url = f"https://graph.facebook.com/v19.0/{PHONE_NUMBER_ID}/messages"
        msg_payload = {"messaging_product": "whatsapp", "to": to, "type": "video", "video": {"id": media_id}}
        msg_res = requests.post(msg_url, headers=headers, json=msg_payload)
        msg_res.raise_for_status()
        return "success"
    except Exception as e:
        print(f"Failed to send video: {e}")
        return "failure"
    
def handle_audio_task(sender, media_id):
    """Worker function to handle audio processing in a background thread."""
    print(f"Starting audio task for {sender}")
    send_whatsapp_message(sender, "Got your voice message, listening... 🎤")
    if media_url := get_media_url(media_id):
        if audio_bytes := download_media_file(media_url):
            if transcribed_text := transcribe_audio(audio_bytes):
                process_message(sender, transcribed_text)
            else:
                send_whatsapp_message(sender, "Sorry, I couldn't understand your audio.")
        else:
            send_whatsapp_message(sender, "Sorry, I couldn't download your voice message.")
    print(f"Audio task for {sender} finished.")
    
def check_topic_relevance(topic):
    """Check if a topic is relevant to the uploaded RAG content with improved accuracy"""
    try:
        # Query the dynamic RAG to see if it has information about this topic
        test_query = f"What information do you have about {topic}? Please be specific about what content exists."
        response = query_dynamic_rag(test_query)
        
        # More robust relevance checking
        if response["source"] in ["dynamic_rag", "no_matches"]:
            response_lower = response["result"].lower()
            
            # Check for clear indicators of no relevant content
            no_content_indicators = [
                "no relevant information", 
                "no materials uploaded", 
                "doesn't contain information",
                "couldn't find information",
                "no information found",
                "not enough information"
            ]
            
            has_content_indicators = [
                "according to",
                "based on the",
                "the document mentions",
                "context provides",
                "sources:",
                "contains information about"
            ]
            
            # Check for negative indicators first
            if any(indicator in response_lower for indicator in no_content_indicators):
                return False, "no_matches"
            
            # Check for positive indicators
            if any(indicator in response_lower for indicator in has_content_indicators):
                return True, "has_content"
            
            # Default to no content if ambiguous
            return False, "ambiguous"
            
        return False, "error"
    except Exception as e:
        print(f"Error checking topic relevance: {e}")
        return False, "error"


def handle_image_task(sender, media_id, prompt):
    """Worker function to handle image processing in a background thread."""
    print(f"Starting image task for {sender}")
    send_whatsapp_message(sender, "🖼️ Analyzing your image...")
    
    if media_url := get_media_url(media_id):
        if image_bytes := download_media_file(media_url):
            # Validate image size before processing
            if len(image_bytes) > 20 * 1024 * 1024:
                send_whatsapp_message(sender, "❌ Image is too large (over 20MB). Please send a smaller image.")
                end_conversation_and_show_menu(sender, None)
                return
                
            if 'podcast' in prompt.lower() or 'audio' in prompt.lower():
                ocr_prompt = "Extract all the text from this image. Do not summarize or explain, just provide the raw text."
                ocr_response = query_gemini_vision(ocr_prompt, image_bytes)
                
                if ocr_response.get("source") == "gemini_vision":
                    extracted_text = ocr_response["result"]
                    if len(extracted_text.strip()) > 50:  # Only create audio if we have substantial text
                        if audio_bytes := generate_voiceover(extracted_text):
                            send_whatsapp_audio(sender, audio_bytes, "podcast.mp3")
                            end_conversation_and_show_menu(sender, "✅ Here is your podcast from the image text!")
                        else:
                            send_whatsapp_message(sender, "❌ I extracted the text, but couldn't create the audio.")
                            end_conversation_and_show_menu(sender, None)
                    else:
                        send_whatsapp_message(sender, "❌ Not enough text found in the image to create a podcast.")
                        end_conversation_and_show_menu(sender, None)
                else:
                    send_whatsapp_message(sender, ocr_response["result"])
                    end_conversation_and_show_menu(sender, None)
            else:
                response = query_gemini_vision(prompt, image_bytes)
                send_whatsapp_message(sender, response["result"])
                end_conversation_and_show_menu(sender, None)
        else:
            send_whatsapp_message(sender, "❌ Sorry, I couldn't download the image. Please try again.")
            end_conversation_and_show_menu(sender, None)
    else:
        send_whatsapp_message(sender, "❌ Sorry, I couldn't access the image. Please try again.")
        end_conversation_and_show_menu(sender, None)
    
    print(f"Image task for {sender} finished.")

# ==============================================================================
# --- FIREBASE CONFIGURATION
# ==============================================================================

# Initialize Firebase
try:
    if not firebase_admin._apps:
        cred = credentials.Certificate('firebase_config.json')
        firebase_admin.initialize_app(cred)
    db = firestore.client()
    print("Firebase initialized successfully")
except Exception as e:
    print(f"Firebase initialization error: {e}")
    db = None

# ==============================================================================
# --- FIREBASE MEMORY MANAGEMENT
# ==============================================================================

def get_user_memory_from_firebase(user_id):
    """Retrieve user memory from Firebase"""
    if not db:
        return []
    
    try:
        doc_ref = db.collection('user_memories').document(user_id)
        doc = doc_ref.get()
        
        if doc.exists:
            data = doc.to_dict()
            return data.get('conversation_history', [])
        else:
            # Create new user document
            doc_ref.set({
                'user_id': user_id,
                'conversation_history': [],
                'created_at': datetime.now(),
                'updated_at': datetime.now()
            })
            return []
    except Exception as e:
        print(f"Error getting user memory from Firebase: {e}")
        return []

def save_user_memory_to_firebase(user_id, conversation_history):
    """Save user memory to Firebase"""
    if not db:
        return False
    
    try:
        doc_ref = db.collection('user_memories').document(user_id)
        doc_ref.set({
            'user_id': user_id,
            'conversation_history': conversation_history,
            'updated_at': datetime.now()
        }, merge=True)
        return True
    except Exception as e:
        print(f"Error saving user memory to Firebase: {e}")
        return False

def append_to_memory(user_id, role, content):
    """
    Appends a message to a user's conversation history and saves to Firebase
    """
    # The Gemini API expects "user" and "model" as roles
    api_role = "model" if role == "assistant" else "user"
    
    # Get existing memory from Firebase
    if user_id not in user_memory:
        user_memory[user_id] = get_user_memory_from_firebase(user_id)
    
    # Append new message
    user_memory[user_id].append({"role": api_role, "parts": [str(content)]})
    
    # Limit history length
    if len(user_memory[user_id]) > MAX_HISTORY:
        user_memory[user_id] = user_memory[user_id][-MAX_HISTORY:]
    
    # Save to Firebase
    save_user_memory_to_firebase(user_id, user_memory[user_id])

def get_user_state_from_firebase(user_id):
    """Retrieve user state from Firebase"""
    if not db:
        return "", {}
    
    try:
        doc_ref = db.collection('user_states').document(user_id)
        doc = doc_ref.get()
        
        if doc.exists:
            data = doc.to_dict()
            return data.get('current_state', ''), data.get('temp_data', {})
        else:
            # Create new user state document
            doc_ref.set({
                'user_id': user_id,
                'current_state': '',
                'temp_data': {},
                'created_at': datetime.now(),
                'updated_at': datetime.now()
            })
            return "", {}
    except Exception as e:
        print(f"Error getting user state from Firebase: {e}")
        return "", {}

def save_user_state_to_firebase(user_id, current_state, temp_data):
    """Save user state to Firebase"""
    if not db:
        return False
    
    try:
        doc_ref = db.collection('user_states').document(user_id)
        doc_ref.set({
            'user_id': user_id,
            'current_state': current_state,
            'temp_data': temp_data,
            'updated_at': datetime.now()
        }, merge=True)
        return True
    except Exception as e:
        print(f"Error saving user state to Firebase: {e}")
        return False

def clear_user_state_from_firebase(user_id):
    """Clear user state from Firebase"""
    if not db:
        return False
    
    try:
        doc_ref = db.collection('user_states').document(user_id)
        doc_ref.set({
            'user_id': user_id,
            'current_state': '',
            'temp_data': {},
            'updated_at': datetime.now()
        }, merge=True)
        return True
    except Exception as e:
        print(f"Error clearing user state from Firebase: {e}")
        return False

def save_uploaded_file_metadata(user_id, filename, file_size, chunks_count):
    """Save uploaded file metadata to Firebase"""
    if not db:
        return False
    
    try:
        file_id = str(uuid.uuid4())
        doc_ref = db.collection('uploaded_files').document(file_id)
        doc_ref.set({
            'file_id': file_id,
            'user_id': user_id,
            'filename': filename,
            'file_size': file_size,
            'chunks_count': chunks_count,
            'uploaded_at': datetime.now()
        })
        return True
    except Exception as e:
        print(f"Error saving file metadata to Firebase: {e}")
        return False

def get_user_uploaded_files(user_id):
    """Get list of files uploaded by a user that exist in BOTH Firebase AND vector DB"""
    if not db:
        return []
    
    try:
        # First get files from Firebase
        docs = db.collection('uploaded_files').where('user_id', '==', user_id).stream()
        firebase_files = []
        for doc in docs:
            file_data = doc.to_dict()
            firebase_files.append({
                'filename': file_data.get('filename'),
                'uploaded_at': file_data.get('uploaded_at'),
                'file_size': file_data.get('file_size'),
                'file_id': file_data.get('file_id')
            })
        
        # Now get files that actually exist in vector DB
        vector_db_files = get_uploaded_files_list()
        vector_db_filenames = [os.path.basename(f) for f in vector_db_files]
        
        # Only return files that exist in both Firebase AND vector DB
        valid_files = []
        for file in firebase_files:
            if file['filename'] in vector_db_filenames:
                valid_files.append(file)
        
        return valid_files
        
    except Exception as e:
        print(f"Error getting user uploaded files: {e}")
        return []

# ==============================================================================
# --- UPDATED STATE MANAGEMENT WITH FIREBASE
# ==============================================================================

def initialize_user_session(sender):
    """Initialize user session with data from Firebase"""
    # Load memory from Firebase
    if sender not in user_memory:
        user_memory[sender] = get_user_memory_from_firebase(sender)
    
    # Load state from Firebase
    if sender not in user_states or sender not in user_temp_data:
        current_state, temp_data = get_user_state_from_firebase(sender)
        user_states[sender] = current_state
        user_temp_data[sender] = temp_data

def update_user_state(sender, state=None, temp_data=None):
    """Update user state and save to Firebase"""
    if state is not None:
        user_states[sender] = state
    
    if temp_data is not None:
        user_temp_data[sender].update(temp_data)
    
    # Save to Firebase
    save_user_state_to_firebase(sender, user_states.get(sender, ''), user_temp_data.get(sender, {}))

def clear_user_session(sender):
    """Clear user session and update Firebase"""
    user_states.pop(sender, None)
    user_temp_data.pop(sender, None)
    clear_user_state_from_firebase(sender)

# ---------------------------------
# --- VOICE INPUT HELPER FUNCTIONS
# ---------------------------------
#
def get_media_url(media_id):
    """Retrieves the downloadable URL for a piece of media from WhatsApp."""
    url = f"https://graph.facebook.com/v19.0/{media_id}/"
    headers = {"Authorization": f"Bearer {ACCESS_TOKEN}"}
    try:
        res = requests.get(url, headers=headers)
        res.raise_for_status()
        return res.json().get("url")
    except Exception as e:
        print(f"Failed to get media URL for ID {media_id}: {e}")
        return None
    
def handle_summary_task(sender, media_id):
    """Worker for the Image-to-Summary feature."""
    print(f"Starting summary task for {sender}")
    send_whatsapp_message(sender, "📝 Processing your image for the summary...")
    
    if media_url := get_media_url(media_id):
        if image_bytes := download_media_file(media_url):
            # Validate image size
            if len(image_bytes) > 20 * 1024 * 1024:
                send_whatsapp_message(sender, "❌ Image is too large (over 20MB). Please send a smaller image.")
                end_conversation_and_show_menu(sender, None)
                return
                
            ocr_prompt = "Extract all text from this image. Do not summarize or explain."
            ocr_response = query_gemini_vision(ocr_prompt, image_bytes)
            
            if ocr_response.get("source") == "gemini_vision":
                extracted_text = ocr_response["result"]
                if len(extracted_text.strip()) > 10:  # Only summarize if we have text
                    summary_prompt = f"Please provide a concise summary of the following text: {extracted_text}"
                    summary_response = query_gemini(summary_prompt, [])
                    end_conversation_and_show_menu(sender, f"📋 Summary:\n\n{summary_response['result']}")
                else:
                    end_conversation_and_show_menu(sender, "❌ No readable text found in the image.")
            else:
                end_conversation_and_show_menu(sender, ocr_response["result"])
        else:
            end_conversation_and_show_menu(sender, "❌ Sorry, I couldn't download the image.")
    else:
        end_conversation_and_show_menu(sender, "❌ Sorry, I couldn't access the image.")
    
    print(f"Summary task for {sender} finished.")

def download_media_file(media_url):
    """Downloads the media file from the given URL."""
    headers = {"Authorization": f"Bearer {ACCESS_TOKEN}"}
    try:
        res = requests.get(media_url, headers=headers)
        res.raise_for_status()
        return res.content
    except Exception as e:
        print(f"Failed to download media from URL: {e}")
        return None

def transcribe_audio(audio_bytes):
    """Transcribes audio bytes using the self-hosted Whisper model."""
    print("Transcribing audio with local Whisper model...")
    
    temp_file_path = None
    try:
        # Create a temporary file, get its path, and then close it
        # so that FFmpeg can access it without a permission error.
        with tempfile.NamedTemporaryFile(suffix=".ogg", delete=False) as temp_file:
            temp_file.write(audio_bytes)
            temp_file_path = temp_file.name
        
        # Load the audio from the file path to get a NumPy array
        audio_np_array = whisper.load_audio(temp_file_path)

        # Transcribe the NumPy array
        result = whisper_model.transcribe(audio_np_array, fp16=False)
        
        transcript = result.get("text", "")
        print(f"Whisper transcription successful: '{transcript}'")
        return transcript

    except Exception as e:
        print(f"Error during local Whisper transcription: {e}")
        traceback.print_exc()
        return None
    finally:
        # Manually delete the temporary file after we're done with it
        if temp_file_path and os.path.exists(temp_file_path):
            os.remove(temp_file_path)

def end_conversation_and_show_menu(sender, final_message):
    """Sends a final message, clears user state safely, and shows the main menu."""
    if final_message:
        send_whatsapp_message(sender, final_message)
    
    clear_user_session(sender)

    menu_text = "What would you like to do next?"
    options = [
        "Ask Question", 
        "Create Worksheet", 
        "Create PPT", 
        "Upload Materials",
        "View Uploaded Files", 
        "Podcast from Image", 
        "Summary from Image"
    ]
    send_menu_message(sender, menu_text, options)
    append_to_memory(sender, "assistant", "Task complete. Displayed main menu.")

def handle_bulk_document_upload(sender, media_items):
    """Handle multiple document uploads at once and send menu only after ALL processing"""
    print(f"Starting bulk upload for {sender}, {len(media_items)} files")
    
    successful_uploads = 0
    failed_uploads = 0
    processed_files = []
    
    send_whatsapp_message(sender, f"📦 Processing {len(media_items)} files... This may take a few minutes.")
    
    for i, media_item in enumerate(media_items, 1):
        media_id = media_item.get('id')
        filename = media_item.get('filename', f'document_{i}')
        
        send_whatsapp_message(sender, f"📄 Processing ({i}/{len(media_items)}): {filename}")
        
        if media_url := get_media_url(media_id):
            if file_bytes := download_media_file(media_url):
                success, message = add_documents_to_vector_index(file_bytes, filename)
                
                if success:
                    successful_uploads += 1
                    processed_files.append(filename)
                    # Save file metadata to Firebase if available
                    if db:
                        save_uploaded_file_metadata(sender, filename, len(file_bytes), message.split(' ')[-2] if 'chunks' in message else 0)
                else:
                    failed_uploads += 1
                    send_whatsapp_message(sender, f"❌ Failed to process {filename}: {message}")
            else:
                failed_uploads += 1
                send_whatsapp_message(sender, f"❌ Couldn't download {filename}")
        else:
            failed_uploads += 1
            send_whatsapp_message(sender, f"❌ Couldn't access {filename}")
    
    # Send summary ONLY after ALL files are processed
    summary_msg = f"✅ Bulk Upload Complete!\n\n"
    summary_msg += f"📊 Success: {successful_uploads} files\n"
    if failed_uploads > 0:
        summary_msg += f"❌ Failed: {failed_uploads} files\n"
    
    if successful_uploads > 0:
        summary_msg += f"\n📚 Added to knowledge base:\n"
        for file in processed_files[:5]:  # Show first 5 files
            summary_msg += f"• {file}\n"
        if len(processed_files) > 5:
            summary_msg += f"• ... and {len(processed_files) - 5} more\n"
        
        summary_msg += f"\nYou can now ask questions based on this content!"
    
    send_whatsapp_message(sender, summary_msg)
    
    # Clear state and show menu ONLY AFTER all processing is done
    if sender in user_states:
        del user_states[sender]
    end_conversation_and_show_menu(sender, None)
    
    print(f"Bulk upload for {sender} finished. Success: {successful_uploads}, Failed: {failed_uploads}")

def upload_file_to_drive(file_path, folder_id=None):
    """Uploads a file to Google Drive, optionally to a specific folder"""
    if not drive_service:
        print("Google Drive API not initialized.")
        return None
    try:
        file_metadata = {'name': os.path.basename(file_path)}
        
        # Only add parent folder if it exists and is valid
        if folder_id and folder_id.strip():
            try:
                # Verify the folder exists before trying to use it
                drive_service.files().get(fileId=folder_id).execute()
                file_metadata['parents'] = [folder_id]
                print(f"Uploading to folder: {folder_id}")
            except HttpError as e:
                print(f"Folder {folder_id} not found, uploading to root directory")
                # Don't add parents if folder doesn't exist
        
        media = MediaFileUpload(file_path, mimetype='application/pdf')
        file = drive_service.files().create(body=file_metadata, media_body=media, fields='id').execute()
        print(f"File uploaded to Drive with ID: {file.get('id')}")
        return file.get('id')
    except Exception as e:
        print(f"Error uploading file to Drive: {e}")
        return None

    
def post_announcement(sender, course_id, text, drive_file_id):
    if not classroom_service:
        send_whatsapp_message(sender, "Google Classroom API not initialized.")
        end_conversation_and_show_menu(sender, None)
        return {"result": "Google Classroom API not initialized.", "source": "api_error"}
    
    announcement = {
        'text': text,
        'materials': [{'driveFile': {'driveFile': {'id': drive_file_id}, 'shareMode': 'VIEW'}}],
        'state': 'PUBLISHED'
    }
    try:
        announcement = classroom_service.courses().announcements().create(courseId=course_id, body=announcement).execute()
        send_whatsapp_message(sender, f"✅ Successfully posted announcement in Google Classroom!")
        end_conversation_and_show_menu(sender, None)
        return {"result": f"Successfully posted announcement in Google Classroom!", "source": "classroom_success"}
        
    except HttpError as api_error:
        print(f"Google Classroom API Error: {api_error.content}")
        send_whatsapp_message(sender, "❌ Couldn't post announcement. Please check permissions and Course ID.")
        end_conversation_and_show_menu(sender, None)
        return {"result": "Couldn't post announcement. Please check permissions and Course ID.", "source": "api_error"}
    except Exception as e:
        print(f"Unexpected error posting announcement: {e}")
        send_whatsapp_message(sender, "❌ An unexpected error occurred while posting the announcement.")
        end_conversation_and_show_menu(sender, None)
        return {"result": "An unexpected error occurred.", "source": "error"}
    
    
def post_assignment(sender, course_id, title, description, drive_file_id, due_date=None):
    """Create a classroom assignment without due dates"""
    if not classroom_service:
        send_whatsapp_message(sender, "Google Classroom API not initialized.")
        end_conversation_and_show_menu(sender, None)
        return {"result": "Google Classroom API not initialized.", "source": "api_error"}
    
    coursework = {
        'title': title,
        'description': description,
        'materials': [{'driveFile': {'driveFile': {'id': drive_file_id}, 'shareMode': 'VIEW'}}],
        'workType': 'ASSIGNMENT',
        'state': 'PUBLISHED'
    }
    
    # Skip due dates entirely as requested
    try:
        assignment = classroom_service.courses().courseWork().create(courseId=course_id, body=coursework).execute()
        send_whatsapp_message(sender, f"✅ Successfully created assignment '{assignment.get('title')}' in Google Classroom!")
        end_conversation_and_show_menu(sender, None)
        return {"result": f"Successfully created assignment '{assignment.get('title')}' in Google Classroom!", "source": "classroom_success"}
    except HttpError as api_error:
        print(f"Google Classroom API Error: {api_error.content}")
        send_whatsapp_message(sender, "❌ Couldn't create Classroom assignment. Please check permissions and Course ID.")
        end_conversation_and_show_menu(sender, None)
        return {"result": "Couldn't create Classroom assignment. Please check permissions and Course ID.", "source": "api_error"}
    except Exception as e:
        print(f"General error creating Classroom assignment: {e}")
        send_whatsapp_message(sender, "❌ An unexpected error occurred while creating the Classroom assignment.")
        end_conversation_and_show_menu(sender, None)
        return {"result": "An unexpected error occurred while creating the Classroom assignment.", "source": "error"}


def query_gemini(question, history):
    try:
        model = genai.GenerativeModel('gemini-2.0-flash')
        chat = model.start_chat(history=history)
        response = chat.send_message(question)
        return {"result": response.text}
    except Exception as e:
        print("Gemini general query error:", e)
        return {"result": "Sorry, I encountered an issue while processing your general question."}
    
def generate_video_script(topic, history):
    print(f"Generating video script for topic: {topic}")
    prompt = f"Create a short, simple video script explaining the basics of {topic} for a beginner. The script should have 3 key points. For each key point, provide a narration sentence and a simple visual description on a new line starting with 'VISUAL:'."
    try:
        model = genai.GenerativeModel('gemini-2.0-flash')
        chat = model.start_chat(history=history)
        response = chat.send_message(prompt)
        return response.text
    except Exception as e:
        print(f"Error generating video script: {e}")
        return None

def generate_images_for_script(script):
    print("Generating images for script...")
    image_paths = []
    visual_prompts = re.findall(r'VISUAL:\s*(.*)', script)
    for i, prompt in enumerate(visual_prompts):
        try:
            img = Image.new('RGB', (1280, 720), color = 'darkblue')
            d = ImageDraw.Draw(img)
            try:
                font = ImageFont.truetype("arial.ttf", 40)
            except IOError:
                font = ImageFont.load_default()
            
            # Simple text wrapping
            lines = []
            words = prompt.split()
            current_line = ""
            for word in words:
                if len(current_line + " " + word) < 50:
                    current_line += " " + word
                else:
                    lines.append(current_line)
                    current_line = word
            lines.append(current_line)

            y_text = 300
            for line in lines:
                d.text((100, y_text), line.strip(), font=font, fill=(255,255,255))
                y_text += 50

            img_path = f"temp_image_{i}.png"
            img.save(img_path)
            image_paths.append(img_path)
        except Exception as e:
            print(f"Error creating image {i}: {e}")
    return image_paths

def generate_video_with_modelscope(prompt, output_path):
    """Generates a video from a text prompt using the ModelScope model."""
    print(f"Generating video for prompt: '{prompt}'...")
    try:
        video_frames = video_pipe(prompt, num_inference_steps=25, num_frames=16).frames
        export_to_video(video_frames, output_path)
        return True
    except Exception as e:
        print(f"Error generating video with ModelScope: {e}")
        return False
    
def generate_voiceover(text):
    """Generates a voiceover from text and speeds it up."""
    print("Generating voiceover...")
    try:
        # Generate the initial audio
        tts = gTTS(text=text, lang='en', slow=False)
        fp = BytesIO()
        tts.write_to_fp(fp)
        fp.seek(0)
        
        # Load the audio with pydub and speed it up
        audio = AudioSegment.from_file(fp, format="mp3")
        sped_up_audio = audio.speedup(playback_speed=1.25)
        
        # Export the sped-up audio to a new BytesIO object
        output_fp = BytesIO()
        sped_up_audio.export(output_fp, format="mp3")
        output_fp.seek(0)
        
        return output_fp.read()
    except Exception as e:
        print(f"Error generating voiceover: {e}")
        return None
    
def create_calendar_event(title, due_date):
    if not calendar_service: return
    try:
        event = {
            'summary': f'Due: {title}',
            'description': 'An assignment is due.',
            'start': {'date': due_date.isoformat()},
            'end': {'date': due_date.isoformat()},
        }
        calendar_service.events().insert(calendarId='primary', body=event).execute()
        print("Calendar event created.")
    except Exception as e:
        print(f"Error creating calendar event: {e}")

#
# ---------------------------------
# --- CONTENT GENERATION FUNCTIONS
# ---------------------------------
#
def generate_mcq_questions_text(topic, num_questions, history):
    print(f"Attempting to generate {num_questions} TEXT MCQs on topic: '{topic}'")
    
    mcq_prompt = f"Generate {num_questions} ONLY multiple-choice questions on the topic of '{topic}'. For each question, provide 4 options (A, B, C, D) and clearly indicate the correct answer. Format each as: 'Question Number. Question Text\\n A) Option A\\n B) Option B\\n C) Option C\\n D) Option D\\n Correct Answer: X'. Provide ONLY the questions and answers."
    
    try:
        model = genai.GenerativeModel('gemini-2.0-flash')
        chat = model.start_chat(history=history)
        response = chat.send_message(mcq_prompt)
        return {"result": response.text, "source": "generated_form_questions_text"}
    except Exception as e:
        print(f"Error generating TEXT MCQs for '{topic}': {e}")
        return {"result": "Sorry, I couldn't generate the quiz questions.", "source": "error"}

def create_google_form_mcq(title, mcq_questions_text):
    """Creates a Google Form quiz from a string of MCQ questions."""
    if not forms_service:
        return {"result": "Google Forms API not initialized.", "source": "api_error"}

    try:
        form_metadata = {'info': {'title': title}}
        form = forms_service.forms().create(body=form_metadata).execute()
        form_id = form['formId']
        
        update = {"requests": [{"updateSettings": {"settings": {"quizSettings": {'isQuiz': True}}, "updateMask": "quizSettings"}}]}
        forms_service.forms().batchUpdate(formId=form_id, body=update).execute()

        question_requests = []
        # A more robust way to split into question blocks
        question_blocks = re.split(r'\n\s*\d+\.\s*', mcq_questions_text)
        
        # Filter out any empty strings or introductory text
        question_blocks = [b.strip() for b in question_blocks if b.strip() and "Instructions:" not in b and "Worksheet" not in b]

        for i, q_block in enumerate(question_blocks):
            lines = [line.strip() for line in q_block.split('\n') if line.strip()]
            if not lines:
                continue

            question_text = lines[0]
            options = []
            correct_answer_text = ""

            for line in lines[1:]:
                if re.match(r'^[a-dA-D][\.\)]', line):
                    options.append(line[2:].strip())
                elif "Correct Answer:" in line:
                    correct_answer_match = re.search(r'Correct Answer:\s*([a-dA-D])', line, re.IGNORECASE)
                    if correct_answer_match:
                        correct_answer_letter = correct_answer_match.group(1).upper()
                        correct_answer_index = ord(correct_answer_letter) - ord('A')
                        if 0 <= correct_answer_index < len(options):
                            correct_answer_text = options[correct_answer_index]

            if not question_text or len(options) < 2 or not correct_answer_text:
                print(f"Skipping malformed block: {q_block[:50]}...")
                continue

            choices = [{"value": opt} for opt in options]
            
            question_requests.append({
                "createItem": {
                    "item": {
                        "title": question_text,
                        "questionItem": {
                            "question": {
                                "required": True,
                                "choiceQuestion": {"type": "RADIO", "options": choices},
                                "grading": {"pointValue": 1, "correctAnswers": {"answers": [{"value": correct_answer_text}]}}
                            }
                        }
                    }, "location": {"index": i}
                }
            })

        if question_requests:
            forms_service.forms().batchUpdate(formId=form_id, body={"requests": question_requests}).execute()

        return {"result": f"I've created a Google Form quiz for you! Access it here: {form['responderUri']}", "source": "google_form_created"}

    except Exception as e:
        print(f"Error creating Google Form: {e}")
        traceback.print_exc()
        return {"result": "Sorry, I encountered an error while creating the Google Form.", "source": "error"}

def generate_worksheet_content_text(topic, num_items, worksheet_type, history):
    """
    Generates worksheet questions AND answers, separated by a unique string.
    """
    print(f"Generating {num_items} {worksheet_type} items and answers on: '{topic}'")
    
    # New prompt that asks for both questions and a clearly separated answer key
    prompt = (
        f"Generate a student worksheet with {num_items} {worksheet_type} questions on the topic of '{topic}'. "
        "The questions should be clear and suitable for a student. Do not provide the answers immediately after each question. "
        "After the questions, add a separator line exactly like this: '--- ANSWERS ---'. "
        "Then, after the separator, provide a numbered list of the corresponding answers. "
        "The final output must contain both the questions and the answers, separated by the '--- ANSWERS ---' line."
    )

    try:
        model = genai.GenerativeModel('gemini-2.0-flash')
        chat = model.start_chat(history=history)
        response = chat.send_message(prompt)
        return {"result": response.text, "source": "generated_worksheet_text"}
    except Exception as e:
        print(f"Error generating worksheet content: {e}")
        return {"result": "Sorry, I couldn't generate the worksheet content.", "source": "error"}
    
def generate_ppt_content(topic):
    prompt = f"Create the content for a 10-slide presentation on the topic of '{topic}'. The first slide should be a title slide. The next 9 slides should be content slides, each with a title and 3-4 bullet points. Format the output clearly, using 'SLIDE:' to mark the beginning of each slide."
    try:
        model = genai.GenerativeModel('gemini-2.0-flash')
        response = model.generate_content(prompt)
        return response.text
    except Exception as e:
        print(f"Error generating PPT content: {e}")
        return None
    
def create_ppt_file(title, content):
    try:
        prs = Presentation()
        # More robustly split the content into slides
        slides = re.split(r'\n*SLIDE:?\s*\d*\s*:?\n*', content)
        slides = [s.strip() for s in slides if s.strip()] # Clean up empty entries

        if not slides:
            print("Could not parse any slides from the generated content.")
            return None

        # Handle the first slide as the title slide
        title_slide_content = slides[0].split('\n')
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title_slide_content[0]
        if len(title_slide_content) > 1:
            slide.placeholders[1].text = "\n".join(title_slide_content[1:])
        
        # Handle the rest as content slides
        for slide_content in slides[1:]:
            lines = [line.strip() for line in slide_content.strip().split('\n') if line.strip()]
            if not lines: continue

            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = lines[0] # Assume first line is the title
            
            body_shape = slide.shapes.placeholders[1]
            tf = body_shape.text_frame
            tf.clear()
            
            for point in lines[1:]:
                # Remove leading bullet points like '*' or '-'
                cleaned_point = re.sub(r'^\s*[\*\-]\s*', '', point)
                p = tf.add_paragraph()
                p.text = cleaned_point
                p.level = 0

        fp = BytesIO()
        prs.save(fp)
        fp.seek(0)
        return fp.read()
    except Exception as e:
        print(f"Error creating PPT file: {e}")
        traceback.print_exc()
        return None

def send_whatsapp_ppt(to, file_bytes, filename):
    media_url = f"https://graph.facebook.com/v19.0/{PHONE_NUMBER_ID}/media"
    headers = {"Authorization": f"Bearer {ACCESS_TOKEN}"}
    files = {'file': (filename, file_bytes, 'application/vnd.openxmlformats-officedocument.presentationml.presentation'), 'messaging_product': (None, 'whatsapp')}
    try:
        media_res = requests.post(media_url, headers=headers, files=files)
        media_res.raise_for_status()
        media_id = media_res.json().get("id")
        if not media_id: return "failure"
        
        msg_url = f"https://graph.facebook.com/v19.0/{PHONE_NUMBER_ID}/messages"
        msg_headers = {"Authorization": f"Bearer {ACCESS_TOKEN}", "Content-Type": "application/json"}
        msg_payload = {"messaging_product": "whatsapp", "to": to, "type": "document", "document": {"id": media_id, "filename": filename}}
        requests.post(msg_url, headers=msg_headers, json=msg_payload).raise_for_status()
        return "success"
    except Exception as e:
        print(f"Failed to send PPT: {e}")
        return "failure"

def create_pdf_locally(title, content):
    """
    Creates a PDF file from a title and content, saving it to a temporary file
    and returning the raw bytes.
    """
    pdf_bytes = None
    temp_pdf_path = None
    try:
        # Create a temporary file to save the PDF to disk first
        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as temp_file:
            temp_pdf_path = temp_file.name

        # Initialize the PDF document
        pdf = FPDF()
        pdf.add_page()
        pdf.set_auto_page_break(auto=True, margin=15)
        
        # Add the title
        pdf.set_font("Helvetica", 'B', 16)
        pdf.multi_cell(0, 10, title, align='C')
        pdf.ln(10)
        
        # Add the main content
        pdf.set_font("Helvetica", size=11)
        # Clean the content to handle special characters that FPDF might not support
        cleaned_content = content.encode('latin-1', 'replace').decode('latin-1')
        pdf.multi_cell(0, 7, cleaned_content)
        
        # Save the PDF to the temporary file path
        pdf.output(temp_pdf_path)

        # Read the bytes from the saved file to ensure it's not empty
        with open(temp_pdf_path, 'rb') as f:
            pdf_bytes = f.read()

        return pdf_bytes
        
    except Exception as e:
        print(f"Error creating PDF: {e}")
        return None
    finally:
        # Clean up the temporary file from disk after we're done with it
        if temp_pdf_path and os.path.exists(temp_pdf_path):
            os.remove(temp_pdf_path)

#
# ---------------------------------
# --- CORE MESSAGE PROCESSOR
# ---------------------------------
#

def handle_final_classroom_post(sender, title):
    """Helper function to handle the final steps of posting to Classroom."""
    try:
        post_choice = user_temp_data[sender]['post_choice']
        questions = user_temp_data[sender]['questions_text']
        answers = user_temp_data[sender]['answers_text']
        
        pdf_content = questions
        pdf_filename = f"{title.replace(' ', '_')}_worksheet.pdf"
        
        if post_choice.lower() == "post with answers":
            pdf_content += f"\n\n--- ANSWERS ---\n{answers}"
            pdf_filename = f"{title.replace(' ', '_')}_with_answers.pdf"

        pdf_bytes = create_pdf_locally(title, pdf_content)
        
        if pdf_bytes:
            with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as temp_file:
                temp_file.write(pdf_bytes)
                temp_file_path = temp_file.name
            
            # Try to upload to the specified folder, but fall back to root if needed
            drive_file_id = upload_file_to_drive(temp_file_path, GOOGLE_FORMS_TARGET_FOLDER_ID)
            os.remove(temp_file_path)

            if drive_file_id:
                if post_choice.lower() != "don't post":
                    # Create assignment with the PDF
                    description = f"Worksheet on {title}"
                    if post_choice.lower() == "post with answers":
                        description += " (Includes answer key)"
                    
                    # Try to post to classroom
                    result = post_assignment(sender, CLASSROOM_COURSE_ID, title, description, drive_file_id, due_date=None)
                    
                    if result["source"] == "classroom_success":
                        send_whatsapp_message(sender, result['result'])
                    else:
                        # If classroom posting fails, provide the Drive link instead
                        drive_link = f"https://drive.google.com/file/d/{drive_file_id}/view"
                        send_whatsapp_message(sender, f"I couldn't post to Classroom due to permission issues. Here's the direct link to the file: {drive_link}")
                else:
                    send_whatsapp_message(sender, "Okay, I won't post to Classroom. The PDFs have been sent to you directly.")
            else:
                send_whatsapp_message(sender, "I created the PDF, but failed to upload it to Google Drive.")
        else:
            send_whatsapp_message(sender, "Sorry, I failed to create the final PDF for upload.")

    except KeyError:
        send_whatsapp_message(sender, "An error occurred. I've lost the worksheet details. Please start over.")
    except Exception as e:
        print(f"Unexpected error in handle_final_classroom_post: {e}")
        send_whatsapp_message(sender, "An unexpected error occurred. Please try again.")
    finally:
        if sender in user_states: del user_states[sender]
        if sender in user_temp_data: del user_temp_data[sender]
        
def query_gemini_vision(prompt, image_bytes):
    """Sends a prompt and an image to the Gemini Vision model with proper error handling."""
    print("Querying Gemini Vision...")
    try:
        # Check if image_bytes is valid
        if not image_bytes or len(image_bytes) == 0:
            return {"result": "❌ The image appears to be empty or corrupted.", "source": "error"}
        
        # Check image size (Gemini has limits)
        if len(image_bytes) > 20 * 1024 * 1024:  # 20MB limit
            return {"result": "❌ Image is too large. Please use images under 20MB.", "source": "error"}
        
        image_parts = [{"mime_type": "image/jpeg", "data": base64.b64encode(image_bytes).decode('utf-8')}]
        prompt_parts = [prompt] + image_parts
        
        model = genai.GenerativeModel('gemini-2.0-flash')
        response = model.generate_content(prompt_parts)
        
        # Better error handling for empty responses
        if not response.candidates:
            return {"result": "❌ Gemini couldn't process this image. The image might be blurry, contain restricted content, or be in an unsupported format.", "source": "error"}
        
        if response.parts:
            return {"result": response.text, "source": "gemini_vision"}
        else:
            # Check for safety blocks
            if response.prompt_feedback and response.prompt_feedback.block_reason:
                block_reason = response.prompt_feedback.block_reason.name
                return {"result": f"❌ Image was blocked due to: {block_reason}. Please try a different image.", "source": "error"}
            
            return {"result": "❌ I'm sorry, I couldn't understand the image. Please try a clearer image.", "source": "error"}
            
    except Exception as e:
        print(f"Error querying Gemini Vision: {e}")
        
        # More specific error messages
        if "image" in str(e).lower() and "size" in str(e).lower():
            return {"result": "❌ Image is too large or in an unsupported format. Please try a smaller image.", "source": "error"}
        elif "invalid" in str(e).lower():
            return {"result": "❌ Invalid image format. Please try a JPEG, PNG, or WebP image.", "source": "error"}
        elif "safety" in str(e).lower() or "block" in str(e).lower():
            return {"result": "❌ This image contains content that cannot be processed. Please try a different image.", "source": "error"}
        else:
            return {"result": "❌ Sorry, there was an error processing the image. Please try again with a different image.", "source": "error"}
    
def send_whatsapp_audio(to, audio_bytes, filename):
    media_url = f"https://graph.facebook.com/v19.0/{PHONE_NUMBER_ID}/media"
    headers = {"Authorization": f"Bearer {ACCESS_TOKEN}"}
    files = {'file': (filename, audio_bytes, 'audio/mpeg'), 'messaging_product': (None, 'whatsapp')}
    try:
        media_res = requests.post(media_url, headers=headers, files=files)
        media_res.raise_for_status()
        media_id = media_res.json().get("id")
        if not media_id: return "failure"

        msg_url = f"https://graph.facebook.com/v19.0/{PHONE_NUMBER_ID}/messages"
        msg_payload = {"messaging_product": "whatsapp", "to": to, "type": "audio", "audio": {"id": media_id}}
        requests.post(msg_url, headers=headers, json=msg_payload).raise_for_status()
        return "success"
    except Exception as e:
        print(f"Failed to send audio: {e}")
        return "failure"
    
def handle_podcast_task(sender, media_id):
    """Worker for the Image-to-Podcast feature."""
    print(f"Starting podcast task for {sender}")
    send_whatsapp_message(sender, "Processing your image for the podcast...")
    if media_url := get_media_url(media_id):
        if image_bytes := download_media_file(media_url):
            ocr_prompt = "Extract all the text from this image. Do not summarize or explain, just provide the raw text."
            ocr_response = query_gemini_vision(ocr_prompt, image_bytes)
            if ocr_response.get("source") == "gemini_vision":
                if audio_bytes := generate_voiceover(ocr_response["result"]):
                    send_whatsapp_audio(sender, audio_bytes, "podcast.mp3")
                    end_conversation_and_show_menu(sender, "Here is your podcast.")
                else:
                    end_conversation_and_show_menu(sender, "I extracted the text, but couldn't create the audio.")
            else:
                end_conversation_and_show_menu(sender, ocr_response["result"])
        else:
            end_conversation_and_show_menu(sender, "Sorry, I couldn't download the image.")
    print(f"Podcast task for {sender} finished.")
    
def generate_voiceover(text):
    print("Generating voiceover...")
    try:
        tts = gTTS(text=text, lang='en')
        fp = BytesIO()
        tts.write_to_fp(fp)
        fp.seek(0)
        return fp.read()
    except Exception as e:
        print(f"Error generating voiceover: {e}")
        return None

def cleanup_orphaned_file_metadata():
    """Remove Firebase entries for files that no longer exist in vector DB"""
    if not db:
        return
    
    try:
        # Get all files from Firebase
        docs = db.collection('uploaded_files').stream()
        vector_db_files = get_uploaded_files_list()
        vector_db_filenames = [os.path.basename(f) for f in vector_db_files]
        
        deleted_count = 0
        for doc in docs:
            file_data = doc.to_dict()
            filename = file_data.get('filename')
            
            # If file doesn't exist in vector DB, delete from Firebase
            if filename not in vector_db_filenames:
                db.collection('uploaded_files').document(doc.id).delete()
                deleted_count += 1
                print(f"Cleaned up orphaned file metadata: {filename}")
        
        if deleted_count > 0:
            print(f"Cleaned up {deleted_count} orphaned file metadata entries")
            
    except Exception as e:
        print(f"Error cleaning up orphaned file metadata: {e}")

def parse_deadline(text):
    text = text.lower()
    today = datetime.now()
    if "tomorrow" in text:
        return today + timedelta(days=1)
    if "next week" in text or "7 days" in text:
        return today + timedelta(days=7)
   
    return today + timedelta(days=2)

def send_interactive_message(to, text, buttons):
    """Sends an interactive message with up to 3 reply buttons with robust error handling."""
    url = f"https://graph.facebook.com/v19.0/{PHONE_NUMBER_ID}/messages"
    headers = {
        "Authorization": f"Bearer {ACCESS_TOKEN}", 
        "Content-Type": "application/json"
    }
    
    try:
        # Validate and format buttons properly
        button_payload = []
        for i, button_title in enumerate(buttons[:3]):  # Max 3 buttons
            # Clean button title - remove emojis and limit length
            clean_title = re.sub(r'[^\w\s\-()]', '', button_title)  # Remove special chars
            clean_title = clean_title[:20].strip()  # Limit to 20 chars
            
            if clean_title:  # Only add if we have valid text
                button_payload.append({
                    "type": "reply",
                    "reply": {
                        "id": f"btn_{i+1}",
                        "title": clean_title
                    }
                })
        
        # If no valid buttons, send as text instead
        if not button_payload:
            send_whatsapp_message(to, f"{text}\n\nPlease respond with your choice.")
            return
        
        payload = {
            "messaging_product": "whatsapp",
            "to": to,
            "type": "interactive",
            "interactive": {
                "type": "button",
                "body": {
                    "text": text[:1024] 
                },
                "action": {
                    "buttons": button_payload
                }
            }
        }
        
        print(f"Sending interactive message with payload: {json.dumps(payload, indent=2)}")
        
        response = requests.post(url, headers=headers, json=payload, timeout=10)
        response.raise_for_status()
        
        print(f"Interactive message sent successfully to {to}")
        
    except requests.exceptions.RequestException as e:
        print(f"Failed to send interactive message: {e}")
        if hasattr(e, 'response') and e.response is not None:
            print(f"Response content: {e.response.text}")
        
        # Fallback: Send as text message with options
        fallback_text = f"{text}\n\nPlease reply with:\n"
        for i, button in enumerate(buttons[:5], 1):  
            fallback_text += f"{i}. {button}\n"
        
        send_whatsapp_message(to, fallback_text)
        
    except Exception as e:
        print(f"Unexpected error in interactive message: {e}")
        # Fallback to simple text
        send_whatsapp_message(to, f"{text}\n\nPlease respond with your choice.")

def send_text_options_message(to, text, options):
    """Alternative: Send options as numbered text message"""
    message = f"{text}\n\nPlease reply with the number:\n"
    for i, option in enumerate(options, 1):
        message += f"{i}. {option}\n"   
    
    send_whatsapp_message(to, message)

def process_message(sender, text):
    # Initialize user session from Firebase
    initialize_user_session(sender)
    
    append_to_memory(sender, "user", text)
    current_state = user_states.get(sender)

    # --- Universal Cancel/Menu ---
    if text.lower().strip() in ['cancel', 'stop', 'menu', 'start', 'exit']:
        if current_state:
            clear_user_session(sender)
            send_whatsapp_message(sender, "Okay, I've canceled the current operation.")
        menu_text = "Hello! I'm your AI teaching assistant. Please choose an option:"
        options = [
            "Ask Question", 
            "Create Worksheet", 
            "Create PPT", 
            "Upload Materials",
            "View Uploaded Files", 
            "Podcast from Image", 
            "Summary from Image"
        ]
        send_menu_message(sender, menu_text, options)
        append_to_memory(sender, "assistant", "Displayed main menu.")
        return

    # --- Initial Triggers for New Conversations (if no active state) ---
    if not current_state:
        lower_text = text.lower().strip()

        if lower_text == "ask question":
            update_user_state(sender, state="awaiting_question")
            send_whatsapp_message(sender, "Of course! What is your question? I'll search through all uploaded documents and knowledge base.")
            return
        
        if lower_text == "create worksheet":
            update_user_state(sender, state="awaiting_worksheet_topic")
            send_whatsapp_message(sender, "Let's create a worksheet! What topic would you like the worksheet to be about?")
            return
            
        if lower_text == "create ppt":
            update_user_state(sender, state="awaiting_ppt_topic")
            send_whatsapp_message(sender, "Excellent! What topic would you like the presentation to be about?")
            return
            
        if lower_text == "upload materials" or lower_text == "upload material":
            update_user_state(sender, state="awaiting_material_file")
            send_whatsapp_message(sender, "Please send the file you'd like to upload (PDF, Word, PowerPoint, or Text). This will be added to my knowledge base for future questions!")
            return
        
        if lower_text == "view uploaded files":
            # Get files that ACTUALLY exist in vector DB
            vector_db_files = get_uploaded_files_list()
            
            if vector_db_files:
                files_list = []
                for file_path in vector_db_files:
                    filename = os.path.basename(file_path)
                    # Try to get upload time from Firebase if available
                    upload_time = "Recently uploaded"
                    if db:
                        try:
                            docs = db.collection('uploaded_files').where('filename', '==', filename).stream()
                            for doc in docs:
                                file_data = doc.to_dict()
                                if 'uploaded_at' in file_data:
                                    upload_time = file_data['uploaded_at'].strftime('%Y-%m-%d')
                                    break
                        except:
                            pass
                    
                    files_list.append(f"• {filename} (Uploaded: {upload_time})")
                
                files_text = "\n".join(files_list[:10])  # Show first 10 files
                send_whatsapp_message(sender, f"Files in your knowledge base:\n{files_text}")
                
                if len(vector_db_files) > 10:
                    send_whatsapp_message(sender, f"... and {len(vector_db_files) - 10} more files")
            else:
                send_whatsapp_message(sender, "No files in knowledge base. Use 'Upload Materials' to add files.")
            
            end_conversation_and_show_menu(sender, None)
            return
            
        if lower_text == "podcast from image":
            update_user_state(sender, state="awaiting_podcast_image")
            send_whatsapp_message(sender, "Please send me an image of the text you'd like to convert to a podcast.")
            return
            
        if lower_text == "summary from image":
            update_user_state(sender, state="awaiting_summary_image")
            send_whatsapp_message(sender, "Please send me an image of the text you'd like me to summarize.")
            return

    # --- State-Based Conversation Flow ---
    if current_state == "awaiting_question":
        response = query_dynamic_rag(text)
        end_conversation_and_show_menu(sender, response["result"])
        return
    
    if current_state == "awaiting_ppt_topic":
        try:
            topic = text.strip()
            update_user_state(sender, temp_data={'ppt_topic': topic})
            send_whatsapp_message(sender, f"Okay, generating a 10-slide presentation on '{topic}'. This may take a moment...")
            ppt_content = generate_ppt_content(topic)
            if ppt_content:
                ppt_bytes = create_ppt_file(topic, ppt_content)
                if ppt_bytes:
                    status = send_whatsapp_ppt(sender, ppt_bytes, f"{topic.replace(' ', '_')}.pptx")
                    if status == "success":
                        end_conversation_and_show_menu(sender, "I've sent the presentation.")
                    else:
                        end_conversation_and_show_menu(sender, "I created the PPT, but there was an error sending it.")
                else:
                    end_conversation_and_show_menu(sender, "Sorry, I generated the content but failed to create the PPT file.")
            else:
                end_conversation_and_show_menu(sender, "Sorry, I couldn't generate content for that topic.")
        except Exception as e:
             print(f"Error in PPT generation flow: {e}")
             end_conversation_and_show_menu(sender, "An unexpected error occurred.")
        return
    
    if current_state == "awaiting_material_file":
        # This will be handled by the webhook when documents are received
        send_whatsapp_message(sender, 
            "Please send the file you'd like to upload:\n"
            "• PDF, Word, PowerPoint, or Text files\n"
            "• I'll add it to my knowledge base\n"
            "• You can send multiple files"
        )
        return
    
    # In the awaiting_topic_confirmation state, enhance the user experience:
    if current_state == "awaiting_topic_confirmation":
        response_text = text.lower().strip()
        
        if response_text in ['yes', 'continue', 'y', 'proceed', '1']:
            # User confirmed to proceed with unrelated topic
            topic = user_temp_data[sender]['pending_topic']
            update_user_state(sender, 
                            state="awaiting_worksheet_format", 
                            temp_data={
                                'topic': topic, 
                                'using_general_knowledge': True
                            })
            
            # Clear the pending topic
            if 'pending_topic' in user_temp_data[sender]:
                del user_temp_data[sender]['pending_topic']
                
            send_whatsapp_message(sender, f"✅ Okay, generating worksheet on '{topic}' using general knowledge...")
            send_interactive_message(sender, f"What format would you like for '{topic}'?", ["PDF Worksheet", "Google Form Quiz"])
            
        elif response_text in ['no', 'change', 'n', 'change topic', '2']:
            # User wants to change topic - return to main menu as requested
            if 'pending_topic' in user_temp_data[sender]:
                del user_temp_data[sender]['pending_topic']
            
            # Clear state and return to main menu
            clear_user_session(sender)
            send_whatsapp_message(sender, "Okay, returning to main menu. You can try a different topic later.")
            
            menu_text = "What would you like to do next?"
            options = [
                "Ask Question", 
                "Create Worksheet", 
                "Create PPT", 
                "Upload Materials",
                "View Uploaded Files", 
                "Podcast from Image", 
                "Summary from Image"
            ]
            send_menu_message(sender, menu_text, options)
            
        else:
            # Handle unexpected responses - show options again
            vector_db_files = get_uploaded_files_list()
            message = (f"⚠️ No materials uploaded yet. The worksheet will be generated from general knowledge."
                    if not vector_db_files 
                    else f"⚠️ Topic doesn't match your uploaded materials. The worksheet will be generated from general knowledge.")
            message += " Continue?"
            
            send_interactive_message(sender, message, ["Yes, Continue", "No, Change Topic"])
        return

    if current_state == "awaiting_worksheet_topic":
        if text.lower() == "other topics":
            update_user_state(sender, state="awaiting_custom_topic")
            send_whatsapp_message(sender, "Please type the custom topic for your worksheet.")
            return
            
        topic = text.strip()
        is_relevant, reason = check_topic_relevance(topic)

        if is_relevant:
            # Topic matches RAG content
            update_user_state(sender, 
                            state="awaiting_worksheet_format", 
                            temp_data={'topic': topic})
            send_interactive_message(sender, f"Great! Topic is '{topic}'.\nWhat format would you like?", ["PDF Worksheet", "Google Form Quiz"])
        else:
            # Topic doesn't match - ask for confirmation
            update_user_state(sender, 
                            state="awaiting_topic_confirmation", 
                            temp_data={'pending_topic': topic})
            
            vector_db_files = get_uploaded_files_list()
            message = (f"⚠️ No materials uploaded yet. The worksheet on '{topic}' will be generated from general knowledge."
                    if not vector_db_files 
                    else f"⚠️ Topic '{topic}' doesn't match your uploaded materials. The worksheet will be generated from general knowledge.")
            message += " Continue?"
            
            send_interactive_message(sender, message, ["Yes, Continue", "No, Change Topic"])
        return
        
    if current_state == "awaiting_custom_topic":
    # Check if custom topic matches RAG content
        topic = text.strip()
        is_relevant, reason = check_topic_relevance(topic)
        
        if is_relevant:
            # ✅ FIXED: Use update_user_state consistently
            update_user_state(sender, 
                            state="awaiting_worksheet_format", 
                            temp_data={'topic': topic})
            send_interactive_message(sender, f"Great! Topic is '{topic}'.\nWhat format would you like?", ["PDF Worksheet", "Google Form Quiz"])
        else:
            # ✅ FIXED: Use update_user_state consistently
            update_user_state(sender, 
                            state="awaiting_topic_confirmation", 
                            temp_data={'pending_topic': topic})
            
            # Check if there are any uploaded files at all
            vector_db_files = get_uploaded_files_list()
            if not vector_db_files:
                message = f"⚠️ No materials uploaded yet. The worksheet on '{topic}' will be generated from general knowledge, not your specific materials. Continue?"
            else:
                message = f"⚠️ Topic '{topic}' doesn't match your uploaded materials. The worksheet will be generated from general knowledge. Continue?"
            
            send_interactive_message(sender, message, ["Yes, Continue", "No, Change Topic"])
        return
        
    if current_state == "awaiting_worksheet_format":
        
        update_user_state(sender, 
                        state="awaiting_worksheet_quantity", 
                        temp_data={'format': text})
        send_interactive_message(sender, f"Perfect, a {text}.\nHow many questions?", ["5", "10", "15"])
        return

    if current_state == "awaiting_worksheet_quantity":
        try:
            user_temp_data[sender]['quantity'] = int(text.strip())
            if "Google Form" in user_temp_data[sender]['format']:
                user_temp_data[sender]['type'] = "mcq"
                topic = user_temp_data[sender]['topic']
                quantity = user_temp_data[sender]['quantity']
                send_whatsapp_message(sender, f"Okay! Generating a {quantity}-question Google Form quiz on '{topic}'. Please wait...")
                
                worksheet_content_result = generate_worksheet_content_text(topic, quantity, "mcq", user_memory[sender])
                if worksheet_content_result.get("source") == "generated_worksheet_text":
                    form_result = create_google_form_mcq(f"Quiz: {topic}", worksheet_content_result["result"])
                    user_temp_data[sender]['form_url'] = form_result.get("result", "").split(' ')[-1]
                    send_whatsapp_message(sender, form_result["result"])
                    
                    
                    update_user_state(sender, state="awaiting_assignment_title")
                    send_whatsapp_message(sender, "The quiz is ready! Now, what should the title of the assignment be for this quiz?")
                else:
                    end_conversation_and_show_menu(sender, "Sorry, I couldn't generate the quiz content.")
                return
            else: # PDF Flow
                
                update_user_state(sender, state="awaiting_worksheet_type")
                send_interactive_message(sender, f"Perfect, {text} questions.\nNow, what type of questions?", ["MCQ", "Short Answer", "Numerical"])
        except ValueError:
            send_whatsapp_message(sender, "Please select a valid number from the buttons.")
        return

    if current_state == "awaiting_worksheet_type":
        try:
            topic = user_temp_data[sender]['topic']
            quantity = user_temp_data[sender]['quantity']
            q_type = text.lower().replace(" answer", "")
            valid_types = ["mcq", "short", "numerical"]
            if q_type not in valid_types:
                send_whatsapp_message(sender, "That's not a valid type."); return

            # Check if we're using general knowledge
            using_general = user_temp_data[sender].get('using_general_knowledge', False)
            
            if using_general:
                send_whatsapp_message(sender, f"🔄 Generating {quantity} {q_type} questions on '{topic}' using general knowledge...")
            else:
                send_whatsapp_message(sender, f"🔄 Generating {quantity} {q_type} questions on '{topic}' based on your uploaded materials...")
            
            worksheet_content_result = generate_worksheet_content_text(topic, quantity, q_type, user_memory[sender])
            
            if worksheet_content_result and worksheet_content_result.get("source") == "generated_worksheet_text":
                full_content = worksheet_content_result["result"]
                questions_text, answers_text = (full_content.split("--- ANSWERS ---", 1) + ["No answer key generated."])[:2]
                user_temp_data[sender]['questions_text'] = questions_text.strip()
                user_temp_data[sender]['answers_text'] = answers_text.strip()
                
                if worksheet_pdf_bytes := create_pdf_locally(f"Worksheet: {topic.title()}", questions_text.strip()):
                    send_whatsapp_document(sender, worksheet_pdf_bytes, f"{topic.replace(' ', '_')}_worksheet.pdf")
                
                if answer_key_pdf_bytes := create_pdf_locally(f"Answer Key: {topic.title()}", answers_text.strip()):
                    send_whatsapp_document(sender, answer_key_pdf_bytes, f"{topic.replace(' ', '_')}_answers.pdf")

                user_states[sender] = "awaiting_classroom_post_type"
                send_interactive_message(sender, "I've sent the PDFs. How should I post this to Classroom?", ["Post only questions", "Post with answers", "Don't post"])
            else:
                end_conversation_and_show_menu(sender, "Sorry, I couldn't generate the worksheet content.")
        except (ValueError, KeyError) as e:
            print(f"Error in worksheet type state: {e}")
            end_conversation_and_show_menu(sender, "An error occurred. Please start over.")
        return

    if current_state == "awaiting_classroom_post_type":
        user_temp_data[sender]['post_choice'] = text
        if text.lower() == "post only questions" or text.lower() == "post with answers":
            # ✅ FIXED: Use update_user_state consistently
            update_user_state(sender, state="awaiting_assignment_title")
            send_whatsapp_message(sender, "Great! What should the assignment title be?")
        else: # Don't Post
            end_conversation_and_show_menu(sender, "Okay, I won't post to Classroom.")
        return

    if current_state == "awaiting_assignment_title":
        user_temp_data[sender]['assignment_title'] = text.strip()
        # SIMPLIFIED: Skip due date question and go straight to posting
        title = user_temp_data[sender]['assignment_title']
        handle_final_classroom_post(sender, title)
        return
    
    if text.lower() in ['reset', 'start over', 'restart', 'clear']:
        clear_user_session(sender)
        send_whatsapp_message(sender, "🔄 Session reset completely. Starting fresh...")
        menu_text = "What would you like to do next?"
        options = [
            "Ask Question", 
            "Create Worksheet", 
            "Create PPT", 
            "Upload Materials",
            "View Uploaded Files", 
            "Podcast from Image", 
            "Summary from Image"
        ]
        send_menu_message(sender, menu_text, options)
        return

    ## --- Fallback to a General Query ---
    print(f"Handling as a general query: '{text}'")
    response = query_dynamic_rag(text)
    send_whatsapp_message(sender, response["result"])
    end_conversation_and_show_menu(sender, None)

#
# ---------------------------------
# --- FLASK WEBHOOK ENDPOINTS (Updated for Document Uploads)
# ---------------------------------
@app.route("/webhook", methods=["GET", "POST"])
def webhook():
    if request.method == "GET":
        if request.args.get("hub.verify_token") == VERIFY_TOKEN:
            return request.args.get("hub.challenge")
        return "Verification token mismatch", 403
    
    data = request.get_json()
    try:
        if data and "entry" in data:
            for entry in data.get("entry", []):
                for change in entry.get("changes", []):
                    value = change.get("value", {})
                    if "messages" in value:
                        for msg in value.get("messages", []):
                            sender = msg["from"]
                            text = ""
                            msg_type = msg.get("type")

                            # Handle document messages for dynamic RAG
                            if msg_type == "document":
                                current_state = user_states.get(sender)
                                
                                # Handle different document upload states
                                if current_state in ["awaiting_bulk_upload", "awaiting_material_file"]:
                                    media_id = msg['document']['id']
                                    filename = msg['document'].get('filename', 'uploaded_file')
                                    
                                    # Initialize bulk upload items if not exists
                                    if 'bulk_upload_items' not in user_temp_data[sender]:
                                        user_temp_data[sender]['bulk_upload_items'] = []
                                    
                                    user_temp_data[sender]['bulk_upload_items'].append({
                                        'id': media_id,
                                        'filename': filename
                                    })
                                    
                                    # Wait for more files then process ALL files
                                    threading.Timer(8.0, process_bulk_upload, args=[sender]).start()
                                    
                                    send_whatsapp_message(sender, f"✅ Added {filename} to upload queue...")
                                else:
                                    # If not in upload state, process as single document
                                    media_id = msg['document']['id']
                                    filename = msg['document'].get('filename', 'uploaded_file')
                                    thread = threading.Thread(target=handle_document_upload, args=(sender, media_id, filename))
                                    thread.start()
                                return "ok", 200

                            # Handle other message types
                            if msg_type == "button":
                                text = msg["button"]["text"]
                            elif msg_type == "text":
                                text = msg["text"]["body"]
                            elif msg_type == "interactive" and msg.get("interactive", {}).get("type") == "button_reply":
                                text = msg["interactive"]["button_reply"]["title"]
                            elif msg_type == "interactive" and msg.get("interactive", {}).get("type") == "list_reply":
                                text = msg["interactive"]["list_reply"]["title"]
                            elif msg_type == "audio":
                                media_id = msg['audio']['id']
                                thread = threading.Thread(target=handle_audio_task, args=(sender, media_id))
                                thread.start()
                            elif msg_type == "image":
                                current_state = user_states.get(sender)
                                media_id = msg['image']['id']
                                
                                if current_state == "awaiting_podcast_image":
                                    thread = threading.Thread(target=handle_podcast_task, args=(sender, media_id))
                                    thread.start()
                                elif current_state == "awaiting_summary_image":
                                    thread = threading.Thread(target=handle_summary_task, args=(sender, media_id))
                                    thread.start()
                                else:
                                    prompt = msg.get("image", {}).get("caption", "Explain this image.")
                                    thread = threading.Thread(target=handle_image_task, args=(sender, media_id, prompt))
                                    thread.start()
                                return "ok", 200
                            
                            if text:
                                process_message(sender, text)
    except Exception as e:
        print(f"Webhook processing error: {e}")
        traceback.print_exc()
    return "ok", 200

def process_bulk_upload(sender):
    """Process accumulated bulk upload items"""
    if sender in user_temp_data and 'bulk_upload_items' in user_temp_data[sender]:
        items = user_temp_data[sender]['bulk_upload_items']
        if items:
            print(f"🔄 Processing {len(items)} files for {sender}")
            
            # Clear the items to avoid reprocessing
            user_temp_data[sender]['bulk_upload_items'] = []
            
            # Process the bulk upload
            handle_bulk_document_upload(sender, items)
        else:
            print(f"ℹ️ No items to process for {sender}")
    else:
        print(f"ℹ️ No bulk upload items found for {sender}")

@app.route("/status", methods=["GET"])
def status():
    return {"status": "running"}, 200

if __name__ == "__main__":
    os.makedirs("data", exist_ok=True)
    os.makedirs("vector_index", exist_ok=True)
    os.makedirs(DYNAMIC_VECTOR_INDEX_PATH, exist_ok=True)
    
    # Initialize APIs
    init_google_apis()
    
    # Check Firebase connection
    if db:
        print("Firebase connected successfully")
        cleanup_orphaned_file_metadata()
    else:
        print("Firebase not connected - running in local mode only")
    
    # --- Send Startup Template Message ---
    print("Sending startup template to users...")
    for number in student_phone_numbers:
        send_start_template(number)
    print("Finished sending templates.")

    # ONLY Dynamic RAG - No static embeddings or databases
    print("Dynamic RAG system ready - using only user-uploaded documents")
    
    print("Starting Flask app...")
    app.run(port=5000, debug=False)